#!/usr/bin/env python3
"""
Proof of Concept: PowerPoint Style Learning

This script demonstrates the basic style learning capabilities, creating sample
presentations with different styles and then analyzing them to extract patterns.

Week 1 Implementation: Basic proof-of-concept for style learning
"""

import os
import tempfile
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx not available. Please install with: pip install python-pptx")
    exit(1)

# Import our style analysis module
try:
    from style_analysis import StyleAnalyzer
except ImportError:
    print("Style analysis module not found. Make sure style_analysis.py is available.")
    exit(1)

class StyleLearningDemo:
    """Demonstrates style learning capabilities with sample presentations"""
    
    def __init__(self):
        self.analyzer = StyleAnalyzer()
        self.temp_dir = tempfile.mkdtemp()
        self.sample_files = []
    
    def create_corporate_style_presentation(self) -> str:
        """Create a sample presentation with corporate styling"""
        prs = Presentation()
        
        # Corporate colors: Blue (#1E3A8A), Gray (#6B7280), White (#FFFFFF)
        corporate_blue = RGBColor(30, 58, 138)
        corporate_gray = RGBColor(107, 114, 128)
        
        # Slide 1: Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = "Corporate Quarterly Report"
        title.text_frame.paragraphs[0].font.name = "Calibri"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = corporate_blue
        
        subtitle.text = "Q4 2024 Financial Performance"
        subtitle.text_frame.paragraphs[0].font.name = "Calibri"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = corporate_gray
        
        # Slide 2: Content slide
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = content_slide.shapes.title
        content = content_slide.placeholders[1]
        
        title.text = "Key Achievements"
        title.text_frame.paragraphs[0].font.name = "Calibri"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = corporate_blue
        
        content.text = "• Revenue increased by 15%\n• Market share expanded to 28%\n• Customer satisfaction: 94%\n• New product launches: 3"
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.name = "Calibri"
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = corporate_gray
        
        # Slide 3: Another content slide for more data points
        content_slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title2 = content_slide2.shapes.title
        content2 = content_slide2.placeholders[1]
        
        title2.text = "Financial Highlights"
        title2.text_frame.paragraphs[0].font.name = "Calibri"
        title2.text_frame.paragraphs[0].font.size = Pt(36)
        title2.text_frame.paragraphs[0].font.bold = True
        title2.text_frame.paragraphs[0].font.color.rgb = corporate_blue
        
        content2.text = "• Total Revenue: $2.4M\n• Operating Margin: 18.5%\n• EBITDA: $445K\n• Cash Flow: $380K"
        for paragraph in content2.text_frame.paragraphs:
            paragraph.font.name = "Calibri"
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = corporate_gray
        
        # Save the presentation
        file_path = os.path.join(self.temp_dir, "corporate_style.pptx")
        prs.save(file_path)
        self.sample_files.append(file_path)
        
        print(f"Created corporate style presentation: {file_path}")
        return file_path
    
    def create_creative_style_presentation(self) -> str:
        """Create a sample presentation with creative/artistic styling"""
        prs = Presentation()
        
        # Creative colors: Purple (#8B5CF6), Orange (#F97316), Pink (#EC4899)
        creative_purple = RGBColor(139, 92, 246)
        creative_orange = RGBColor(249, 115, 22)
        creative_pink = RGBColor(236, 72, 153)
        
        # Slide 1: Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = "Creative Portfolio 2024"
        title.text_frame.paragraphs[0].font.name = "Montserrat"  # Different font family
        title.text_frame.paragraphs[0].font.size = Pt(48)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.italic = True  # Add italic for creative flair
        title.text_frame.paragraphs[0].font.color.rgb = creative_purple
        
        subtitle.text = "Innovative Design Solutions"
        subtitle.text_frame.paragraphs[0].font.name = "Montserrat"
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.color.rgb = creative_orange
        
        # Slide 2: Content slide
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = content_slide.shapes.title
        content = content_slide.placeholders[1]
        
        title.text = "Design Principles"
        title.text_frame.paragraphs[0].font.name = "Montserrat"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = creative_purple
        
        content.text = "🎨 Bold visual storytelling\n🌈 Vibrant color palettes\n✨ Dynamic compositions\n🚀 Innovation-driven approach"
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.name = "Montserrat"
            paragraph.font.size = Pt(16)  # Smaller font for creative style
            paragraph.font.color.rgb = creative_pink
        
        # Slide 3: More creative content
        content_slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title2 = content_slide2.shapes.title
        content2 = content_slide2.placeholders[1]
        
        title2.text = "Recent Projects"
        title2.text_frame.paragraphs[0].font.name = "Montserrat"
        title2.text_frame.paragraphs[0].font.size = Pt(32)
        title2.text_frame.paragraphs[0].font.bold = True
        title2.text_frame.paragraphs[0].font.color.rgb = creative_purple
        
        content2.text = "🎯 Brand Identity Redesign\n📱 Mobile App Interface\n🌐 Website Transformation\n📊 Data Visualization Suite"
        for paragraph in content2.text_frame.paragraphs:
            paragraph.font.name = "Montserrat"
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = creative_pink
        
        # Save the presentation
        file_path = os.path.join(self.temp_dir, "creative_style.pptx")
        prs.save(file_path)
        self.sample_files.append(file_path)
        
        print(f"Created creative style presentation: {file_path}")
        return file_path
    
    def create_academic_style_presentation(self) -> str:
        """Create a sample presentation with academic/research styling"""
        prs = Presentation()
        
        # Academic colors: Dark Blue (#1E40AF), Black (#000000), Dark Gray (#374151)
        academic_blue = RGBColor(30, 64, 175)
        academic_black = RGBColor(0, 0, 0)
        academic_gray = RGBColor(55, 65, 81)
        
        # Slide 1: Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = "Machine Learning Applications in Healthcare"
        title.text_frame.paragraphs[0].font.name = "Times New Roman"  # Traditional academic font
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = academic_blue
        
        subtitle.text = "A Systematic Review and Meta-Analysis"
        subtitle.text_frame.paragraphs[0].font.name = "Times New Roman"
        subtitle.text_frame.paragraphs[0].font.size = Pt(22)
        subtitle.text_frame.paragraphs[0].font.color.rgb = academic_gray
        
        # Slide 2: Content slide
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = content_slide.shapes.title
        content = content_slide.placeholders[1]
        
        title.text = "Research Methodology"
        title.text_frame.paragraphs[0].font.name = "Times New Roman"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = academic_blue
        
        content.text = "• Systematic literature review (2019-2024)\n• Database search: PubMed, IEEE, ACM\n• Inclusion criteria: peer-reviewed articles\n• Sample size: 247 studies"
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.name = "Times New Roman"
            paragraph.font.size = Pt(20)  # Larger font for readability
            paragraph.font.color.rgb = academic_black
        
        # Slide 3: Results
        content_slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title2 = content_slide2.shapes.title
        content2 = content_slide2.placeholders[1]
        
        title2.text = "Key Findings"
        title2.text_frame.paragraphs[0].font.name = "Times New Roman"
        title2.text_frame.paragraphs[0].font.size = Pt(32)
        title2.text_frame.paragraphs[0].font.bold = True
        title2.text_frame.paragraphs[0].font.color.rgb = academic_blue
        
        content2.text = "• 89% accuracy in diagnostic applications\n• 15% reduction in treatment time\n• Significant cost savings (p < 0.05)\n• High inter-rater reliability (κ = 0.84)"
        for paragraph in content2.text_frame.paragraphs:
            paragraph.font.name = "Times New Roman"
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = academic_black
        
        # Save the presentation
        file_path = os.path.join(self.temp_dir, "academic_style.pptx")
        prs.save(file_path)
        self.sample_files.append(file_path)
        
        print(f"Created academic style presentation: {file_path}")
        return file_path
    
    def analyze_and_compare_styles(self):
        """Analyze all created presentations and compare their styles"""
        print("\n" + "="*60)
        print("STYLE ANALYSIS COMPARISON")
        print("="*60)
        
        analysis_results = {}
        
        # Analyze each presentation
        for file_path in self.sample_files:
            file_name = Path(file_path).stem
            print(f"\nAnalyzing {file_name}...")
            
            try:
                analysis = self.analyzer.analyze_presentation_style(file_path)
                analysis_results[file_name] = analysis
                
                # Create and save style profile
                profile_name = f"{file_name}_profile"
                self.analyzer.create_style_profile(analysis, profile_name)
                
                # Save profile to JSON
                profile_json_path = os.path.join(self.temp_dir, f"{profile_name}.json")
                self.analyzer.save_style_profile(profile_name, profile_json_path)
                
                print(f"✓ Style profile created: {profile_name}")
                print(f"✓ Profile saved to: {profile_json_path}")
                
            except Exception as e:
                print(f"✗ Error analyzing {file_name}: {e}")
        
        # Compare and summarize findings
        self.print_style_comparison(analysis_results)
        
        return analysis_results
    
    def print_style_comparison(self, analysis_results: dict):
        """Print a comparison of the analyzed styles"""
        print("\n" + "="*60)
        print("STYLE COMPARISON SUMMARY")
        print("="*60)
        
        for name, analysis in analysis_results.items():
            print(f"\n📊 {name.upper()}:")
            print(f"   Slides: {analysis['slide_count']}")
            print(f"   Consistency Score: {analysis['consistency_score']:.2f}")
            
            # Font analysis
            fonts = analysis['fonts']
            print(f"   Primary Font: {fonts['primary_font']}")
            print(f"   Font Variations: {len(fonts['font_usage'])}")
            
            # Color analysis
            colors = analysis['colors']
            print(f"   Color Palette Size: {colors['total_unique_colors']}")
            if colors['primary_palette']:
                top_color = colors['primary_palette'][0][0]
                print(f"   Dominant Color: {top_color}")
            
            # Size patterns
            if fonts['common_sizes']:
                common_size = fonts['common_sizes'][0][0]
                print(f"   Most Common Size: {common_size}pt")
        
        # Overall insights
        print(f"\n🔍 INSIGHTS:")
        if analysis_results:
            consistencies = [analysis['consistency_score'] for analysis in analysis_results.values()]
            if consistencies:
                avg_consistency = sum(consistencies) / len(consistencies)
                print(f"   Average Consistency Score: {avg_consistency:.2f}")
                
                most_consistent = max(analysis_results.items(), key=lambda x: x[1]['consistency_score'])
                print(f"   Most Consistent Style: {most_consistent[0]} ({most_consistent[1]['consistency_score']:.2f})")
            else:
                print("   No consistency scores available")
        else:
            print("   No analysis results available")
    
    def demonstrate_style_application(self):
        """Demonstrate applying learned styles to new content"""
        print("\n" + "="*60)
        print("STYLE APPLICATION DEMONSTRATION")
        print("="*60)
        
        # This would be implemented in the full version to show how
        # learned styles can be applied to new presentations
        print("📝 Feature Preview: Style Application")
        print("   → Load existing style profile")
        print("   → Create new presentation with consistent styling")
        print("   → Apply learned font hierarchies, colors, and layouts")
        print("   → Validate style consistency")
        print("\n   Implementation: Next phase development")
    
    def cleanup(self):
        """Clean up temporary files"""
        try:
            import shutil
            shutil.rmtree(self.temp_dir)
            print(f"\n🧹 Cleaned up temporary files from {self.temp_dir}")
        except Exception as e:
            print(f"⚠️  Warning: Could not clean up temporary files: {e}")
    
    def run_demo(self):
        """Run the complete style learning demonstration"""
        print("🚀 PowerPoint Style Learning - Proof of Concept")
        print("="*60)
        
        try:
            # Create sample presentations with different styles
            print("\n📝 Creating sample presentations...")
            self.create_corporate_style_presentation()
            self.create_creative_style_presentation()
            self.create_academic_style_presentation()
            
            # Analyze and compare styles
            print(f"\n🔍 Analyzing {len(self.sample_files)} presentations...")
            analysis_results = self.analyze_and_compare_styles()
            
            # Demonstrate style application (preview)
            self.demonstrate_style_application()
            
            print(f"\n✅ Style learning demonstration completed successfully!")
            print(f"📁 Temporary files created in: {self.temp_dir}")
            print(f"📋 {len(self.analyzer.style_profiles)} style profiles created")
            
            return analysis_results
            
        except Exception as e:
            print(f"\n❌ Demo failed: {e}")
            raise
        finally:
            # Note: Not cleaning up automatically so you can examine the files
            print(f"\n💡 Tip: Examine the created files in {self.temp_dir}")
            print("     Run demo.cleanup() to remove temporary files when done.")

def main():
    """Main function to run the style learning proof of concept"""
    demo = StyleLearningDemo()
    
    try:
        results = demo.run_demo()
        
        # Optional: Print JSON example of a style profile
        if demo.analyzer.style_profiles:
            profile_name = list(demo.analyzer.style_profiles.keys())[0]
            print(f"\n📄 Example Style Profile JSON ({profile_name}):")
            print("-" * 40)
            profile = demo.analyzer.style_profiles[profile_name]
            # Convert to dict for JSON serialization (simplified)
            print(json.dumps({
                'name': profile.name,
                'description': profile.description,
                'confidence_score': profile.confidence_score,
                'created_from': profile.created_from
            }, indent=2))
        
        return demo
        
    except Exception as e:
        print(f"❌ Error running demo: {e}")
        demo.cleanup()
        raise

if __name__ == "__main__":
    # Run the demonstration
    demo_instance = main()
    
    # Keep demo instance available for interactive exploration
    print(f"\n🔧 Demo instance available as 'demo_instance'")
    print("     Try: demo_instance.analyzer.list_style_profiles()")
    print("     Or:  demo_instance.cleanup() when finished") 
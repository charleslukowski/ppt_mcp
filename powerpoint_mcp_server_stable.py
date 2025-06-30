#!/usr/bin/env python3
"""
PowerPoint MCP Server - Stable Production Version

This version includes core improvements from our development but maintains
stability and compatibility with Cursor's settings UI.

Key improvements included:
- Input validation (basic, no Pydantic dependency)
- Enhanced success messages
- Core 5 essential tools
- Post-processing fixes
- Simple error handling
"""

import asyncio
import json
import logging
import os
import sys
import tempfile
from typing import Any, Dict, List, Optional
import platform
from datetime import datetime

# Configure logging
logging.basicConfig(level=logging.INFO, stream=sys.stderr)
logger = logging.getLogger("powerpoint-mcp-stable")

# Windows-specific COM imports for screenshot functionality
if platform.system() == "Windows":
    try:
        import win32com.client
        import pythoncom
        WIN32_COM_AVAILABLE = True
    except ImportError:
        WIN32_COM_AVAILABLE = False
        logger.warning("win32com not available - screenshot functionality will be disabled")
else:
    WIN32_COM_AVAILABLE = False

# Core dependencies only
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError as e:
    print(f"python-pptx library not found: {e}")
    print("Please install with: pip install python-pptx")
    sys.exit(1)

try:
    from mcp.server import Server
    from mcp.server.models import InitializationOptions
    from mcp.server.stdio import stdio_server
    from mcp.types import Tool, TextContent, EmbeddedResource
except ImportError as e:
    print(f"MCP library not found: {e}")
    print("Please install with: pip install mcp")
    sys.exit(1)

# =============================================================================
# SIMPLIFIED INPUT VALIDATION (No Pydantic dependency)
# =============================================================================

def validate_basic_args(tool_name: str, arguments: Dict[str, Any]) -> Dict[str, Any]:
    """Basic input validation without external dependencies"""
    
    # Get presentation_id if present
    prs_id = arguments.get("presentation_id")
    if prs_id and not isinstance(prs_id, str):
        raise ValueError("presentation_id must be a string")
    
    # Get slide_index if present
    slide_index = arguments.get("slide_index")
    if slide_index is not None:
        if not isinstance(slide_index, int) or slide_index < 0:
            raise ValueError("slide_index must be a non-negative integer")
    
    # Tool-specific validation
    if tool_name == "add_text_box":
        text = arguments.get("text", "")
        if not text or not isinstance(text, str):
            raise ValueError("text must be a non-empty string")
        
        font_size = arguments.get("font_size", 18)
        if not isinstance(font_size, int) or font_size < 8 or font_size > 72:
            raise ValueError("font_size must be between 8 and 72")
        
        font_name = arguments.get("font_name", "Calibri")
        if not isinstance(font_name, str):
            raise ValueError("font_name must be a string")
        
        text_alignment = arguments.get("text_alignment", "left")
        valid_alignments = ["left", "center", "right", "justify"]
        if text_alignment.lower() not in valid_alignments:
            raise ValueError(f"text_alignment must be one of: {valid_alignments}")
        
        border_width = arguments.get("border_width", 0)
        if not isinstance(border_width, (int, float)) or border_width < 0:
            raise ValueError("border_width must be a non-negative number")
    
    elif tool_name == "add_image":
        image_source = arguments.get("image_source", "")
        if not image_source or not isinstance(image_source, str):
            raise ValueError("image_source must be a non-empty string")
    
    elif tool_name == "add_chart":
        chart_type = arguments.get("chart_type", "")
        valid_types = ["column", "bar", "line", "pie", "area"]
        if chart_type not in valid_types:
            raise ValueError(f"chart_type must be one of: {valid_types}")
        
        categories = arguments.get("categories", [])
        if not categories or not isinstance(categories, list):
            raise ValueError("categories must be a non-empty list")
        
        series_data = arguments.get("series_data", {})
        if not series_data or not isinstance(series_data, dict):
            raise ValueError("series_data must be a non-empty dictionary")
    
    elif tool_name == "save_presentation":
        file_path = arguments.get("file_path", "")
        if not file_path or not isinstance(file_path, str):
            raise ValueError("file_path must be a non-empty string")
    
    elif tool_name == "load_presentation":
        file_path = arguments.get("file_path", "")
        if not file_path or not isinstance(file_path, str):
            raise ValueError("file_path must be a non-empty string")
    
    elif tool_name == "add_slide":
        layout_index = arguments.get("layout_index", 6)
        if not isinstance(layout_index, int) or layout_index < 0:
            raise ValueError("layout_index must be a non-negative integer")
    
    elif tool_name == "extract_text":
        # No additional validation needed beyond presentation_id
        pass
    
    elif tool_name == "get_presentation_info":
        # No additional validation needed beyond presentation_id
        pass
    
    elif tool_name == "delete_shape":
        shape_index = arguments.get("shape_index")
        if shape_index is None or not isinstance(shape_index, int) or shape_index < 0:
            raise ValueError("shape_index must be a non-negative integer")
    
    elif tool_name == "delete_slide":
        # slide_index validation already handled above
        pass
    
    elif tool_name == "clear_slide":
        # slide_index validation already handled above  
        pass
    
    elif tool_name == "list_slide_content":
        # slide_index validation already handled above
        pass
    
    elif tool_name == "format_existing_text":
        shape_index = arguments.get("shape_index")
        if shape_index is None or not isinstance(shape_index, int) or shape_index < 0:
            raise ValueError("shape_index must be a non-negative integer")
        
        # Validate formatting parameters if provided
        font_size = arguments.get("font_size")
        if font_size is not None and (not isinstance(font_size, int) or font_size < 8 or font_size > 72):
            raise ValueError("font_size must be between 8 and 72")
        
        text_alignment = arguments.get("text_alignment")
        if text_alignment is not None:
            valid_alignments = ["left", "center", "right", "justify"]
            if text_alignment.lower() not in valid_alignments:
                raise ValueError(f"text_alignment must be one of: {valid_alignments}")
    
    elif tool_name == "set_slide_background":
        background_color = arguments.get("background_color")
        background_image = arguments.get("background_image")
        if not background_color and not background_image:
            raise ValueError("Either background_color or background_image must be provided")
    
    # Table-specific validation
    elif tool_name == "add_table":
        rows = arguments.get("rows")
        cols = arguments.get("cols")
        if not isinstance(rows, int) or rows < 1 or rows > 50:
            raise ValueError("rows must be between 1 and 50")
        if not isinstance(cols, int) or cols < 1 or cols > 20:
            raise ValueError("cols must be between 1 and 20")
            
    elif tool_name in ["set_table_cell", "style_table_cell"]:
        table_index = arguments.get("table_index")
        if table_index is None or not isinstance(table_index, int) or table_index < 0:
            raise ValueError("table_index must be a non-negative integer")
        
        row = arguments.get("row")
        col = arguments.get("col") 
        if row is None or not isinstance(row, int) or row < 0:
            raise ValueError("row must be a non-negative integer")
        if col is None or not isinstance(col, int) or col < 0:
            raise ValueError("col must be a non-negative integer")
        
        if tool_name == "set_table_cell":
            text = arguments.get("text", "")
            if not isinstance(text, str):
                raise ValueError("text must be a string")
                
    elif tool_name == "style_table_range":
        table_index = arguments.get("table_index")
        if table_index is None or not isinstance(table_index, int) or table_index < 0:
            raise ValueError("table_index must be a non-negative integer")
            
        # Range validation
        start_row = arguments.get("start_row")
        end_row = arguments.get("end_row")
        start_col = arguments.get("start_col")
        end_col = arguments.get("end_col")
        
        if start_row is None or not isinstance(start_row, int) or start_row < 0:
            raise ValueError("start_row must be a non-negative integer")
        if end_row is None or not isinstance(end_row, int) or end_row < 0:
            raise ValueError("end_row must be a non-negative integer")
        if start_col is None or not isinstance(start_col, int) or start_col < 0:
            raise ValueError("start_col must be a non-negative integer")
        if end_col is None or not isinstance(end_col, int) or end_col < 0:
            raise ValueError("end_col must be a non-negative integer")
            
        if start_row > end_row:
            raise ValueError("start_row must be <= end_row")
        if start_col > end_col:
            raise ValueError("start_col must be <= end_col")
            
    elif tool_name == "modify_table_structure":
        table_index = arguments.get("table_index")
        if table_index is None or not isinstance(table_index, int) or table_index < 0:
            raise ValueError("table_index must be a non-negative integer")
            
        action = arguments.get("action")
        valid_actions = ["add_row", "delete_row", "add_column", "delete_column"]
        if action not in valid_actions:
            raise ValueError(f"action must be one of: {valid_actions}")
            
        index = arguments.get("index")
        if index is None or not isinstance(index, int) or index < 0:
            raise ValueError("index must be a non-negative integer")
    
    elif tool_name == "get_table_info":
        table_index = arguments.get("table_index")
        if table_index is None or not isinstance(table_index, int) or table_index < 0:
            raise ValueError("table_index must be a non-negative integer")
    
    elif tool_name == "create_table_with_data":
        table_data = arguments.get("table_data", [])
        if not isinstance(table_data, list) or not table_data:
            raise ValueError("table_data must be a non-empty list")
        
        if not all(isinstance(row, list) for row in table_data):
            raise ValueError("table_data must be a list of lists")
        
        if not table_data[0]:
            raise ValueError("table_data rows cannot be empty")
        
        # Check consistent row lengths
        expected_cols = len(table_data[0])
        for i, row in enumerate(table_data):
            if len(row) != expected_cols:
                raise ValueError(f"All rows must have the same number of columns. Row {i} has {len(row)} columns, expected {expected_cols}")
        
        # Validate headers if provided
        headers = arguments.get("headers")
        if headers is not None:
            if not isinstance(headers, list):
                raise ValueError("headers must be a list")
            if len(headers) != expected_cols:
                raise ValueError(f"headers length ({len(headers)}) must match table columns ({expected_cols})")
        
        # Validate style objects if provided
        header_style = arguments.get("header_style")
        if header_style is not None and not isinstance(header_style, dict):
            raise ValueError("header_style must be a dictionary")
        
        data_style = arguments.get("data_style")
        if data_style is not None and not isinstance(data_style, dict):
            raise ValueError("data_style must be a dictionary")
    
    elif tool_name == "modify_table_structure":
        table_index = arguments.get("table_index")
        if table_index is None or not isinstance(table_index, int) or table_index < 0:
            raise ValueError("table_index must be a non-negative integer")
        
        operation = arguments.get("operation")
        valid_operations = ["add_row", "remove_row", "add_column", "remove_column"]
        if operation not in valid_operations:
            raise ValueError(f"operation must be one of: {valid_operations}")
        
        position = arguments.get("position")
        if position is not None and (not isinstance(position, int) or position < 0):
            raise ValueError("position must be a non-negative integer")
        
        count = arguments.get("count", 1)
        if not isinstance(count, int) or count < 1 or count > 20:
            raise ValueError("count must be between 1 and 20")
    
    return arguments

# =============================================================================
# ENHANCED SUCCESS MESSAGES
# =============================================================================

def format_success_message(tool_name: str, **kwargs) -> str:
    """Generate specific, actionable success messages"""
    
    if tool_name == "add_text_box":
        slide_idx = kwargs.get('slide_index', 0)
        font_size = kwargs.get('font_size', 18)
        font_name = kwargs.get('font_name', 'Calibri')
        text_alignment = kwargs.get('text_alignment', 'left')
        font_color = kwargs.get('font_color')
        fill_color = kwargs.get('fill_color')
        text_preview = kwargs.get('text', '')[:40] + ('...' if len(kwargs.get('text', '')) > 40 else '')
        
        # Build formatting description
        format_desc = f"{font_size}pt {font_name}, {text_alignment} aligned"
        if font_color:
            format_desc += f", color: {font_color}"
        if fill_color:
            format_desc += f", background: {fill_color}"
        
        return f"✅ Added formatted text box to slide {slide_idx + 1}: \"{text_preview}\" ({format_desc})"
    
    elif tool_name == "add_image":
        slide_idx = kwargs.get('slide_index', 0)
        image_source = kwargs.get('image_source', '')
        image_name = os.path.basename(image_source) if image_source else 'image'
        return f"✅ Added image to slide {slide_idx + 1}: {image_name}"
    
    elif tool_name == "add_chart":
        slide_idx = kwargs.get('slide_index', 0)
        chart_type = kwargs.get('chart_type', 'chart')
        categories = kwargs.get('categories', [])
        series_data = kwargs.get('series_data', {})
        return f"✅ Added {chart_type} chart to slide {slide_idx + 1}: {len(categories)} categories, {len(series_data)} series"
    
    elif tool_name == "save_presentation":
        file_path = kwargs.get('file_path', '')
        if file_path:
            # Show both filename and full path for clarity
            file_name = os.path.basename(file_path)
            # Normalize path for display
            display_path = os.path.normpath(file_path)
            
            # Check if it's in Documents folder and mention it prominently
            if "Documents" in display_path:
                return f"✅ Saved presentation: {file_name}\n📁 Location: Documents folder\n📍 Full path: {display_path}"
            else:
                # Truncate very long paths for readability but keep them informative
                if len(display_path) > 80:
                    display_path = f"...{display_path[-77:]}"
                return f"✅ Saved presentation: {file_name}\n📁 Full path: {display_path}"
        return f"✅ Saved presentation → Ready for use!"
    
    elif tool_name == "create_presentation":
        prs_id = kwargs.get('presentation_id', 'new')
        return f"✅ Created presentation {prs_id} → Ready to add slides!"
    
    elif tool_name == "load_presentation":
        prs_id = kwargs.get('presentation_id', 'loaded')
        file_path = kwargs.get('file_path', '')
        file_name = os.path.basename(file_path) if file_path else 'presentation'
        slide_count = kwargs.get('slide_count', 'unknown')
        return f"📂 Loaded presentation {prs_id} from {file_name} → {slide_count} slides available"
    
    elif tool_name == "add_slide":
        slide_idx = kwargs.get('slide_index', 0)
        layout_idx = kwargs.get('layout_index', 6)
        layout_name = kwargs.get('layout_name', f'Layout {layout_idx}')
        total_slides = kwargs.get('total_slides', 'unknown')
        return f"➕ Added slide {slide_idx + 1} using {layout_name} → {total_slides} slides total"
    
    elif tool_name == "extract_text":
        slide_count = kwargs.get('slide_count', 0)
        text_items = kwargs.get('text_items', 0)
        return f"📝 Extracted text from {slide_count} slides → Found {text_items} text items"
    
    elif tool_name == "get_presentation_info":
        slide_count = kwargs.get('slide_count', 0)
        total_shapes = kwargs.get('total_shapes', 0)
        return f"ℹ️ Presentation info: {slide_count} slides, {total_shapes} total shapes"
    
    elif tool_name == "delete_shape":
        slide_idx = kwargs.get('slide_index', 0)
        shape_idx = kwargs.get('shape_index', 0)
        shape_type = kwargs.get('shape_type', 'shape')
        return f"🗑️ Deleted {shape_type} (index {shape_idx}) from slide {slide_idx + 1}"
    
    elif tool_name == "delete_slide":
        slide_idx = kwargs.get('slide_index', 0)
        remaining_slides = kwargs.get('remaining_slides', 'unknown')
        return f"🗑️ Deleted slide {slide_idx + 1} → {remaining_slides} slides remaining"
    
    elif tool_name == "clear_slide":
        slide_idx = kwargs.get('slide_index', 0)
        shapes_cleared = kwargs.get('shapes_cleared', 0)
        return f"🧹 Cleared slide {slide_idx + 1} → Removed {shapes_cleared} shapes"
    
    elif tool_name == "list_slide_content":
        slide_idx = kwargs.get('slide_index', 0)
        shape_count = kwargs.get('shape_count', 0)
        return f"📋 Slide {slide_idx + 1} contents: {shape_count} shapes found"
    
    elif tool_name == "format_existing_text":
        slide_idx = kwargs.get('slide_index', 0)
        shape_idx = kwargs.get('shape_index', 0)
        formatted_props = []
        if kwargs.get('font_size'):
            formatted_props.append(f"size: {kwargs['font_size']}pt")
        if kwargs.get('font_name'):
            formatted_props.append(f"font: {kwargs['font_name']}")
        if kwargs.get('font_color'):
            formatted_props.append(f"color: {kwargs['font_color']}")
        if kwargs.get('text_alignment'):
            formatted_props.append(f"align: {kwargs['text_alignment']}")
        props_desc = ", ".join(formatted_props) if formatted_props else "basic formatting"
        return f"🎨 Updated text formatting for shape {shape_idx} on slide {slide_idx + 1}: {props_desc}"
    
    elif tool_name == "set_slide_background":
        slide_idx = kwargs.get('slide_index', 0)
        bg_color = kwargs.get('background_color')
        bg_image = kwargs.get('background_image')
        if bg_color:
            return f"🎨 Set slide {slide_idx + 1} background color: {bg_color}"
        elif bg_image:
            image_name = os.path.basename(bg_image) if bg_image else 'image'
            return f"🎨 Set slide {slide_idx + 1} background image: {image_name}"
        return f"🎨 Updated slide {slide_idx + 1} background"
    
    # Table-specific success messages
    elif tool_name == "add_table":
        slide_idx = kwargs.get('slide_index', 0)
        rows = kwargs.get('rows', 0)
        cols = kwargs.get('cols', 0)
        header_row = kwargs.get('header_row', False)
        header_note = " (with header)" if header_row else ""
        return f"📊 Added {rows}×{cols} table to slide {slide_idx + 1}{header_note}"
    
    elif tool_name == "set_table_cell":
        slide_idx = kwargs.get('slide_index', 0)
        table_idx = kwargs.get('table_index', 0)
        row = kwargs.get('row', 0)
        col = kwargs.get('col', 0)
        text_preview = kwargs.get('text', '')[:30] + ('...' if len(kwargs.get('text', '')) > 30 else '')
        return f"✅ Updated table {table_idx} cell [{row},{col}] on slide {slide_idx + 1}: \"{text_preview}\""
    
    elif tool_name == "style_table_cell":
        slide_idx = kwargs.get('slide_index', 0)
        table_idx = kwargs.get('table_index', 0)
        row = kwargs.get('row', 0)
        col = kwargs.get('col', 0)
        style_changes = []
        if kwargs.get('fill_color'):
            style_changes.append(f"fill: {kwargs['fill_color']}")
        if kwargs.get('border_color'):
            style_changes.append(f"border: {kwargs['border_color']}")
        style_desc = f" ({', '.join(style_changes)})" if style_changes else ""
        return f"🎨 Styled table {table_idx} cell [{row},{col}] on slide {slide_idx + 1}{style_desc}"
    
    elif tool_name == "style_table_range":
        slide_idx = kwargs.get('slide_index', 0)
        table_idx = kwargs.get('table_index', 0)
        start_row = kwargs.get('start_row', 0)
        start_col = kwargs.get('start_col', 0)
        end_row = kwargs.get('end_row', 0)
        end_col = kwargs.get('end_col', 0)
        cell_count = (end_row - start_row + 1) * (end_col - start_col + 1)
        return f"🎨 Styled table {table_idx} range [{start_row},{start_col}] to [{end_row},{end_col}] on slide {slide_idx + 1} ({cell_count} cells)"
    
    elif tool_name == "modify_table_structure":
        slide_idx = kwargs.get('slide_index', 0)
        table_idx = kwargs.get('table_index', 0)
        operation = kwargs.get('operation', '')
        position = kwargs.get('position', 0)
        count = kwargs.get('count', 1)
        operation_desc = operation.replace('_', ' ')
        count_desc = f" ({count} {'rows' if 'row' in operation else 'columns'})" if count > 1 else ""
        return f"🔧 Table {table_idx} on slide {slide_idx + 1}: {operation_desc} at position {position}{count_desc}"
    
    elif tool_name == "get_table_info":
        slide_idx = kwargs.get('slide_index', 0)
        table_idx = kwargs.get('table_index', 0)
        rows = kwargs.get('rows', 0)
        cols = kwargs.get('cols', 0)
        total_cells = kwargs.get('total_cells', 0)
        return f"ℹ️ Table {table_idx} info on slide {slide_idx + 1}: {rows}×{cols} table with {total_cells} cells"
    
    return f"✅ {tool_name} completed successfully"

# =============================================================================
# SIMPLIFIED POWERPOINT MANAGER
# =============================================================================

class StablePowerPointManager:
    """Simplified PowerPoint manager focused on core functionality"""
    
    def __init__(self):
        self.presentations: Dict[str, Presentation] = {}
        self.temp_files: List[str] = []  # Track temporary files for cleanup
        logger.info("PowerPoint manager initialized")
    
    def create_presentation(self) -> str:
        """Create a new blank presentation"""
        prs = Presentation()
        prs_id = f"ppt_{len(self.presentations)}"
        self.presentations[prs_id] = prs
        logger.info(f"Created presentation: {prs_id}")
        return prs_id
    
    def load_presentation(self, file_path: str) -> str:
        """Load an existing PowerPoint presentation from file"""
        # Ensure .pptx extension if not provided
        if not file_path.lower().endswith('.pptx'):
            file_path += '.pptx'
        
        # Convert to absolute path for consistency
        if not os.path.isabs(file_path):
            file_path = os.path.abspath(file_path)
        
        # Check if file exists
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Presentation file not found: {file_path}")
        
        try:
            # Load the presentation
            prs = Presentation(file_path)
            prs_id = f"ppt_{len(self.presentations)}"
            self.presentations[prs_id] = prs
            
            logger.info(f"Loaded presentation: {prs_id} from {file_path}")
            return prs_id
            
        except Exception as e:
            logger.error(f"Failed to load presentation: {e}")
            raise RuntimeError(f"Failed to load presentation from {file_path}: {e}")
    
    def add_slide(self, prs_id: str, layout_index: int = 6) -> int:
        """Add a new slide to the presentation with specified layout"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        # Validate layout index
        if layout_index < 0 or layout_index >= len(prs.slide_layouts):
            available_layouts = len(prs.slide_layouts)
            raise ValueError(f"Layout index {layout_index} is invalid. Available layouts: 0-{available_layouts-1}")
        
        # Add the slide
        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)
        slide_index = len(prs.slides) - 1
        
        logger.info(f"Added slide {slide_index} with layout {layout_index} to {prs_id}")
        return slide_index
    
    def add_text_box(self, prs_id: str, slide_index: int, text: str, 
                     left: float = 1, top: float = 1, width: float = 8, height: float = 1,
                     font_size: int = 18, font_name: str = "Calibri", font_color: Optional[str] = None,
                     bold: bool = False, italic: bool = False, underline: bool = False,
                     text_alignment: str = "left", fill_color: Optional[str] = None,
                     border_color: Optional[str] = None, border_width: float = 0) -> bool:
        """Add a text box to a slide with comprehensive formatting"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        # Add slide if needed
        while len(prs.slides) <= slide_index:
            layout = prs.slide_layouts[6]  # Blank layout
            prs.slides.add_slide(layout)
        
        slide = prs.slides[slide_index]
        
        # Add text box
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        
        # Map text alignment
        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY
        }
        paragraph_alignment = alignment_map.get(text_alignment.lower(), PP_ALIGN.LEFT)
        
        # Apply text formatting
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = paragraph_alignment
            for run in paragraph.runs:
                # Font properties
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.italic = italic
                run.font.underline = underline
                
                # Font color
                if font_color:
                    try:
                        rgb = self._parse_color(font_color)
                        run.font.color.rgb = RGBColor(*rgb)
                    except Exception as e:
                        logger.warning(f"Invalid font color '{font_color}': {e}")
        
        # Apply shape formatting
        try:
            # Fill color
            if fill_color:
                try:
                    rgb = self._parse_color(fill_color)
                    textbox.fill.solid()
                    textbox.fill.fore_color.rgb = RGBColor(*rgb)
                except Exception as e:
                    logger.warning(f"Invalid fill color '{fill_color}': {e}")
            
            # Border formatting
            if border_width > 0:
                textbox.line.width = Pt(border_width)
                if border_color:
                    try:
                        rgb = self._parse_color(border_color)
                        textbox.line.color.rgb = RGBColor(*rgb)
                    except Exception as e:
                        logger.warning(f"Invalid border color '{border_color}': {e}")
            else:
                # No border
                textbox.line.fill.background()
                
        except Exception as e:
            logger.warning(f"Failed to apply shape formatting: {e}")
        
        logger.info(f"Added formatted text box to slide {slide_index}")
        return True
    
    def _parse_color(self, color_str: str) -> tuple:
        """Parse color string to RGB tuple. Supports hex (#RRGGBB) and RGB (r,g,b) formats"""
        color_str = color_str.strip()
        
        # Hex format: #RRGGBB or RRGGBB
        if color_str.startswith('#'):
            color_str = color_str[1:]
        
        if len(color_str) == 6:
            try:
                r = int(color_str[0:2], 16)
                g = int(color_str[2:4], 16) 
                b = int(color_str[4:6], 16)
                return (r, g, b)
            except ValueError:
                pass
        
        # RGB format: "r,g,b" or "(r,g,b)"
        if ',' in color_str:
            color_str = color_str.strip('()')
            try:
                parts = [int(x.strip()) for x in color_str.split(',')]
                if len(parts) == 3 and all(0 <= x <= 255 for x in parts):
                    return tuple(parts)
            except ValueError:
                pass
        
        # Predefined colors
        color_map = {
            'black': (0, 0, 0),
            'white': (255, 255, 255),
            'red': (255, 0, 0),
            'darkred': (139, 0, 0),
            'green': (0, 128, 0),
            'darkgreen': (0, 100, 0),
            'blue': (0, 0, 255),
            'darkblue': (0, 0, 139),
            'yellow': (255, 255, 0),
            'orange': (255, 165, 0),
            'purple': (128, 0, 128),
            'gray': (128, 128, 128),
            'grey': (128, 128, 128),
            'lightgray': (211, 211, 211),
            'lightgrey': (211, 211, 211),
            'darkgray': (64, 64, 64),
            'darkgrey': (64, 64, 64)
        }
        
        if color_str.lower() in color_map:
            return color_map[color_str.lower()]
        
        raise ValueError(f"Invalid color format: {color_str}. Use hex (#RRGGBB), RGB (r,g,b), or predefined color names")
    
    def format_existing_text(self, prs_id: str, slide_index: int, shape_index: int,
                           font_size: Optional[int] = None, font_name: Optional[str] = None,
                           font_color: Optional[str] = None, bold: Optional[bool] = None,
                           italic: Optional[bool] = None, underline: Optional[bool] = None,
                           text_alignment: Optional[str] = None) -> bool:
        """Modify formatting of existing text shape"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide {slide_index} does not exist")
        
        slide = prs.slides[slide_index]
        
        if shape_index >= len(slide.shapes):
            raise ValueError(f"Shape {shape_index} does not exist")
        
        shape = slide.shapes[shape_index]
        
        # Check if it's a text shape
        if not hasattr(shape, 'text_frame'):
            raise ValueError(f"Shape {shape_index} is not a text shape")
        
        text_frame = shape.text_frame
        
        # Apply text alignment if specified
        if text_alignment:
            alignment_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY
            }
            paragraph_alignment = alignment_map.get(text_alignment.lower(), PP_ALIGN.LEFT)
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = paragraph_alignment
        
        # Apply text formatting
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if font_name is not None:
                    run.font.name = font_name
                if font_size is not None:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
                if italic is not None:
                    run.font.italic = italic
                if underline is not None:
                    run.font.underline = underline
                
                if font_color:
                    try:
                        rgb = self._parse_color(font_color)
                        run.font.color.rgb = RGBColor(*rgb)
                    except Exception as e:
                        logger.warning(f"Invalid font color '{font_color}': {e}")
        
        logger.info(f"Updated formatting for text shape {shape_index} on slide {slide_index}")
        return True
    
    def set_slide_background(self, prs_id: str, slide_index: int, 
                           background_color: Optional[str] = None,
                           background_image: Optional[str] = None) -> bool:
        """Set slide background color or image"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide {slide_index} does not exist")
        
        slide = prs.slides[slide_index]
        
        try:
            if background_color:
                # Set background color
                background = slide.background
                fill = background.fill
                fill.solid()
                rgb = self._parse_color(background_color)
                fill.fore_color.rgb = RGBColor(*rgb)
                logger.info(f"Set slide {slide_index} background color to {background_color}")
            
            if background_image:
                # Set background image
                if background_image.startswith(('http://', 'https://')):
                    # Download image temporarily
                    import urllib.request
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                        urllib.request.urlretrieve(background_image, tmp_file.name)
                        image_path = tmp_file.name
                else:
                    image_path = background_image
                    if not os.path.exists(image_path):
                        raise FileNotFoundError(f"Background image not found: {image_path}")
                
                # Apply background image
                background = slide.background
                fill = background.fill
                fill.picture(image_path)
                
                # Clean up temporary file if downloaded
                if background_image.startswith(('http://', 'https://')):
                    os.unlink(image_path)
                
                logger.info(f"Set slide {slide_index} background image")
            
            return True
            
        except Exception as e:
            logger.error(f"Failed to set slide background: {e}")
            raise RuntimeError(f"Failed to set slide background: {e}")
    
    def add_image(self, prs_id: str, slide_index: int, image_source: str,
                  left: float = 1, top: float = 1, width: Optional[float] = None, 
                  height: Optional[float] = None) -> bool:
        """Add an image to a slide"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        # Add slide if needed
        while len(prs.slides) <= slide_index:
            layout = prs.slide_layouts[6]
            prs.slides.add_slide(layout)
        
        slide = prs.slides[slide_index]
        
        try:
            # Handle different image sources
            if image_source.startswith(('http://', 'https://')):
                # Download image temporarily
                import urllib.request
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
                    urllib.request.urlretrieve(image_source, tmp_file.name)
                    image_path = tmp_file.name
            else:
                # Local file
                image_path = image_source
                if not os.path.exists(image_path):
                    raise FileNotFoundError(f"Image file not found: {image_path}")
            
            # Add image to slide
            if width and height:
                picture = slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top), Inches(width), Inches(height)
                )
            else:
                picture = slide.shapes.add_picture(
                    image_path, Inches(left), Inches(top)
                )
            
            # Clean up temporary file if downloaded
            if image_source.startswith(('http://', 'https://')):
                os.unlink(image_path)
            
            logger.info(f"Added image to slide {slide_index}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to add image: {e}")
            raise
    
    def add_chart(self, prs_id: str, slide_index: int, chart_type: str, 
                  categories: List[str], series_data: Dict[str, List[float]],
                  left: float = 2, top: float = 2, width: float = 6, height: float = 4.5) -> bool:
        """Add a chart to a slide"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        # Add slide if needed
        while len(prs.slides) <= slide_index:
            layout = prs.slide_layouts[6]
            prs.slides.add_slide(layout)
        
        slide = prs.slides[slide_index]
        
        try:
            # Create chart data
            chart_data = CategoryChartData()
            chart_data.categories = categories
            
            for series_name, values in series_data.items():
                if len(values) != len(categories):
                    raise ValueError(f"Series '{series_name}' has {len(values)} values but {len(categories)} categories")
                chart_data.add_series(series_name, values)
            
            # Map chart type
            chart_type_map = {
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
                "area": XL_CHART_TYPE.AREA
            }
            
            xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Add chart to slide
            chart = slide.shapes.add_chart(
                xl_chart_type, Inches(left), Inches(top), 
                Inches(width), Inches(height), chart_data
            ).chart
            
            # Basic chart formatting
            chart.has_legend = True
            chart.legend.position = 2  # Right
            
            logger.info(f"Added {chart_type} chart to slide {slide_index}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to add chart: {e}")
            raise
    
    def save_presentation(self, prs_id: str, file_path: str) -> str:
        """Save presentation to file and return file info - handles Windows paths properly"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        # Ensure .pptx extension first (before path processing)
        if not file_path.lower().endswith('.pptx'):
            file_path += '.pptx'
        
        # Handle relative paths intelligently
        if not os.path.isabs(file_path):
            try:
                cwd = os.getcwd()
                logger.info(f"Current working directory: {cwd}")
                
                # Check if we're in a system/application directory
                problematic_paths = ['AppData', 'cursor', 'Program Files', 'Windows']
                is_system_dir = any(path_part in cwd for path_part in problematic_paths)
                
                if is_system_dir:
                    # Use user's Documents folder for better accessibility
                    documents_dir = os.path.join(os.path.expanduser("~"), "Documents")
                    file_path = os.path.join(documents_dir, file_path)
                    logger.info(f"Using Documents directory instead of system directory: {documents_dir}")
                else:
                    # Use the current working directory
                    file_path = os.path.join(cwd, file_path)
                    logger.info(f"Using current working directory: {cwd}")
                    
            except Exception as e:
                # Fallback to Documents folder
                logger.warning(f"Error determining working directory: {e}")
                documents_dir = os.path.join(os.path.expanduser("~"), "Documents")
                file_path = os.path.join(documents_dir, file_path)
                logger.info(f"Fallback to Documents directory: {documents_dir}")
        
        # Normalize path for Windows (handles both / and \ separators)
        file_path = os.path.normpath(file_path)
        
        # Create directory if needed (always try to create parent directory)
        dir_path = os.path.dirname(file_path)
        if dir_path and not os.path.exists(dir_path):
            try:
                os.makedirs(dir_path, exist_ok=True)
                logger.info(f"Created directory: {dir_path}")
            except OSError as e:
                raise RuntimeError(f"Failed to create directory {dir_path}: {e}")
        
        # Save presentation
        try:
            self.presentations[prs_id].save(file_path)
            logger.info(f"PowerPoint saved to: {file_path}")
        except Exception as e:
            logger.error(f"Save failed: {e}")
            raise RuntimeError(f"Failed to save PowerPoint file to {file_path}: {e}")
        
        # Verify the file was created
        if os.path.exists(file_path):
            file_size = os.path.getsize(file_path)
            logger.info(f"Saved {prs_id} to {file_path} ({file_size} bytes)")
            return file_path
        else:
            raise RuntimeError(f"File was not created at {file_path}")
    
    def screenshot_slides(self, file_path: str, output_dir: Optional[str] = None, 
                         image_format: str = "PNG", width: int = 1920, height: int = 1080) -> List[str]:
        """Screenshot each slide of a PowerPoint presentation (Windows only)
        
        Args:
            file_path: Path to the PowerPoint file
            output_dir: Directory to save screenshots (defaults to temp directory)
            image_format: Image format (PNG, JPG, etc.)
            width: Screenshot width in pixels
            height: Screenshot height in pixels
            
        Returns:
            List of paths to the generated screenshot files
        """
        if not WIN32_COM_AVAILABLE:
            raise RuntimeError("Screenshot feature is only available on Windows with win32com installed")
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            # Initialize COM
            pythoncom.CoInitialize()
            
            # Create PowerPoint application instance
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = True  # Make visible for screenshot
            
            # Open the presentation
            presentation = ppt_app.Presentations.Open(os.path.abspath(file_path))
            
            # Set up output directory
            if output_dir is None:
                output_dir = tempfile.mkdtemp(prefix="ppt_screenshots_")
            else:
                os.makedirs(output_dir, exist_ok=True)
            
            screenshot_paths = []
            
            # Export each slide as image
            for i, slide in enumerate(presentation.Slides):
                slide_num = i + 1
                output_file = os.path.join(output_dir, f"slide_{slide_num:03d}.{image_format.lower()}")
                
                # Export slide as image
                slide.Export(output_file, image_format, width, height)
                screenshot_paths.append(output_file)
                
                logger.info(f"Exported slide {slide_num} to {output_file}")
            
            # Close presentation and quit PowerPoint
            presentation.Close()
            ppt_app.Quit()
            
            # Add to temp files for cleanup if using temp directory
            if output_dir.startswith(tempfile.gettempdir()):
                self.temp_files.extend(screenshot_paths)
                self.temp_files.append(output_dir)
            
            logger.info(f"Successfully created {len(screenshot_paths)} slide screenshots")
            return screenshot_paths
            
        except Exception as e:
            logger.error(f"Error creating slide screenshots: {e}")
            # Try to cleanup COM objects
            try:
                if 'presentation' in locals():
                    presentation.Close()
                if 'ppt_app' in locals():
                    ppt_app.Quit()
            except:
                pass
            finally:
                pythoncom.CoUninitialize()
            raise
        finally:
            # Ensure COM is uninitialized
            try:
                pythoncom.CoUninitialize()
            except:
                pass
    
    async def screenshot_slides_async(self, file_path: str, output_dir: Optional[str] = None, 
                                    image_format: str = "PNG", width: int = 1920, height: int = 1080) -> List[str]:
        """Take screenshots asynchronously to prevent blocking the event loop"""
        return await asyncio.to_thread(self.screenshot_slides, file_path, output_dir, image_format, width, height)
    
    def cleanup(self):
        """Clean up temporary files and resources"""
        for temp_file in self.temp_files:
            try:
                if os.path.isfile(temp_file):
                    os.remove(temp_file)
                elif os.path.isdir(temp_file):
                    import shutil
                    shutil.rmtree(temp_file)
            except Exception as e:
                logger.warning(f"Could not clean up {temp_file}: {e}")
        self.temp_files.clear()
        logger.info("Cleanup completed")
    
    def delete_shape(self, prs_id: str, slide_index: int, shape_index: int) -> bool:
        """Delete a specific shape from a slide by index"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide {slide_index} does not exist (presentation has {len(prs.slides)} slides)")
        
        slide = prs.slides[slide_index]
        
        if shape_index >= len(slide.shapes):
            raise ValueError(f"Shape {shape_index} does not exist (slide has {len(slide.shapes)} shapes)")
        
        # Get shape info for logging before deletion
        shape = slide.shapes[shape_index]
        shape_type = "unknown"
        try:
            if hasattr(shape, 'text_frame'):
                shape_type = "text box"
            elif hasattr(shape, 'chart'):
                shape_type = "chart"
            elif hasattr(shape, 'image'):
                shape_type = "image"
        except:
            pass
        
        # Delete the shape
        shape_element = shape.element
        shape_element.getparent().remove(shape_element)
        
        logger.info(f"Deleted {shape_type} (index {shape_index}) from slide {slide_index}")
        return True
    
    def delete_slide(self, prs_id: str, slide_index: int) -> bool:
        """Delete an entire slide from the presentation"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        if len(prs.slides) <= 1:
            raise ValueError("Cannot delete slide - presentation must have at least one slide")
        
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide {slide_index} does not exist (presentation has {len(prs.slides)} slides)")
        
        # Remove the slide
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[slide_index])
        
        logger.info(f"Deleted slide {slide_index} from presentation {prs_id}")
        return True
    
    def clear_slide(self, prs_id: str, slide_index: int) -> bool:
        """Clear all content from a slide but keep the slide"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide {slide_index} does not exist (presentation has {len(prs.slides)} slides)")
        
        slide = prs.slides[slide_index]
        
        # Count shapes before deletion
        shape_count = len(slide.shapes)
        
        # Delete all shapes (in reverse order to avoid index issues)
        for i in range(len(slide.shapes) - 1, -1, -1):
            try:
                shape = slide.shapes[i]
                shape_element = shape.element
                shape_element.getparent().remove(shape_element)
            except Exception as e:
                logger.warning(f"Could not delete shape {i}: {e}")
        
        logger.info(f"Cleared {shape_count} shapes from slide {slide_index}")
        return True
    
    def list_slide_content(self, prs_id: str, slide_index: int) -> Dict[str, Any]:
        """List all content on a slide for easier deletion targeting"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide {slide_index} does not exist (presentation has {len(prs.slides)} slides)")
        
        slide = prs.slides[slide_index]
        content = []
        
        for i, shape in enumerate(slide.shapes):
            shape_info = {
                "index": i,
                "type": "unknown",
                "description": ""
            }
            
            try:
                if hasattr(shape, 'text_frame') and shape.text_frame.text:
                    shape_info["type"] = "text"
                    shape_info["description"] = f"Text: '{shape.text_frame.text[:50]}...'" if len(shape.text_frame.text) > 50 else f"Text: '{shape.text_frame.text}'"
                elif hasattr(shape, 'table'):
                    shape_info["type"] = "table"
                    table = shape.table
                    rows = len(table.rows)
                    cols = len(table.columns)
                    # Get a preview of table content
                    try:
                        first_cell = table.cell(0, 0).text.strip() if rows > 0 and cols > 0 else ""
                        preview = f"'{first_cell[:20]}...'" if len(first_cell) > 20 and first_cell else "(empty)"
                        shape_info["description"] = f"Table ({rows}×{cols}) - {preview}"
                    except:
                        shape_info["description"] = f"Table ({rows}×{cols})"
                elif hasattr(shape, 'chart'):
                    shape_info["type"] = "chart"
                    shape_info["description"] = "Chart"
                elif hasattr(shape, 'image'):
                    shape_info["type"] = "image"
                    shape_info["description"] = "Image"
                else:
                    shape_info["type"] = "shape"
                    shape_info["description"] = "Shape"
            except:
                pass
            
            content.append(shape_info)
        
        return {
            "slide_index": slide_index,
            "shape_count": len(content),
            "shapes": content
        }
    
    def extract_text(self, prs_id: str) -> List[Dict[str, Any]]:
        """Extract all text content from the presentation"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        extracted_text = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = {
                "slide_index": slide_idx,
                "slide_number": slide_idx + 1,
                "text_content": []
            }
            
            for shape_idx, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                        shape_text = {
                            "shape_index": shape_idx,
                            "shape_type": "text",
                            "text": shape.text_frame.text.strip()
                        }
                        slide_text["text_content"].append(shape_text)
                    elif hasattr(shape, 'table'):
                        # Extract text from table cells using enhanced method
                        table_text = self._extract_table_text(shape.table, shape_idx)
                        if table_text:
                            slide_text["text_content"].append(table_text)
                except Exception as e:
                    logger.warning(f"Could not extract text from shape {shape_idx} on slide {slide_idx}: {e}")
            
            extracted_text.append(slide_text)
        
        logger.info(f"Extracted text from {len(extracted_text)} slides in {prs_id}")
        return extracted_text

    def get_presentation_info(self, prs_id: str) -> Dict[str, Any]:
        """Get comprehensive presentation information"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        # Count different types of content
        total_text_boxes = 0
        total_images = 0
        total_charts = 0
        total_shapes = 0
        slide_details = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_info = {
                "slide_index": slide_idx,
                "slide_number": slide_idx + 1,
                "shape_count": len(slide.shapes),
                "has_text": False,
                "has_images": False,
                "has_charts": False
            }
            
            for shape in slide.shapes:
                total_shapes += 1
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                        total_text_boxes += 1
                        slide_info["has_text"] = True
                    elif hasattr(shape, 'chart'):
                        total_charts += 1
                        slide_info["has_charts"] = True
                    elif hasattr(shape, 'image'):
                        total_images += 1
                        slide_info["has_images"] = True
                except:
                    pass
            
            slide_details.append(slide_info)
        
        # Get available slide layouts
        available_layouts = []
        for i, layout in enumerate(prs.slide_layouts):
            try:
                layout_name = layout.name if hasattr(layout, 'name') else f"Layout {i}"
                available_layouts.append({
                    "index": i,
                    "name": layout_name
                })
            except:
                available_layouts.append({
                    "index": i,
                    "name": f"Layout {i}"
                })
        
        return {
            "presentation_id": prs_id,
            "slide_count": len(prs.slides),
            "total_shapes": total_shapes,
            "content_summary": {
                "text_boxes": total_text_boxes,
                "images": total_images, 
                "charts": total_charts,
                "other_shapes": total_shapes - total_text_boxes - total_images - total_charts
            },
            "available_layouts": available_layouts,
            "slide_details": slide_details,
            "status": "ready"
        }
    
    def _post_process_slide(self, prs_id: str, slide_index: int):
        """Basic post-processing to fix common issues"""
        try:
            prs = self.presentations[prs_id]
            slide = prs.slides[slide_index]
            
            # Fix green rectangle fills on placeholder shapes
            for shape in slide.shapes:
                if hasattr(shape, 'fill'):
                    try:
                        if hasattr(shape.fill, 'solid'):
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
                    except:
                        pass  # Skip if can't modify
                        
        except Exception as e:
            logger.warning(f"Post-processing failed for slide {slide_index}: {e}")
    
    # =============================================================================
    # TABLE OPERATIONS - Phase 1: Foundation
    # =============================================================================
    
    def _get_table_shape(self, prs_id: str, slide_index: int, table_index: int):
        """Get table shape object with validation"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        if slide_index >= len(prs.slides):
            raise ValueError(f"Slide {slide_index} does not exist")
        
        slide = prs.slides[slide_index]
        tables = [shape for shape in slide.shapes if hasattr(shape, 'table')]
        
        if table_index >= len(tables):
            raise ValueError(f"Table {table_index} does not exist (found {len(tables)} tables)")
        
        return tables[table_index]
    
    def _get_table(self, prs_id: str, slide_index: int, table_index: int):
        """Get table object with validation"""
        return self._get_table_shape(prs_id, slide_index, table_index).table
    
    def add_table(self, prs_id: str, slide_index: int, rows: int, cols: int,
                  left: float = 1, top: float = 1, width: float = 8, height: float = 4,
                  header_row: bool = False) -> int:
        """Add a table to a slide and return table index"""
        if prs_id not in self.presentations:
            raise ValueError(f"Presentation {prs_id} not found")
        
        prs = self.presentations[prs_id]
        
        # Add slide if needed
        while len(prs.slides) <= slide_index:
            layout = prs.slide_layouts[6]
            prs.slides.add_slide(layout)
        
        slide = prs.slides[slide_index]
        
        try:
            # Add table
            table_shape = slide.shapes.add_table(
                rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
            )
            
            table = table_shape.table
            
            # Style header row if requested
            if header_row and rows > 0:
                for col_idx in range(cols):
                    cell = table.cell(0, col_idx)
                    # Make header bold
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                    # Add header background - blue header
                    try:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(79, 129, 189)  # Blue header
                        # Set text color to white for contrast
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(255, 255, 255)  # White text
                    except Exception as e:
                        logger.warning(f"Failed to style header row: {e}")
            
            # Return the index of the newly created table
            table_index = len([shape for shape in slide.shapes if hasattr(shape, 'table')]) - 1
            
            logger.info(f"Added {rows}×{cols} table to slide {slide_index}")
            return table_index
            
        except Exception as e:
            logger.error(f"Failed to add table: {e}")
            raise RuntimeError(f"Failed to add table to slide {slide_index}: {e}")
    
    def set_table_cell(self, prs_id: str, slide_index: int, table_index: int,
                       row: int, col: int, text: str,
                       font_size: Optional[int] = None, font_name: Optional[str] = None,
                       font_color: Optional[str] = None, bold: Optional[bool] = None,
                       italic: Optional[bool] = None, underline: Optional[bool] = None,
                       text_alignment: Optional[str] = None) -> bool:
        """Set cell content and formatting"""
        table = self._get_table(prs_id, slide_index, table_index)
        
        # Validate cell coordinates
        if row >= len(table.rows) or col >= len(table.columns):
            raise ValueError(f"Cell [{row},{col}] is out of bounds for table with {len(table.rows)} rows and {len(table.columns)} columns")
        
        try:
            cell = table.cell(row, col)
            cell.text = text
            
            # Apply text formatting if specified
            if any([font_size, font_name, font_color, bold, italic, underline, text_alignment]):
                text_frame = cell.text_frame
                
                # Map text alignment
                if text_alignment:
                    alignment_map = {
                        "left": PP_ALIGN.LEFT,
                        "center": PP_ALIGN.CENTER,
                        "right": PP_ALIGN.RIGHT,
                        "justify": PP_ALIGN.JUSTIFY
                    }
                    paragraph_alignment = alignment_map.get(text_alignment.lower(), PP_ALIGN.LEFT)
                    for paragraph in text_frame.paragraphs:
                        paragraph.alignment = paragraph_alignment
                
                # Apply formatting to all runs
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if font_name:
                            run.font.name = font_name
                        if font_size:
                            run.font.size = Pt(font_size)
                        if bold is not None:
                            run.font.bold = bold
                        if italic is not None:
                            run.font.italic = italic
                        if underline is not None:
                            run.font.underline = underline
                        if font_color:
                            rgb = self._parse_color(font_color)
                            run.font.color.rgb = RGBColor(*rgb)
            
            logger.info(f"Set table cell [{row},{col}] content and formatting")
            return True
            
        except Exception as e:
            logger.error(f"Failed to set table cell: {e}")
            raise RuntimeError(f"Failed to set cell [{row},{col}] in table {table_index}: {e}")
    
    def get_table_info(self, prs_id: str, slide_index: int, table_index: int) -> Dict[str, Any]:
        """Get comprehensive table information"""
        table = self._get_table(prs_id, slide_index, table_index)
        
        try:
            # Extract cell contents
            cell_data = []
            for row_idx, row in enumerate(table.rows):
                row_data = []
                for col_idx, cell in enumerate(row.cells):
                    row_data.append({
                        "text": cell.text.strip(),
                        "row": row_idx,
                        "col": col_idx
                    })
                cell_data.append(row_data)
            
            return {
                "table_index": table_index,
                "rows": len(table.rows),
                "columns": len(table.columns),
                "cell_data": cell_data,
                "total_cells": len(table.rows) * len(table.columns)
            }
            
        except Exception as e:
            logger.error(f"Failed to get table info: {e}")
            raise RuntimeError(f"Failed to get info for table {table_index}: {e}")
    
    def _extract_table_text(self, table, shape_idx):
        """Extract text content from table cells for enhanced text extraction"""
        try:
            table_content = []
            for row_idx, row in enumerate(table.rows):
                row_content = []
                for col_idx, cell in enumerate(row.cells):
                    if cell.text.strip():
                        row_content.append(cell.text.strip())
                if row_content:
                    table_content.append(" | ".join(row_content))
            
            if table_content:
                return {
                    "shape_index": shape_idx,
                    "shape_type": "table",
                    "text": "\n".join(table_content),
                    "rows": len(table.rows),
                    "columns": len(table.columns)
                }
            return None
        except Exception as e:
            logger.warning(f"Failed to extract table text: {e}")
            return None
    
    # =============================================================================
    # TABLE OPERATIONS - Phase 2: Advanced Styling
    # =============================================================================
    
    def style_table_cell(self, prs_id: str, slide_index: int, table_index: int,
                         row: int, col: int, fill_color: Optional[str] = None,
                         border_color: Optional[str] = None, border_width: Optional[float] = None,
                         margin_left: Optional[float] = None, margin_right: Optional[float] = None,
                         margin_top: Optional[float] = None, margin_bottom: Optional[float] = None) -> bool:
        """Apply styling to a table cell (background, borders, margins)"""
        table = self._get_table(prs_id, slide_index, table_index)
        
        # Validate cell coordinates
        if row >= len(table.rows) or col >= len(table.columns):
            raise ValueError(f"Cell [{row},{col}] is out of bounds for table with {len(table.rows)} rows and {len(table.columns)} columns")
        
        try:
            cell = table.cell(row, col)
            
            # Apply fill color
            if fill_color:
                rgb = self._parse_color(fill_color)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*rgb)
                logger.info(f"Applied fill color {fill_color} to cell [{row},{col}]")
            
            # Apply margins
            margin_applied = False
            if margin_left is not None:
                cell.margin_left = Inches(margin_left)
                margin_applied = True
            if margin_right is not None:
                cell.margin_right = Inches(margin_right)
                margin_applied = True
            if margin_top is not None:
                cell.margin_top = Inches(margin_top)
                margin_applied = True
            if margin_bottom is not None:
                cell.margin_bottom = Inches(margin_bottom)
                margin_applied = True
            
            if margin_applied:
                logger.info(f"Applied margins to cell [{row},{col}]")
            
            # Apply borders (simplified approach - python-pptx has limited border support)
            if border_color and border_width:
                try:
                    # Note: python-pptx border handling is complex and limited
                    # This is a simplified implementation that may not work perfectly
                    rgb = self._parse_color(border_color)
                    # Set border on the cell (this may not work as expected due to python-pptx limitations)
                    logger.info(f"Attempted to apply border to cell [{row},{col}] - limited support in python-pptx")
                except Exception as e:
                    logger.warning(f"Border styling not fully supported: {e}")
            
            logger.info(f"Applied styling to table cell [{row},{col}]")
            return True
            
        except Exception as e:
            logger.error(f"Failed to style table cell: {e}")
            raise RuntimeError(f"Failed to style cell [{row},{col}] in table {table_index}: {e}")
    
    def style_table_range(self, prs_id: str, slide_index: int, table_index: int,
                         start_row: int, start_col: int, end_row: int, end_col: int,
                         fill_color: Optional[str] = None, border_color: Optional[str] = None,
                         border_width: Optional[float] = None, margin_left: Optional[float] = None,
                         margin_right: Optional[float] = None, margin_top: Optional[float] = None,
                         margin_bottom: Optional[float] = None) -> bool:
        """Apply styling to a range of cells"""
        table = self._get_table(prs_id, slide_index, table_index)
        
        # Validate range
        if start_row > end_row or start_col > end_col:
            raise ValueError("Invalid range: start coordinates must be <= end coordinates")
        
        if end_row >= len(table.rows) or end_col >= len(table.columns):
            raise ValueError(f"Range end [{end_row},{end_col}] is out of bounds for table with {len(table.rows)} rows and {len(table.columns)} columns")
        
        try:
            cells_styled = 0
            for row in range(start_row, end_row + 1):
                for col in range(start_col, end_col + 1):
                    self.style_table_cell(
                        prs_id, slide_index, table_index, row, col,
                        fill_color, border_color, border_width,
                        margin_left, margin_right, margin_top, margin_bottom
                    )
                    cells_styled += 1
            
            logger.info(f"Applied styling to {cells_styled} cells in range [{start_row},{start_col}] to [{end_row},{end_col}]")
            return True
            
        except Exception as e:
            logger.error(f"Failed to style table range: {e}")
            raise RuntimeError(f"Failed to style range [{start_row},{start_col}] to [{end_row},{end_col}] in table {table_index}: {e}")
    
    def create_table_with_data(self, prs_id: str, slide_index: int, 
                              table_data: List[List[str]], headers: Optional[List[str]] = None,
                              left: float = 1, top: float = 1, width: float = 8, height: float = 4,
                              header_style: Optional[Dict[str, Any]] = None,
                              data_style: Optional[Dict[str, Any]] = None,
                              alternating_rows: bool = False) -> int:
        """Create a table with data and optional styling (convenience method)"""
        if not table_data or not table_data[0]:
            raise ValueError("table_data must be a non-empty list of lists")
        
        # Determine table dimensions
        rows = len(table_data)
        cols = len(table_data[0])
        
        # Add header row if provided
        if headers:
            if len(headers) != cols:
                raise ValueError(f"Headers length ({len(headers)}) must match table columns ({cols})")
            rows += 1
        
        # Create the table
        table_index = self.add_table(
            prs_id, slide_index, rows, cols, left, top, width, height, header_row=bool(headers)
        )
        
        try:
            current_row = 0
            
            # Set headers if provided
            if headers:
                for col_idx, header_text in enumerate(headers):
                    self.set_table_cell(
                        prs_id, slide_index, table_index, current_row, col_idx, header_text,
                        **(header_style or {})
                    )
                current_row += 1
            
            # Set data
            for row_data in table_data:
                if len(row_data) != cols:
                    raise ValueError(f"All data rows must have {cols} columns")
                
                for col_idx, cell_text in enumerate(row_data):
                    # Determine cell style
                    cell_style = data_style or {}
                    
                    # Apply alternating row styling
                    if alternating_rows and current_row % 2 == 1:
                        # Odd rows get light gray background
                        if 'fill_color' not in cell_style:
                            self.style_table_cell(
                                prs_id, slide_index, table_index, current_row, col_idx,
                                fill_color="#F2F2F2"
                            )
                    
                    self.set_table_cell(
                        prs_id, slide_index, table_index, current_row, col_idx, str(cell_text),
                        **cell_style
                    )
                
                current_row += 1
            
            logger.info(f"Created table with data: {rows}×{cols} table with {len(table_data)} data rows")
            return table_index
            
        except Exception as e:
            logger.error(f"Failed to create table with data: {e}")
            raise RuntimeError(f"Failed to populate table with data: {e}")
    
    def modify_table_structure(self, prs_id: str, slide_index: int, table_index: int,
                               operation: str, position: Optional[int] = None, count: int = 1) -> bool:
        """
        Modify table structure using a workaround approach.
        
        NOTE: python-pptx library does not natively support adding/removing rows and columns
        from existing tables. This method provides a workaround by creating a new table
        with the desired structure and copying content from the original table.
        
        Args:
            prs_id: Presentation ID
            slide_index: Slide index (0-based)
            table_index: Table index (0-based)
            operation: Operation type - "add_row", "remove_row", "add_column", "remove_column"
            position: Position to insert/remove at (None = end for add, last for remove)
            count: Number of rows/columns to add/remove (default: 1)
        
        Returns:
            bool: True if operation succeeded
        """
        try:
            # Get the original table
            table_shape = self._get_table_shape(prs_id, slide_index, table_index)
            table = table_shape.table
            
            # Current dimensions
            current_rows = len(table.rows)
            current_cols = len(table.columns)
            
            # Calculate new dimensions
            if operation == "add_row":
                new_rows = current_rows + count
                new_cols = current_cols
                if position is None:
                    position = current_rows
                elif position < 0 or position > current_rows:
                    raise ValueError(f"Position {position} invalid for add_row (valid range: 0-{current_rows})")
                    
            elif operation == "remove_row":
                if count >= current_rows:
                    raise ValueError(f"Cannot remove {count} rows from table with only {current_rows} rows")
                new_rows = current_rows - count
                new_cols = current_cols
                if position is None:
                    position = current_rows - count
                elif position < 0 or position + count > current_rows:
                    raise ValueError(f"Position {position} invalid for remove_row with count {count}")
                    
            elif operation == "add_column":
                new_rows = current_rows
                new_cols = current_cols + count
                if position is None:
                    position = current_cols
                elif position < 0 or position > current_cols:
                    raise ValueError(f"Position {position} invalid for add_column (valid range: 0-{current_cols})")
                    
            elif operation == "remove_column":
                if count >= current_cols:
                    raise ValueError(f"Cannot remove {count} columns from table with only {current_cols} columns")
                new_rows = current_rows
                new_cols = current_cols - count
                if position is None:
                    position = current_cols - count
                elif position < 0 or position + count > current_cols:
                    raise ValueError(f"Position {position} invalid for remove_column with count {count}")
                    
            else:
                raise ValueError(f"Unknown operation: {operation}")
            
            # Extract all current cell data
            cell_data = []
            for row_idx in range(current_rows):
                row_data = []
                for col_idx in range(current_cols):
                    cell = table.cell(row_idx, col_idx)
                    row_data.append(cell.text)
                cell_data.append(row_data)
            
            # Get table position and size
            left = table_shape.left
            top = table_shape.top
            width = table_shape.width
            height = table_shape.height
            
            # Get slide reference
            prs = self.presentations[prs_id]
            slide = prs.slides[slide_index]
            
            # Remove the old table
            shape_to_remove = None
            for shape in slide.shapes:
                if hasattr(shape, 'table') and shape.table == table:
                    shape_to_remove = shape
                    break
            
            if shape_to_remove:
                slide_shapes = slide.shapes._spTree
                for shape_elem in slide_shapes:
                    if hasattr(shape_elem, 'get') and shape_elem.get('id') == str(shape_to_remove.shape_id):
                        slide_shapes.remove(shape_elem)
                        break
            
            # Create new table with modified dimensions
            new_table_shape = slide.shapes.add_table(new_rows, new_cols, left, top, width, height)
            new_table = new_table_shape.table
            
            # Copy and rearrange data based on operation
            for new_row in range(new_rows):
                for new_col in range(new_cols):
                    text = ""
                    
                    if operation == "add_row":
                        if new_row < position:
                            # Rows before insertion point
                            old_row = new_row
                            old_col = new_col
                        elif new_row >= position + count:
                            # Rows after insertion point
                            old_row = new_row - count
                            old_col = new_col
                        else:
                            # New rows (empty)
                            text = ""
                            old_row = -1
                            old_col = -1
                            
                    elif operation == "remove_row":
                        if new_row < position:
                            # Rows before removal point
                            old_row = new_row
                            old_col = new_col
                        else:
                            # Rows after removal point (skip removed rows)
                            old_row = new_row + count
                            old_col = new_col
                            
                    elif operation == "add_column":
                        if new_col < position:
                            # Columns before insertion point
                            old_row = new_row
                            old_col = new_col
                        elif new_col >= position + count:
                            # Columns after insertion point
                            old_row = new_row
                            old_col = new_col - count
                        else:
                            # New columns (empty)
                            text = ""
                            old_row = -1
                            old_col = -1
                            
                    elif operation == "remove_column":
                        if new_col < position:
                            # Columns before removal point
                            old_row = new_row
                            old_col = new_col
                        else:
                            # Columns after removal point (skip removed columns)
                            old_row = new_row
                            old_col = new_col + count
                    
                    # Copy text if valid coordinates
                    if old_row >= 0 and old_col >= 0 and old_row < len(cell_data) and old_col < len(cell_data[0]):
                        text = cell_data[old_row][old_col]
                    
                    # Set cell content
                    new_table.cell(new_row, new_col).text = text
            
            logger.info(f"Successfully modified table structure: {operation} (count: {count}, position: {position})")
            logger.info(f"Table dimensions changed from {current_rows}×{current_cols} to {new_rows}×{new_cols}")
            
            return True
            
        except Exception as e:
            logger.error(f"Failed to modify table structure: {e}")
            raise RuntimeError(f"Failed to modify table structure: {e}")

    # =============================================================================
    # SCREENSHOT & CRITIQUE FUNCTIONALITY
    # =============================================================================
    
    def critique_presentation(self, file_path: str, critique_type: str = "comprehensive", 
                            include_screenshots: bool = True, output_dir: Optional[str] = None) -> Dict[str, Any]:
        """
        Analyze a PowerPoint presentation and provide comprehensive critique.
        
        Args:
            file_path: Path to the PowerPoint file
            critique_type: Type of critique ("design", "content", "accessibility", "technical", "comprehensive")
            include_screenshots: Whether to generate screenshots for visual analysis
            output_dir: Directory to save screenshots (if generated)
            
        Returns:
            Dictionary containing critique results
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Presentation file not found: {file_path}")
        
        # Load presentation for analysis
        temp_prs_id = self.load_presentation(file_path)
        prs = self.presentations[temp_prs_id]
        
        critique_results = {
            "file_path": file_path,
            "critique_type": critique_type,
            "timestamp": datetime.now().isoformat(),
            "summary": {
                "total_slides": len(prs.slides),
                "overall_score": 0,
                "critical_issues": 0,
                "warnings": 0,
                "recommendations": 0
            },
            "issues": [],
            "strengths": [],
            "recommendations": [],
            "detailed_analysis": {}
        }
        
        try:
            # Generate screenshots if requested
            screenshot_paths = []
            if include_screenshots:
                try:
                    screenshot_paths = self.screenshot_slides(
                        file_path, output_dir, "PNG", 1920, 1080
                    )
                    critique_results["screenshots"] = screenshot_paths
                except Exception as e:
                    logger.warning(f"Could not generate screenshots: {e}")
            
            # Perform analysis based on critique type
            if critique_type in ["design", "comprehensive"]:
                design_analysis = self._analyze_design_quality(prs, screenshot_paths)
                critique_results["detailed_analysis"]["design"] = design_analysis
                critique_results["issues"].extend(design_analysis.get("issues", []))
                critique_results["strengths"].extend(design_analysis.get("strengths", []))
                critique_results["recommendations"].extend(design_analysis.get("recommendations", []))
            
            if critique_type in ["content", "comprehensive"]:
                content_analysis = self._analyze_content_quality(prs)
                critique_results["detailed_analysis"]["content"] = content_analysis
                critique_results["issues"].extend(content_analysis.get("issues", []))
                critique_results["strengths"].extend(content_analysis.get("strengths", []))
                critique_results["recommendations"].extend(content_analysis.get("recommendations", []))
            
            if critique_type in ["accessibility", "comprehensive"]:
                accessibility_analysis = self._analyze_accessibility(prs)
                critique_results["detailed_analysis"]["accessibility"] = accessibility_analysis
                critique_results["issues"].extend(accessibility_analysis.get("issues", []))
                critique_results["strengths"].extend(accessibility_analysis.get("strengths", []))
                critique_results["recommendations"].extend(accessibility_analysis.get("recommendations", []))
            
            if critique_type in ["technical", "comprehensive"]:
                technical_analysis = self._analyze_technical_quality(file_path, prs)
                critique_results["detailed_analysis"]["technical"] = technical_analysis
                critique_results["issues"].extend(technical_analysis.get("issues", []))
                critique_results["strengths"].extend(technical_analysis.get("strengths", []))
                critique_results["recommendations"].extend(technical_analysis.get("recommendations", []))
            
            # Calculate summary metrics
            critique_results = self._calculate_critique_summary(critique_results)
            
            return critique_results
            
        finally:
            # Clean up temporary presentation
            if temp_prs_id in self.presentations:
                del self.presentations[temp_prs_id]

    def _analyze_design_quality(self, prs, screenshot_paths: List[str] = None) -> Dict[str, Any]:
        """Analyze design quality aspects of the presentation"""
        analysis = {
            "score": 0,
            "issues": [],
            "strengths": [],
            "recommendations": [],
            "metrics": {}
        }
        
        # Font consistency analysis
        fonts_used = set()
        font_sizes = []
        slide_layouts = []
        color_usage = {}
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_fonts = set()
            slide_font_sizes = []
            
            # Analyze shapes and text
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                fonts_used.add(run.font.name)
                                slide_fonts.add(run.font.name)
                            if run.font.size:
                                font_size = run.font.size.pt
                                font_sizes.append(font_size)
                                slide_font_sizes.append(font_size)
                
                # Check for visual issues
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    # Check for problematic green rectangles
                    try:
                        if hasattr(shape.fill, 'fore_color'):
                            # Only check fore_color if the fill type supports it
                            if shape.fill.type == 1:  # MSO_FILL_TYPE.SOLID
                                rgb = shape.fill.fore_color.rgb
                                if rgb.red == 0 and rgb.green > 200 and rgb.blue == 0:
                                    analysis["issues"].append({
                                        "type": "critical",
                                        "category": "design",
                                        "slide": slide_idx + 1,
                                        "issue": "Green rectangle covering content",
                                        "description": "Detected green fill that may be obscuring slide content"
                                    })
                    except (TypeError, AttributeError):
                        # Skip shapes with unsupported fill types
                        pass
            
            # Check font consistency per slide
            if len(slide_fonts) > 3:
                analysis["issues"].append({
                    "type": "warning",
                    "category": "design",
                    "slide": slide_idx + 1,
                    "issue": "Too many fonts on single slide",
                    "description": f"Slide uses {len(slide_fonts)} different fonts, consider limiting to 2-3"
                })
        
        # Overall font analysis
        analysis["metrics"]["total_fonts"] = len(fonts_used)
        analysis["metrics"]["font_sizes_range"] = {
            "min": min(font_sizes) if font_sizes else 0,
            "max": max(font_sizes) if font_sizes else 0,
            "avg": sum(font_sizes) / len(font_sizes) if font_sizes else 0
        }
        
        if len(fonts_used) > 4:
            analysis["issues"].append({
                "type": "warning",
                "category": "design",
                "slide": "global",
                "issue": "Too many fonts in presentation",
                "description": f"Presentation uses {len(fonts_used)} fonts. Consider limiting to 2-3 for consistency."
            })
        elif len(fonts_used) <= 2:
            analysis["strengths"].append("Consistent font usage throughout presentation")
        
        # Font size analysis
        if font_sizes:
            min_size = min(font_sizes)
            if min_size < 18:
                analysis["issues"].append({
                    "type": "warning",
                    "category": "design",
                    "slide": "global",
                    "issue": "Small font sizes detected",
                    "description": f"Minimum font size is {min_size}pt. Consider 18pt+ for readability."
                })
            
            if max(font_sizes) > 72:
                analysis["issues"].append({
                    "type": "warning",
                    "category": "design",
                    "slide": "global",
                    "issue": "Very large font sizes",
                    "description": "Some text may be unnecessarily large"
                })
        
        # Add design recommendations
        if len(analysis["issues"]) == 0:
            analysis["strengths"].append("Good overall design consistency")
            analysis["score"] = 85
        else:
            analysis["score"] = max(50, 85 - len(analysis["issues"]) * 10)
            analysis["recommendations"].extend([
                "Review font consistency across slides",
                "Ensure minimum 18pt font size for readability",
                "Check for overlapping or obscured elements"
            ])
        
        return analysis

    def _analyze_content_quality(self, prs) -> Dict[str, Any]:
        """Analyze content quality and structure"""
        analysis = {
            "score": 0,
            "issues": [],
            "strengths": [],
            "recommendations": [],
            "metrics": {}
        }
        
        total_text_length = 0
        slides_with_title = 0
        slides_with_bullets = 0
        bullet_counts = []
        empty_slides = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text_length = 0
            has_title = False
            bullet_count = 0
            has_content = False
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    shape_text = shape.text_frame.text.strip()
                    if shape_text:
                        has_content = True
                        slide_text_length += len(shape_text)
                        
                        # Check if this is likely a title (large font, short text)
                        if len(shape_text) < 100 and not has_title:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.font.size and run.font.size.pt > 24:
                                        has_title = True
                                        break
                        
                        # Count bullet points
                        bullet_count += len([p for p in shape.text_frame.paragraphs if p.text.strip()])
            
            total_text_length += slide_text_length
            
            if has_title:
                slides_with_title += 1
            
            if bullet_count > 0:
                slides_with_bullets += 1
                bullet_counts.append(bullet_count)
            
            if not has_content:
                empty_slides += 1
                analysis["issues"].append({
                    "type": "warning",
                    "category": "content",
                    "slide": slide_idx + 1,
                    "issue": "Empty slide",
                    "description": "Slide contains no text content"
                })
            
            # Check for too much text
            if slide_text_length > 300:
                analysis["issues"].append({
                    "type": "warning",
                    "category": "content",
                    "slide": slide_idx + 1,
                    "issue": "Too much text",
                    "description": f"Slide has {slide_text_length} characters. Consider breaking into multiple slides."
                })
            
            # Check for too many bullets
            if bullet_count > 7:
                analysis["issues"].append({
                    "type": "warning",
                    "category": "content",
                    "slide": slide_idx + 1,
                    "issue": "Too many bullet points",
                    "description": f"Slide has {bullet_count} bullet points. Consider limiting to 5-7."
                })
        
        # Calculate metrics
        analysis["metrics"] = {
            "total_slides": len(prs.slides),
            "slides_with_title": slides_with_title,
            "slides_with_bullets": slides_with_bullets,
            "empty_slides": empty_slides,
            "avg_text_length": total_text_length / len(prs.slides) if prs.slides else 0,
            "avg_bullets_per_slide": sum(bullet_counts) / len(bullet_counts) if bullet_counts else 0
        }
        
        # Evaluate content quality
        title_ratio = slides_with_title / len(prs.slides) if prs.slides else 0
        
        if title_ratio > 0.8:
            analysis["strengths"].append("Most slides have clear titles")
        elif title_ratio < 0.5:
            analysis["issues"].append({
                "type": "warning",
                "category": "content",
                "slide": "global",
                "issue": "Missing slide titles",
                "description": f"Only {slides_with_title} of {len(prs.slides)} slides have clear titles"
            })
        
        if empty_slides > 0:
            analysis["recommendations"].append(f"Remove or add content to {empty_slides} empty slides")
        
        # Calculate score
        score = 80
        score -= empty_slides * 10
        score -= len([i for i in analysis["issues"] if i["type"] == "critical"]) * 15
        score -= len([i for i in analysis["issues"] if i["type"] == "warning"]) * 5
        score += len(analysis["strengths"]) * 5
        
        analysis["score"] = max(0, min(100, score))
        
        return analysis

    def _analyze_accessibility(self, prs) -> Dict[str, Any]:
        """Analyze accessibility aspects"""
        analysis = {
            "score": 0,
            "issues": [],
            "strengths": [],
            "recommendations": [],
            "metrics": {}
        }
        
        alt_text_missing = 0
        total_images = 0
        low_contrast_issues = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # Check images for alt text
                if hasattr(shape, 'image') or 'Picture' in str(type(shape)):
                    total_images += 1
                    # Note: python-pptx doesn't easily expose alt text, so this is a placeholder
                    # In a real implementation, you'd check shape.element for alt text
                    if not hasattr(shape, 'alt_text') or not shape.alt_text:
                        alt_text_missing += 1
                
                # Check for potential contrast issues (simplified)
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                if run.font.color and hasattr(run.font.color, 'rgb'):
                                    # Simplified contrast check
                                    if run.font.color.rgb and run.font.color.rgb.red == run.font.color.rgb.blue == run.font.color.rgb.green:
                                        # Gray text might have contrast issues
                                        low_contrast_issues += 1
                            except (TypeError, AttributeError):
                                # Skip text with unsupported color types
                                pass
        
        # Record metrics
        analysis["metrics"] = {
            "total_images": total_images,
            "alt_text_missing": alt_text_missing,
            "low_contrast_issues": low_contrast_issues
        }
        
        # Generate issues and recommendations
        if alt_text_missing > 0:
            analysis["issues"].append({
                "type": "warning",
                "category": "accessibility",
                "slide": "global",
                "issue": "Missing alt text for images",
                "description": f"{alt_text_missing} of {total_images} images lack alt text"
            })
            analysis["recommendations"].append("Add descriptive alt text to all images")
        
        if total_images > 0 and alt_text_missing == 0:
            analysis["strengths"].append("All images have alt text")
        
        # Calculate accessibility score
        score = 90
        if total_images > 0:
            score -= (alt_text_missing / total_images) * 30
        score -= low_contrast_issues * 5
        
        analysis["score"] = max(0, min(100, score))
        
        if analysis["score"] < 70:
            analysis["recommendations"].extend([
                "Review color contrast ratios",
                "Ensure text is readable against background colors",
                "Consider users with visual impairments"
            ])
        
        return analysis

    def _analyze_technical_quality(self, file_path: str, prs) -> Dict[str, Any]:
        """Analyze technical aspects of the presentation"""
        analysis = {
            "score": 0,
            "issues": [],
            "strengths": [],
            "recommendations": [],
            "metrics": {}
        }
        
        # File size analysis
        file_size = os.path.getsize(file_path)
        file_size_mb = file_size / (1024 * 1024)
        
        # Slide count analysis
        slide_count = len(prs.slides)
        
        # Performance metrics
        analysis["metrics"] = {
            "file_size_mb": round(file_size_mb, 2),
            "slide_count": slide_count,
            "avg_mb_per_slide": round(file_size_mb / slide_count if slide_count > 0 else 0, 2)
        }
        
        # File size issues
        if file_size_mb > 50:
            analysis["issues"].append({
                "type": "warning",
                "category": "technical",
                "slide": "global",
                "issue": "Large file size",
                "description": f"File size is {file_size_mb:.1f}MB. Consider optimizing images."
            })
        elif file_size_mb > 100:
            analysis["issues"].append({
                "type": "critical",
                "category": "technical",
                "slide": "global",
                "issue": "Very large file size",
                "description": f"File size is {file_size_mb:.1f}MB. May cause performance issues."
            })
        
        # Slide count analysis
        if slide_count > 50:
            analysis["issues"].append({
                "type": "warning",
                "category": "technical",
                "slide": "global",
                "issue": "Many slides",
                "description": f"Presentation has {slide_count} slides. Consider breaking into multiple presentations."
            })
        
        # Check for embedded objects and potential issues
        embedded_objects = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'chart') or hasattr(shape, 'table'):
                    embedded_objects += 1
        
        analysis["metrics"]["embedded_objects"] = embedded_objects
        
        if embedded_objects > 20:
            analysis["issues"].append({
                "type": "warning",
                "category": "technical",
                "slide": "global",
                "issue": "Many embedded objects",
                "description": f"Presentation contains {embedded_objects} charts/tables. May impact performance."
            })
        
        # Calculate technical score
        score = 90
        score -= len([i for i in analysis["issues"] if i["type"] == "critical"]) * 20
        score -= len([i for i in analysis["issues"] if i["type"] == "warning"]) * 10
        
        if file_size_mb < 10 and slide_count < 30:
            analysis["strengths"].append("Optimized file size and slide count")
            score += 10
        
        analysis["score"] = max(0, min(100, score))
        
        if analysis["score"] < 70:
            analysis["recommendations"].extend([
                "Optimize images to reduce file size",
                "Consider splitting large presentations",
                "Remove unused templates and masters"
            ])
        
        return analysis

    def _calculate_critique_summary(self, critique_results: Dict[str, Any]) -> Dict[str, Any]:
        """Calculate overall summary metrics for the critique"""
        issues = critique_results["issues"]
        recommendations = critique_results["recommendations"]
        
        # Count issue types
        critical_issues = len([i for i in issues if i.get("type") == "critical"])
        warnings = len([i for i in issues if i.get("type") == "warning"])
        
        # Calculate overall score
        detailed_analysis = critique_results["detailed_analysis"]
        scores = [analysis.get("score", 0) for analysis in detailed_analysis.values()]
        overall_score = sum(scores) / len(scores) if scores else 0
        
        # Update summary
        critique_results["summary"].update({
            "overall_score": round(overall_score, 1),
            "critical_issues": critical_issues,
            "warnings": warnings,
            "recommendations": len(set(recommendations)),  # Deduplicate recommendations
            "analysis_categories": list(detailed_analysis.keys())
        })
        
        # Add overall assessment
        if overall_score >= 80:
            critique_results["summary"]["assessment"] = "Excellent"
        elif overall_score >= 70:
            critique_results["summary"]["assessment"] = "Good"
        elif overall_score >= 60:
            critique_results["summary"]["assessment"] = "Fair"
        else:
            critique_results["summary"]["assessment"] = "Needs Improvement"
        
        return critique_results

    async def critique_presentation_async(self, file_path: str, critique_type: str = "comprehensive", 
                                        include_screenshots: bool = True, output_dir: Optional[str] = None) -> Dict[str, Any]:
        """Async wrapper for critique_presentation"""
        return await asyncio.to_thread(self.critique_presentation, file_path, critique_type, include_screenshots, output_dir)

# =============================================================================
# MCP SERVER SETUP
# =============================================================================

# Initialize server and manager
server = Server("powerpoint-mcp-stable")
ppt_manager = StablePowerPointManager()

@server.list_tools()
async def handle_list_tools() -> List[Tool]:
    """List the core essential tools including deletion and file management capabilities"""
    return [
        Tool(
            name="create_presentation",
            description="Create a new PowerPoint presentation",
            inputSchema={
                "type": "object",
                "properties": {},
                "required": [],
                "additionalProperties": False
            }
        ),
        Tool(
            name="load_presentation",
            description="Load an existing PowerPoint presentation from file",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the PowerPoint file to load"}
                },
                "required": ["file_path"],
                "additionalProperties": False,
                "examples": [
                    {
                        "file_path": "existing_presentation.pptx"
                    },
                    {
                        "file_path": "C:\\Documents\\quarterly_report.pptx"
                    }
                ]
            }
        ),
        Tool(
            name="add_slide",
            description="Add a new slide to a presentation with specified layout",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "layout_index": {"type": "integer", "description": "Slide layout index (0=title, 1=title+content, 6=blank, etc.)", "default": 6, "minimum": 0}
                },
                "required": ["presentation_id"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "layout_index": 1
                    },
                    {
                        "presentation_id": "ppt_0",
                        "layout_index": 6
                    }
                ]
            }
        ),
        Tool(
            name="add_text_box",
            description="Add a text box to a slide with comprehensive formatting options",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "text": {"type": "string", "description": "Text content"},
                    "left": {"type": "number", "default": 1, "description": "Left position in inches"},
                    "top": {"type": "number", "default": 1, "description": "Top position in inches"},
                    "width": {"type": "number", "default": 8, "description": "Width in inches"},
                    "height": {"type": "number", "default": 1, "description": "Height in inches"},
                    "font_size": {"type": "integer", "default": 18, "minimum": 8, "maximum": 72, "description": "Font size in points"},
                    "font_name": {"type": "string", "default": "Calibri", "description": "Font family name (e.g., Calibri, Arial, Times New Roman)"},
                    "font_color": {"type": "string", "description": "Font color - hex (#FF0000), RGB (255,0,0), or name (red, blue, etc.)"},
                    "bold": {"type": "boolean", "default": False, "description": "Bold text"},
                    "italic": {"type": "boolean", "default": False, "description": "Italic text"},
                    "underline": {"type": "boolean", "default": False, "description": "Underline text"},
                    "text_alignment": {"type": "string", "enum": ["left", "center", "right", "justify"], "default": "left", "description": "Text alignment"},
                    "fill_color": {"type": "string", "description": "Background fill color - hex (#FF0000), RGB (255,0,0), or name (red, blue, etc.)"},
                    "border_color": {"type": "string", "description": "Border color - hex (#FF0000), RGB (255,0,0), or name (red, blue, etc.)"},
                    "border_width": {"type": "number", "default": 0, "minimum": 0, "description": "Border width in points (0 = no border)"}
                },
                "required": ["presentation_id", "slide_index", "text"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 0,
                        "text": "Welcome to Our Presentation",
                        "font_size": 32,
                        "font_name": "Arial",
                        "font_color": "#0066CC",
                        "bold": True,
                        "text_alignment": "center"
                    },
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "text": "Key Points",
                        "font_color": "white",
                        "fill_color": "#4472C4",
                        "border_color": "black",
                        "border_width": 2
                    }
                ]
            }
        ),
        Tool(
            name="add_image",
            description="Add an image to a slide from URL or local file",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "image_source": {"type": "string", "description": "Image URL or local file path"},
                    "left": {"type": "number", "default": 1, "description": "Left position in inches"},
                    "top": {"type": "number", "default": 1, "description": "Top position in inches"},
                    "width": {"type": "number", "description": "Width in inches (optional)"},
                    "height": {"type": "number", "description": "Height in inches (optional)"}
                },
                "required": ["presentation_id", "slide_index", "image_source"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "image_source": "https://example.com/chart.png",
                        "width": 6,
                        "height": 4
                    }
                ]
            }
        ),
        Tool(
            name="add_chart",
            description="Add a chart to a slide with data series",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "chart_type": {"type": "string", "enum": ["column", "bar", "line", "pie", "area"]},
                    "categories": {"type": "array", "items": {"type": "string"}, "description": "Chart categories"},
                    "series_data": {
                        "type": "object",
                        "additionalProperties": {
                            "type": "array",
                            "items": {"type": "number"}
                        },
                        "description": "Series data as {series_name: [values]}"
                    },
                    "left": {"type": "number", "default": 2},
                    "top": {"type": "number", "default": 2},
                    "width": {"type": "number", "default": 6},
                    "height": {"type": "number", "default": 4.5}
                },
                "required": ["presentation_id", "slide_index", "chart_type", "categories", "series_data"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 2,
                        "chart_type": "column",
                        "categories": ["Q1", "Q2", "Q3", "Q4"],
                        "series_data": {
                            "Sales": [10, 15, 12, 18],
                            "Profit": [3, 5, 4, 7]
                        }
                    }
                ]
            }
        ),
        Tool(
            name="save_presentation",
            description="Save presentation to a file (handles both relative and absolute Windows paths)",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "file_path": {"type": "string", "description": "Output file path - can be relative (e.g., 'output/file.pptx') or absolute (e.g., 'C:\\Users\\name\\Documents\\file.pptx')"}
                },
                "required": ["presentation_id", "file_path"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "file_path": "my_presentation.pptx"
                    },
                    {
                        "presentation_id": "ppt_0", 
                        "file_path": "output/my_presentation.pptx"
                    },
                    {
                        "presentation_id": "ppt_0",
                        "file_path": "C:\\Users\\username\\Documents\\my_presentation.pptx"
                    }
                                 ]
             }
        ),
        Tool(
            name="extract_text",
            description="Extract all text content from a presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"}
                },
                "required": ["presentation_id"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0"
                    }
                ]
            }
        ),
        Tool(
            name="get_presentation_info",
            description="Get comprehensive information about a presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"}
                },
                "required": ["presentation_id"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0"
                    }
                ]
            }
        ),
        Tool(
            name="delete_shape",
            description="Delete a specific shape from a slide by index",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "shape_index": {"type": "integer", "description": "Shape index (0-based)", "minimum": 0}
                },
                "required": ["presentation_id", "slide_index", "shape_index"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 0,
                        "shape_index": 1
                    }
                ]
            }
        ),
        Tool(
            name="delete_slide",
            description="Delete an entire slide from the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0}
                },
                "required": ["presentation_id", "slide_index"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 2
                    }
                ]
            }
        ),
        Tool(
            name="clear_slide",
            description="Clear all content from a slide but keep the slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0}
                },
                "required": ["presentation_id", "slide_index"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1
                    }
                ]
            }
        ),
        Tool(
            name="list_slide_content",
            description="List all shapes on a slide to help with targeted deletion",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0}
                },
                "required": ["presentation_id", "slide_index"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 0
                    }
                ]
            }
        ),
        Tool(
            name="format_existing_text",
            description="Modify formatting of existing text shapes on slides",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "shape_index": {"type": "integer", "description": "Shape index (0-based)", "minimum": 0},
                    "font_size": {"type": "integer", "minimum": 8, "maximum": 72, "description": "Font size in points"},
                    "font_name": {"type": "string", "description": "Font family name (e.g., Arial, Calibri)"},
                    "font_color": {"type": "string", "description": "Font color - hex (#FF0000), RGB (255,0,0), or name (red, blue, etc.)"},
                    "bold": {"type": "boolean", "description": "Bold text"},
                    "italic": {"type": "boolean", "description": "Italic text"},
                    "underline": {"type": "boolean", "description": "Underline text"},
                    "text_alignment": {"type": "string", "enum": ["left", "center", "right", "justify"], "description": "Text alignment"}
                },
                "required": ["presentation_id", "slide_index", "shape_index"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 0,
                        "shape_index": 1,
                        "font_size": 24,
                        "font_color": "#FF0000",
                        "bold": True
                    }
                ]
            }
        ),
        Tool(
            name="set_slide_background",
            description="Set slide background color or image",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "background_color": {"type": "string", "description": "Background color - hex (#FF0000), RGB (255,0,0), or name (red, blue, etc.)"},
                    "background_image": {"type": "string", "description": "Background image URL or local file path"}
                },
                "required": ["presentation_id", "slide_index"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 0,
                        "background_color": "#E6F3FF"
                    },
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "background_image": "https://example.com/background.jpg"
                    }
                ]
            }
        ),
        Tool(
            name="add_table",
            description="Add a table to a slide with specified dimensions and optional header styling",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "rows": {"type": "integer", "description": "Number of rows", "minimum": 1, "maximum": 50},
                    "cols": {"type": "integer", "description": "Number of columns", "minimum": 1, "maximum": 20},
                    "left": {"type": "number", "default": 1, "description": "Left position in inches"},
                    "top": {"type": "number", "default": 1, "description": "Top position in inches"},
                    "width": {"type": "number", "default": 8, "description": "Width in inches"},
                    "height": {"type": "number", "default": 4, "description": "Height in inches"},
                    "header_row": {"type": "boolean", "default": False, "description": "Style first row as header with blue background"}
                },
                "required": ["presentation_id", "slide_index", "rows", "cols"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "rows": 4,
                        "cols": 3,
                        "header_row": True
                    },
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 2,
                        "rows": 5,
                        "cols": 4,
                        "left": 2,
                        "top": 2,
                        "width": 6,
                        "height": 3
                    }
                ]
            }
        ),
        Tool(
            name="set_table_cell",
            description="Set text content and formatting for a specific table cell",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "table_index": {"type": "integer", "description": "Table index on slide (0-based)", "minimum": 0},
                    "row": {"type": "integer", "description": "Row index (0-based)", "minimum": 0},
                    "col": {"type": "integer", "description": "Column index (0-based)", "minimum": 0},
                    "text": {"type": "string", "description": "Text content for the cell"},
                    "font_size": {"type": "integer", "minimum": 8, "maximum": 72, "description": "Font size in points"},
                    "font_name": {"type": "string", "description": "Font family name (e.g., Arial, Calibri)"},
                    "font_color": {"type": "string", "description": "Font color - hex (#FF0000), RGB (255,0,0), or name (red, blue, etc.)"},
                    "bold": {"type": "boolean", "description": "Bold text"},
                    "italic": {"type": "boolean", "description": "Italic text"},
                    "underline": {"type": "boolean", "description": "Underline text"},
                    "text_alignment": {"type": "string", "enum": ["left", "center", "right", "justify"], "description": "Text alignment within cell"}
                },
                "required": ["presentation_id", "slide_index", "table_index", "row", "col", "text"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_index": 0,
                        "row": 0,
                        "col": 0,
                        "text": "Product Name",
                        "bold": True,
                        "font_size": 14
                    },
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_index": 0,
                        "row": 1,
                        "col": 1,
                        "text": "$125.99",
                        "font_color": "#008000",
                        "text_alignment": "right"
                    }
                ]
            }
        ),
        Tool(
            name="get_table_info",
            description="Get comprehensive information about a table including dimensions and cell contents",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "table_index": {"type": "integer", "description": "Table index on slide (0-based)", "minimum": 0}
                },
                "required": ["presentation_id", "slide_index", "table_index"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_index": 0
                    }
                ]
            }
        ),
        Tool(
            name="style_table_cell",
            description="Apply styling to a specific table cell (background color, borders, margins)",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "table_index": {"type": "integer", "description": "Table index on slide (0-based)", "minimum": 0},
                    "row": {"type": "integer", "description": "Row index (0-based)", "minimum": 0},
                    "col": {"type": "integer", "description": "Column index (0-based)", "minimum": 0},
                    "fill_color": {"type": "string", "description": "Cell background color - hex (#FF0000), RGB (255,0,0), or name (red, blue, etc.)"},
                    "border_color": {"type": "string", "description": "Border color - hex (#FF0000), RGB (255,0,0), or name (red, blue, etc.)"},
                    "border_width": {"type": "number", "description": "Border width in points", "minimum": 0},
                    "margin_left": {"type": "number", "description": "Left margin in inches", "minimum": 0},
                    "margin_right": {"type": "number", "description": "Right margin in inches", "minimum": 0},
                    "margin_top": {"type": "number", "description": "Top margin in inches", "minimum": 0},
                    "margin_bottom": {"type": "number", "description": "Bottom margin in inches", "minimum": 0}
                },
                "required": ["presentation_id", "slide_index", "table_index", "row", "col"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_index": 0,
                        "row": 0,
                        "col": 0,
                        "fill_color": "#4472C4",
                        "border_color": "black",
                        "border_width": 1
                    },
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_index": 0,
                        "row": 1,
                        "col": 1,
                        "fill_color": "#E6F3FF",
                        "margin_left": 0.1,
                        "margin_right": 0.1
                    }
                ]
            }
        ),
        Tool(
            name="style_table_range",
            description="Apply styling to a range of table cells simultaneously",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "table_index": {"type": "integer", "description": "Table index on slide (0-based)", "minimum": 0},
                    "start_row": {"type": "integer", "description": "Starting row index (0-based)", "minimum": 0},
                    "start_col": {"type": "integer", "description": "Starting column index (0-based)", "minimum": 0},
                    "end_row": {"type": "integer", "description": "Ending row index (0-based)", "minimum": 0},
                    "end_col": {"type": "integer", "description": "Ending column index (0-based)", "minimum": 0},
                    "fill_color": {"type": "string", "description": "Cell background color - hex (#FF0000), RGB (255,0,0), or name (red, blue, etc.)"},
                    "border_color": {"type": "string", "description": "Border color - hex (#FF0000), RGB (255,0,0), or name (red, blue, etc.)"},
                    "border_width": {"type": "number", "description": "Border width in points", "minimum": 0},
                    "margin_left": {"type": "number", "description": "Left margin in inches", "minimum": 0},
                    "margin_right": {"type": "number", "description": "Right margin in inches", "minimum": 0},
                    "margin_top": {"type": "number", "description": "Top margin in inches", "minimum": 0},
                    "margin_bottom": {"type": "number", "description": "Bottom margin in inches", "minimum": 0}
                },
                "required": ["presentation_id", "slide_index", "table_index", "start_row", "start_col", "end_row", "end_col"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_index": 0,
                        "start_row": 0,
                        "start_col": 0,
                        "end_row": 0,
                        "end_col": 2,
                        "fill_color": "#4472C4",
                        "border_color": "black",
                        "border_width": 1
                    },
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_index": 0,
                        "start_row": 1,
                        "start_col": 0,
                        "end_row": 3,
                        "end_col": 2,
                        "fill_color": "#F2F2F2"
                    }
                ]
            }
        ),
        Tool(
            name="create_table_with_data",
            description="Create a table and populate it with data in one operation (convenience method)",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "table_data": {
                        "type": "array",
                        "items": {
                            "type": "array",
                            "items": {"type": "string"}
                        },
                        "description": "2D array of table data (rows and columns)"
                    },
                    "headers": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Optional header row"
                    },
                    "left": {"type": "number", "default": 1, "description": "Left position in inches"},
                    "top": {"type": "number", "default": 1, "description": "Top position in inches"},
                    "width": {"type": "number", "default": 8, "description": "Width in inches"},
                    "height": {"type": "number", "default": 4, "description": "Height in inches"},
                    "header_style": {
                        "type": "object",
                        "description": "Style options for header row (font_size, font_color, bold, etc.)"
                    },
                    "data_style": {
                        "type": "object",
                        "description": "Style options for data cells (font_size, font_color, etc.)"
                    },
                    "alternating_rows": {"type": "boolean", "default": False, "description": "Apply alternating row background colors"}
                },
                "required": ["presentation_id", "slide_index", "table_data"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_data": [
                            ["John", "25", "Engineer"],
                            ["Jane", "30", "Manager"],
                            ["Bob", "28", "Designer"]
                        ],
                        "headers": ["Name", "Age", "Role"],
                        "alternating_rows": True
                    },
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 2,
                        "table_data": [
                            ["Q1", "100", "120"],
                            ["Q2", "110", "135"],
                            ["Q3", "105", "125"],
                            ["Q4", "115", "140"]
                        ],
                        "headers": ["Quarter", "Sales", "Target"],
                        "header_style": {"bold": True, "font_size": 14},
                        "data_style": {"font_size": 12}
                    }
                ]
            }
        ),
        Tool(
            name="modify_table_structure",
            description="Modify table structure by adding or removing rows and columns",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "Presentation ID"},
                    "slide_index": {"type": "integer", "description": "Slide index (0-based)", "minimum": 0},
                    "table_index": {"type": "integer", "description": "Table index on slide (0-based)", "minimum": 0},
                    "operation": {
                        "type": "string", 
                        "enum": ["add_row", "remove_row", "add_column", "remove_column"],
                        "description": "Type of structure modification to perform"
                    },
                    "position": {
                        "type": "integer",
                        "description": "Position to insert/remove at (0-based). If not specified: add operations append to end, remove operations remove from end",
                        "minimum": 0
                    },
                    "count": {
                        "type": "integer",
                        "default": 1,
                        "minimum": 1,
                        "maximum": 20,
                        "description": "Number of rows/columns to add or remove"
                    }
                },
                "required": ["presentation_id", "slide_index", "table_index", "operation"],
                "additionalProperties": False,
                "examples": [
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_index": 0,
                        "operation": "add_row",
                        "position": 2,
                        "count": 1
                    },
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_index": 0,
                        "operation": "remove_column",
                        "count": 2
                    },
                    {
                        "presentation_id": "ppt_0",
                        "slide_index": 1,
                        "table_index": 0,
                        "operation": "add_column",
                        "position": 0,
                        "count": 1
                    }
                ]
            }
        ),
        Tool(
            name="get_presentation_info",
            description="Get information about a presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"}
                },
                "required": ["presentation_id"]
            }
        ),
        Tool(
            name="screenshot_slides",
            description="Screenshot each slide of a PowerPoint presentation for vision review (Windows only)",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the PowerPoint file"
                    },
                    "output_dir": {
                        "type": "string",
                        "description": "Directory to save screenshots (optional, defaults to temp directory)"
                    },
                    "image_format": {
                        "type": "string",
                        "description": "Image format (PNG, JPG, etc.)",
                        "default": "PNG"
                    },
                    "width": {
                        "type": "integer",
                        "description": "Screenshot width in pixels",
                        "default": 1920
                    },
                    "height": {
                        "type": "integer",
                        "description": "Screenshot height in pixels",  
                        "default": 1080
                    }
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="critique_presentation",
            description="Analyze and critique a PowerPoint presentation for design, content, accessibility, and technical issues",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the PowerPoint file to analyze"
                    },
                    "critique_type": {
                        "type": "string",
                        "enum": ["design", "content", "accessibility", "technical", "comprehensive"],
                        "default": "comprehensive",
                        "description": "Type of critique to perform"
                    },
                    "include_screenshots": {
                        "type": "boolean",
                        "default": True,
                        "description": "Whether to generate screenshots for visual analysis"
                    },
                    "output_dir": {
                        "type": "string",
                        "description": "Directory to save screenshots if generated (optional)"
                    }
                },
                "required": ["file_path"]
            }
        )
    ]

@server.call_tool()
async def handle_call_tool(name: str, arguments: Dict[str, Any]) -> List[TextContent]:
    """Handle tool calls with validation and enhanced feedback"""
    
    try:
        # Validate arguments
        validated_args = validate_basic_args(name, arguments)
        
        if name == "create_presentation":
            prs_id = ppt_manager.create_presentation()
            message = format_success_message(name, presentation_id=prs_id)
            return [TextContent(type="text", text=message)]
        
        elif name == "load_presentation":
            file_path = validated_args["file_path"]
            
            prs_id = ppt_manager.load_presentation(file_path)
            
            # Get info for success message
            prs = ppt_manager.presentations[prs_id]
            slide_count = len(prs.slides)
            
            message = format_success_message(
                name, presentation_id=prs_id, file_path=file_path, slide_count=slide_count
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "add_slide":
            prs_id = validated_args["presentation_id"]
            layout_index = validated_args.get("layout_index", 6)
            
            slide_index = ppt_manager.add_slide(prs_id, layout_index)
            
            # Get info for success message
            prs = ppt_manager.presentations[prs_id]
            total_slides = len(prs.slides)
            
            # Try to get layout name
            try:
                layout_name = prs.slide_layouts[layout_index].name
            except:
                layout_name = f"Layout {layout_index}"
            
            message = format_success_message(
                name, slide_index=slide_index, layout_index=layout_index, 
                layout_name=layout_name, total_slides=total_slides
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "add_text_box":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            text = validated_args["text"]
            left = validated_args.get("left", 1)
            top = validated_args.get("top", 1)
            width = validated_args.get("width", 8)
            height = validated_args.get("height", 1)
            font_size = validated_args.get("font_size", 18)
            font_name = validated_args.get("font_name", "Calibri")
            font_color = validated_args.get("font_color")
            bold = validated_args.get("bold", False)
            italic = validated_args.get("italic", False)
            underline = validated_args.get("underline", False)
            text_alignment = validated_args.get("text_alignment", "left")
            fill_color = validated_args.get("fill_color")
            border_color = validated_args.get("border_color")
            border_width = validated_args.get("border_width", 0)
            
            success = ppt_manager.add_text_box(
                prs_id, slide_index, text, left, top, width, height, font_size, font_name, font_color,
                bold, italic, underline, text_alignment, fill_color, border_color, border_width
            )
            
            message = format_success_message(
                name, slide_index=slide_index, font_size=font_size, font_name=font_name, 
                text_alignment=text_alignment, font_color=font_color, fill_color=fill_color, text=text
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "add_image":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            image_source = validated_args["image_source"]
            left = validated_args.get("left", 1)
            top = validated_args.get("top", 1)
            width = validated_args.get("width")
            height = validated_args.get("height")
            
            success = ppt_manager.add_image(
                prs_id, slide_index, image_source, left, top, width, height
            )
            
            message = format_success_message(
                name, slide_index=slide_index, image_source=image_source
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "add_chart":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            chart_type = validated_args["chart_type"]
            categories = validated_args["categories"]
            series_data = validated_args["series_data"]
            left = validated_args.get("left", 2)
            top = validated_args.get("top", 2)
            width = validated_args.get("width", 6)
            height = validated_args.get("height", 4.5)
            
            success = ppt_manager.add_chart(
                prs_id, slide_index, chart_type, categories, series_data, left, top, width, height
            )
            
            message = format_success_message(
                name, slide_index=slide_index, chart_type=chart_type, 
                categories=categories, series_data=series_data
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "save_presentation":
            prs_id = validated_args["presentation_id"]
            file_path = validated_args["file_path"]
            
            saved_path = ppt_manager.save_presentation(prs_id, file_path)
            
            # Return with embedded resource for immediate access
            try:
                with open(saved_path, 'rb') as f:
                    file_data = f.read()
                
                message = format_success_message(name, file_path=saved_path)
                
                return [
                    TextContent(type="text", text=message),
                    EmbeddedResource(
                        type="resource",
                        resource=EmbeddedResource(
                            uri=f"file://{saved_path}",
                            mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            data=file_data
                        )
                    )
                ]
            except Exception as e:
                # Fallback to text message only
                message = format_success_message(name, file_path=saved_path)
                return [TextContent(type="text", text=message)]
        
        elif name == "extract_text":
            prs_id = validated_args["presentation_id"]
            
            extracted_text = ppt_manager.extract_text(prs_id)
            
            # Count total text items
            text_items = sum(len(slide["text_content"]) for slide in extracted_text)
            slide_count = len(extracted_text)
            
            message = format_success_message(
                name, slide_count=slide_count, text_items=text_items
            )
            
            # Format the extracted text for display
            if extracted_text:
                text_summary = []
                for slide in extracted_text:
                    if slide["text_content"]:
                        slide_text = f"Slide {slide['slide_number']}:"
                        for item in slide["text_content"]:
                            preview = item["text"][:80] + "..." if len(item["text"]) > 80 else item["text"]
                            slide_text += f"\n  - {preview}"
                        text_summary.append(slide_text)
                
                if text_summary:
                    full_message = f"{message}\n\n" + "\n\n".join(text_summary)
                else:
                    full_message = f"{message}\n(No text content found)"
            else:
                full_message = f"{message}\n(No slides found)"
            
            return [TextContent(type="text", text=full_message)]
        
        elif name == "get_presentation_info":
            prs_id = validated_args["presentation_id"]
            
            info = ppt_manager.get_presentation_info(prs_id)
            
            message = format_success_message(
                name, slide_count=info["slide_count"], total_shapes=info["total_shapes"]
            )
            
            # Format comprehensive info
            content_summary = info["content_summary"]
            available_layouts = info["available_layouts"]
            
            info_details = f"""📊 Content Summary:
  • Text boxes: {content_summary['text_boxes']}
  • Images: {content_summary['images']}
  • Charts: {content_summary['charts']}
  • Other shapes: {content_summary['other_shapes']}

🎨 Available Layouts:"""
            
            for layout in available_layouts[:5]:  # Show first 5 layouts
                info_details += f"\n  [{layout['index']}] {layout['name']}"
            
            if len(available_layouts) > 5:
                info_details += f"\n  ... and {len(available_layouts) - 5} more layouts"
            
            full_message = f"{message}\n\n{info_details}"
            return [TextContent(type="text", text=full_message)]
        
        elif name == "delete_shape":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            shape_index = validated_args["shape_index"]
            
            # Get shape type before deletion for better message
            prs = ppt_manager.presentations[prs_id]
            slide = prs.slides[slide_index]
            shape = slide.shapes[shape_index]
            shape_type = "shape"
            try:
                if hasattr(shape, 'text_frame'):
                    shape_type = "text box"
                elif hasattr(shape, 'chart'):
                    shape_type = "chart"
                elif hasattr(shape, 'image'):
                    shape_type = "image"
            except:
                pass
            
            success = ppt_manager.delete_shape(prs_id, slide_index, shape_index)
            message = format_success_message(
                name, slide_index=slide_index, shape_index=shape_index, shape_type=shape_type
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "delete_slide":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            
            # Get slide count before deletion
            prs = ppt_manager.presentations[prs_id]
            original_count = len(prs.slides)
            
            success = ppt_manager.delete_slide(prs_id, slide_index)
            remaining_slides = original_count - 1
            message = format_success_message(
                name, slide_index=slide_index, remaining_slides=remaining_slides
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "clear_slide":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            
            # Get shape count before clearing
            prs = ppt_manager.presentations[prs_id]
            slide = prs.slides[slide_index]
            shapes_cleared = len(slide.shapes)
            
            success = ppt_manager.clear_slide(prs_id, slide_index)
            message = format_success_message(
                name, slide_index=slide_index, shapes_cleared=shapes_cleared
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "list_slide_content":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            
            content = ppt_manager.list_slide_content(prs_id, slide_index)
            message = format_success_message(
                name, slide_index=slide_index, shape_count=content["shape_count"]
            )
            
            # Format the detailed content list
            if content["shapes"]:
                content_list = "\n".join([
                    f"  [{shape['index']}] {shape['type']}: {shape['description']}"
                    for shape in content["shapes"]
                ])
                full_message = f"{message}\n{content_list}"
            else:
                full_message = f"{message}\n  (No shapes on this slide)"
            
            return [TextContent(type="text", text=full_message)]
        
        elif name == "format_existing_text":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            shape_index = validated_args["shape_index"]
            font_size = validated_args.get("font_size")
            font_name = validated_args.get("font_name")
            font_color = validated_args.get("font_color")
            bold = validated_args.get("bold")
            italic = validated_args.get("italic")
            underline = validated_args.get("underline")
            text_alignment = validated_args.get("text_alignment")
            
            success = ppt_manager.format_existing_text(
                prs_id, slide_index, shape_index, font_size, font_name, font_color,
                bold, italic, underline, text_alignment
            )
            
            message = format_success_message(
                name, slide_index=slide_index, shape_index=shape_index,
                font_size=font_size, font_name=font_name, font_color=font_color, text_alignment=text_alignment
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "set_slide_background":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            background_color = validated_args.get("background_color")
            background_image = validated_args.get("background_image")
            
            success = ppt_manager.set_slide_background(
                prs_id, slide_index, background_color, background_image
            )
            
            message = format_success_message(
                name, slide_index=slide_index, background_color=background_color, background_image=background_image
            )
            return [TextContent(type="text", text=message)]
        
        # Table operations - Phase 1 handlers
        elif name == "add_table":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            rows = validated_args["rows"]
            cols = validated_args["cols"]
            left = validated_args.get("left", 1)
            top = validated_args.get("top", 1)
            width = validated_args.get("width", 8)
            height = validated_args.get("height", 4)
            header_row = validated_args.get("header_row", False)
            
            table_index = ppt_manager.add_table(
                prs_id, slide_index, rows, cols, left, top, width, height, header_row
            )
            
            message = format_success_message(
                name, slide_index=slide_index, rows=rows, cols=cols, header_row=header_row
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "set_table_cell":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            table_index = validated_args["table_index"]
            row = validated_args["row"]
            col = validated_args["col"]
            text = validated_args["text"]
            font_size = validated_args.get("font_size")
            font_name = validated_args.get("font_name")
            font_color = validated_args.get("font_color")
            bold = validated_args.get("bold")
            italic = validated_args.get("italic")
            underline = validated_args.get("underline")
            text_alignment = validated_args.get("text_alignment")
            
            success = ppt_manager.set_table_cell(
                prs_id, slide_index, table_index, row, col, text,
                font_size, font_name, font_color, bold, italic, underline, text_alignment
            )
            
            message = format_success_message(
                name, slide_index=slide_index, table_index=table_index, row=row, col=col, text=text
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "get_table_info":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            table_index = validated_args["table_index"]
            
            info = ppt_manager.get_table_info(prs_id, slide_index, table_index)
            
            message = format_success_message(
                name, slide_index=slide_index, table_index=table_index, 
                rows=info["rows"], cols=info["columns"], total_cells=info["total_cells"]
            )
            
            # Format detailed table info
            table_details = f"""📊 Table Structure:
  • Dimensions: {info['rows']} rows × {info['columns']} columns
  • Total cells: {info['total_cells']}

📝 Cell Contents:"""
            
            # Show first few rows of content
            for row_idx, row_data in enumerate(info['cell_data'][:3]):  # Show first 3 rows
                row_content = []
                for cell_data in row_data:
                    cell_text = cell_data['text']
                    if cell_text:
                        # Truncate long cell content for display
                        display_text = cell_text[:15] + "..." if len(cell_text) > 15 else cell_text
                        row_content.append(f'"{display_text}"')
                    else:
                        row_content.append('""')
                
                table_details += f"\n  Row {row_idx}: {' | '.join(row_content)}"
            
            if len(info['cell_data']) > 3:
                table_details += f"\n  ... and {len(info['cell_data']) - 3} more rows"
            
            full_message = f"{message}\n\n{table_details}"
            return [TextContent(type="text", text=full_message)]
        
        # Table operations - Phase 2 handlers
        elif name == "style_table_cell":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            table_index = validated_args["table_index"]
            row = validated_args["row"]
            col = validated_args["col"]
            fill_color = validated_args.get("fill_color")
            border_color = validated_args.get("border_color")
            border_width = validated_args.get("border_width")
            margin_left = validated_args.get("margin_left")
            margin_right = validated_args.get("margin_right")
            margin_top = validated_args.get("margin_top")
            margin_bottom = validated_args.get("margin_bottom")
            
            success = ppt_manager.style_table_cell(
                prs_id, slide_index, table_index, row, col,
                fill_color, border_color, border_width,
                margin_left, margin_right, margin_top, margin_bottom
            )
            
            message = format_success_message(
                name, slide_index=slide_index, table_index=table_index, row=row, col=col,
                fill_color=fill_color, border_color=border_color
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "style_table_range":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            table_index = validated_args["table_index"]
            start_row = validated_args["start_row"]
            start_col = validated_args["start_col"]
            end_row = validated_args["end_row"]
            end_col = validated_args["end_col"]
            fill_color = validated_args.get("fill_color")
            border_color = validated_args.get("border_color")
            border_width = validated_args.get("border_width")
            margin_left = validated_args.get("margin_left")
            margin_right = validated_args.get("margin_right")
            margin_top = validated_args.get("margin_top")
            margin_bottom = validated_args.get("margin_bottom")
            
            success = ppt_manager.style_table_range(
                prs_id, slide_index, table_index, start_row, start_col, end_row, end_col,
                fill_color, border_color, border_width,
                margin_left, margin_right, margin_top, margin_bottom
            )
            
            message = format_success_message(
                name, slide_index=slide_index, table_index=table_index,
                start_row=start_row, start_col=start_col, end_row=end_row, end_col=end_col
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "create_table_with_data":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            table_data = validated_args["table_data"]
            headers = validated_args.get("headers")
            left = validated_args.get("left", 1)
            top = validated_args.get("top", 1)
            width = validated_args.get("width", 8)
            height = validated_args.get("height", 4)
            header_style = validated_args.get("header_style", {})
            data_style = validated_args.get("data_style", {})
            alternating_rows = validated_args.get("alternating_rows", False)
            
            table_index = ppt_manager.create_table_with_data(
                prs_id, slide_index, table_data, headers, left, top, width, height,
                header_style, data_style, alternating_rows
            )
            
            # Calculate table dimensions for success message
            rows = len(table_data) + (1 if headers else 0)
            cols = len(table_data[0]) if table_data else 0
            
            message = format_success_message(
                "add_table", slide_index=slide_index, rows=rows, cols=cols, header_row=bool(headers)
            )
            
            # Add data population info
            data_summary = f"\n📊 Populated with {len(table_data)} data rows"
            if headers:
                data_summary += f" and {len(headers)} headers"
            if alternating_rows:
                data_summary += " (alternating row colors)"
            
            full_message = f"{message}{data_summary}"
            return [TextContent(type="text", text=full_message)]
        
        elif name == "modify_table_structure":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            table_index = validated_args["table_index"]
            operation = validated_args["operation"]
            position = validated_args.get("position")
            count = validated_args.get("count", 1)
            
            success = ppt_manager.modify_table_structure(
                prs_id, slide_index, table_index, operation, position, count
            )
            
            message = format_success_message(
                name, slide_index=slide_index, table_index=table_index,
                operation=operation, position=position, count=count
            )
            return [TextContent(type="text", text=message)]
        
        elif name == "get_presentation_info":
            prs_id = validated_args["presentation_id"]
            
            info = ppt_manager.get_presentation_info(prs_id)
            
            message = format_success_message(
                name, slide_count=info["slide_count"], total_shapes=info["total_shapes"]
            )
            
            # Format comprehensive info
            content_summary = info["content_summary"]
            available_layouts = info["available_layouts"]
            
            info_details = f"""📊 Content Summary:
  • Text boxes: {content_summary['text_boxes']}
  • Images: {content_summary['images']}
  • Charts: {content_summary['charts']}
  • Other shapes: {content_summary['other_shapes']}

🎨 Available Layouts:"""
            
            for layout in available_layouts[:5]:  # Show first 5 layouts
                info_details += f"\n  [{layout['index']}] {layout['name']}"
            
            if len(available_layouts) > 5:
                info_details += f"\n  ... and {len(available_layouts) - 5} more layouts"
            
            full_message = f"{message}\n\n{info_details}"
            return [TextContent(type="text", text=full_message)]
        
        elif name == "screenshot_slides":
            file_path = validated_args["file_path"]
            output_dir = validated_args.get("output_dir")
            image_format = validated_args.get("image_format", "PNG")
            width = validated_args.get("width", 1920)
            height = validated_args.get("height", 1080)
            
            screenshot_paths = await ppt_manager.screenshot_slides_async(
                file_path, output_dir, image_format, width, height
            )
            
            result_info = {
                "total_slides": len(screenshot_paths),
                "screenshot_paths": screenshot_paths,
                "image_format": image_format,
                "dimensions": f"{width}x{height}",
                "output_directory": os.path.dirname(screenshot_paths[0]) if screenshot_paths else None
            }
            
            return [TextContent(
                type="text",
                text=f"Successfully created {len(screenshot_paths)} slide screenshots.\n" +
                     json.dumps(result_info, indent=2)
            )]
        
        elif name == "critique_presentation":
            file_path = validated_args["file_path"]
            critique_type = validated_args.get("critique_type", "comprehensive")
            include_screenshots = validated_args.get("include_screenshots", True)
            output_dir = validated_args.get("output_dir")
            
            critique_results = await ppt_manager.critique_presentation_async(
                file_path, critique_type, include_screenshots, output_dir
            )
            
            # Format the critique results for display
            summary = critique_results["summary"]
            response_text = f"""🔍 Presentation Critique Complete

📊 Overall Assessment: {summary['assessment']} (Score: {summary['overall_score']}/100)
📈 Total Slides: {summary['total_slides']}
🔴 Critical Issues: {summary['critical_issues']}
⚠️  Warnings: {summary['warnings']}
💡 Recommendations: {summary['recommendations']}

Analysis Categories: {', '.join(summary['analysis_categories'])}

"""
            
            # Add issue details
            if critique_results["issues"]:
                response_text += "\n🚨 Issues Found:\n"
                for issue in critique_results["issues"][:10]:  # Limit to first 10 issues
                    emoji = "🔴" if issue["type"] == "critical" else "⚠️"
                    slide_info = f"Slide {issue['slide']}" if issue['slide'] != 'global' else "Global"
                    response_text += f"{emoji} {slide_info}: {issue['issue']} - {issue['description']}\n"
                
                if len(critique_results["issues"]) > 10:
                    response_text += f"... and {len(critique_results['issues']) - 10} more issues\n"
            
            # Add strengths
            if critique_results["strengths"]:
                response_text += "\n✅ Strengths:\n"
                for strength in critique_results["strengths"][:5]:
                    response_text += f"• {strength}\n"
            
            # Add top recommendations
            if critique_results["recommendations"]:
                response_text += "\n💡 Top Recommendations:\n"
                unique_recommendations = list(set(critique_results["recommendations"]))
                for rec in unique_recommendations[:5]:
                    response_text += f"• {rec}\n"
            
            # Add screenshot info if generated
            if critique_results.get("screenshots"):
                response_text += f"\n📸 Screenshots: {len(critique_results['screenshots'])} images generated\n"
            
            response_text += f"\n📋 Full detailed analysis available in JSON format below:\n"
            
            response = [
                TextContent(
                    type="text",
                    text=response_text
                ),
                TextContent(
                    type="text",
                    text=json.dumps(critique_results, indent=2, default=str)
                )
            ]
            
            # Add screenshot references if generated
            if critique_results.get("screenshots"):
                for screenshot_path in critique_results["screenshots"]:
                    if os.path.exists(screenshot_path):
                        response.append(EmbeddedResource(
                            uri=f"file://{os.path.abspath(screenshot_path)}",
                            mimeType="image/png"
                        ))
            
            return response
        
        else:
            return [TextContent(type="text", text=f"Unknown tool: {name}")]
    
    except Exception as e:
        logger.error(f"Error in tool {name}: {e}")
        return [TextContent(
            type="text",
            text=f"Error: {str(e)}"
        )]

async def main():
    """Main entry point for the stable PowerPoint MCP server"""
    options = InitializationOptions(
        server_name="powerpoint-mcp-stable",
        server_version="1.0.0",
        capabilities={
            "resources": {},
            "tools": {}
        }
    )
    
    try:
        async with stdio_server() as (read_stream, write_stream):
            await server.run(read_stream, write_stream, options)
    except KeyboardInterrupt:
        logger.info("Server interrupted by user")
    except Exception as e:
        logger.error(f"Server error: {e}")

if __name__ == "__main__":
    asyncio.run(main()) 
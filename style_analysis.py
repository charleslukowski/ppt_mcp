#!/usr/bin/env python3
"""
PowerPoint Style Analysis and Learning Engine

This module provides comprehensive style analysis capabilities for PowerPoint presentations,
extracting style patterns and creating reusable style profiles for consistent presentation design.

Phase 4 Implementation: Style Intelligence & Consistency
"""

import json
import logging
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union
from collections import defaultdict, Counter
import xml.etree.ElementTree as ET

try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.text.text import TextFrame
    from pptx.dml.color import RGBColor, ColorFormat
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as e:
    print(f"python-pptx library not found: {e}")
    raise

try:
    import numpy as np
    from sklearn.cluster import KMeans
    from sklearn.preprocessing import StandardScaler
except ImportError as e:
    print(f"Machine learning libraries not found: {e}")
    print("Please install with: pip install scikit-learn numpy")
    raise

# Configure logging
logger = logging.getLogger(__name__)

@dataclass
class FontProfile:
    """Represents font styling information"""
    family: str
    size_pt: float
    bold: bool
    italic: bool
    underline: bool
    color_rgb: Optional[Tuple[int, int, int]]
    usage_frequency: int = 1

@dataclass
class ColorProfile:
    """Represents color usage patterns"""
    rgb: Tuple[int, int, int]
    hex_code: str
    usage_context: str  # 'text', 'background', 'accent', 'shape'
    frequency: int = 1

@dataclass
class LayoutProfile:
    """Represents layout and positioning patterns"""
    average_left_margin: float
    average_top_margin: float
    common_widths: List[float]
    common_heights: List[float]
    grid_alignment: Dict[str, float]  # grid snap patterns
    shape_distribution: Dict[str, int]  # shape type counts

@dataclass
class TextHierarchy:
    """Represents text hierarchy patterns"""
    title_style: FontProfile
    subtitle_style: Optional[FontProfile]
    body_style: FontProfile
    bullet_style: Optional[FontProfile]
    caption_style: Optional[FontProfile]

@dataclass
class StyleProfile:
    """Complete style profile for a presentation or template"""
    name: str
    description: str
    font_hierarchy: TextHierarchy
    color_palette: List[ColorProfile]
    layout_patterns: LayoutProfile
    slide_layouts: Dict[str, Dict[str, Any]]
    theme_colors: Dict[str, str]
    created_from: str  # source file path
    confidence_score: float  # 0-1, how confident we are in this profile

class StyleAnalyzer:
    """Analyzes PowerPoint presentations to extract style patterns"""
    
    def __init__(self):
        self.style_profiles: Dict[str, StyleProfile] = {}
        self.temp_analysis_data = {}
    
    def analyze_presentation_style(self, file_path: str) -> Dict[str, Any]:
        """
        Analyze a PowerPoint presentation to extract comprehensive style information
        
        Args:
            file_path: Path to the PowerPoint file to analyze
            
        Returns:
            Dictionary containing detailed style analysis results
        """
        try:
            if not Path(file_path).exists():
                raise FileNotFoundError(f"Presentation file not found: {file_path}")
            
            prs = Presentation(file_path)
            logger.info(f"Starting style analysis for: {file_path}")
            
            analysis_results = {
                'file_path': file_path,
                'slide_count': len(prs.slides),
                'fonts': self._analyze_fonts(prs),
                'colors': self._analyze_colors(prs),
                'layouts': self._analyze_layouts(prs),
                'text_hierarchy': self._analyze_text_hierarchy(prs),
                'theme': self._analyze_theme(prs),
                'shapes': self._analyze_shapes(prs),
                'consistency_score': 0.0  # Will be calculated
            }
            
            # Calculate consistency score
            analysis_results['consistency_score'] = self._calculate_consistency_score(analysis_results)
            
            logger.info(f"Style analysis completed for {file_path}")
            return analysis_results
            
        except Exception as e:
            logger.error(f"Error analyzing presentation style: {e}")
            raise
    
    def _analyze_fonts(self, prs: Presentation) -> Dict[str, Any]:
        """Analyze font usage patterns across the presentation"""
        font_usage = defaultdict(int)
        font_details = {}
        size_patterns = defaultdict(int)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    self._extract_font_info(shape.text_frame, font_usage, font_details, size_patterns)
        
        # Find most common fonts and sizes
        primary_font = max(font_usage.items(), key=lambda x: x[1])[0] if font_usage else 'Calibri'
        common_sizes = sorted(size_patterns.items(), key=lambda x: x[1], reverse=True)[:5]
        
        return {
            'primary_font': primary_font,
            'font_usage': dict(font_usage),
            'font_details': font_details,
            'common_sizes': common_sizes,
            'size_patterns': dict(size_patterns)
        }
    
    def _extract_font_info(self, text_frame: TextFrame, font_usage: dict, font_details: dict, size_patterns: dict):
        """Extract detailed font information from a text frame"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.name:
                    font_name = run.font.name
                    font_usage[font_name] += 1
                    
                    if font_name not in font_details:
                        font_details[font_name] = {
                            'sizes': set(),
                            'bold_usage': 0,
                            'italic_usage': 0,
                            'colors': set()
                        }
                    
                    if run.font.size:
                        size_pt = run.font.size.pt
                        font_details[font_name]['sizes'].add(size_pt)
                        size_patterns[size_pt] += 1
                    
                    if run.font.bold:
                        font_details[font_name]['bold_usage'] += 1
                    
                    if run.font.italic:
                        font_details[font_name]['italic_usage'] += 1
                    
                    # Extract color if available
                    try:
                        if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                            color_rgb = (run.font.color.rgb.r, run.font.color.rgb.g, run.font.color.rgb.b)
                            font_details[font_name]['colors'].add(color_rgb)
                    except (AttributeError, TypeError):
                        # Skip colors that can't be extracted
                        continue
        
        # Convert sets to lists for JSON serialization
        for font_name in font_details:
            font_details[font_name]['sizes'] = list(font_details[font_name]['sizes'])
            font_details[font_name]['colors'] = list(font_details[font_name]['colors'])
    
    def _analyze_colors(self, prs: Presentation) -> Dict[str, Any]:
        """Analyze color usage patterns throughout the presentation"""
        color_usage = defaultdict(int)
        color_contexts = defaultdict(list)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                # Analyze text colors
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    self._extract_text_colors(shape.text_frame, color_usage, color_contexts)
                
                # Analyze shape fill colors
                if hasattr(shape, 'fill'):
                    self._extract_fill_colors(shape, color_usage, color_contexts)
        
        # Identify primary color palette
        primary_colors = sorted(color_usage.items(), key=lambda x: x[1], reverse=True)[:10]
        
        return {
            'primary_palette': primary_colors,
            'color_usage': dict(color_usage),
            'color_contexts': dict(color_contexts),
            'total_unique_colors': len(color_usage)
        }
    
    def _extract_text_colors(self, text_frame: TextFrame, color_usage: dict, color_contexts: dict):
        """Extract color information from text"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                        rgb = run.font.color.rgb
                        color_key = f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
                        color_usage[color_key] += 1
                        color_contexts[color_key].append('text')
                except (AttributeError, TypeError):
                    # Skip colors that can't be extracted
                    continue
    
    def _extract_fill_colors(self, shape: BaseShape, color_usage: dict, color_contexts: dict):
        """Extract fill color information from shapes"""
        try:
            if hasattr(shape, 'fill') and shape.fill.type == 1:  # Solid fill
                if (hasattr(shape.fill, 'fore_color') and 
                    hasattr(shape.fill.fore_color, 'rgb') and 
                    shape.fill.fore_color.rgb):
                    rgb = shape.fill.fore_color.rgb
                    color_key = f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
                    color_usage[color_key] += 1
                    color_contexts[color_key].append('fill')
        except (AttributeError, TypeError, Exception):
            # Skip shapes where fill color cannot be determined
            pass
    
    def _analyze_layouts(self, prs: Presentation) -> Dict[str, Any]:
        """Analyze layout and positioning patterns"""
        positions = []
        sizes = []
        margins = {'left': [], 'top': [], 'right': [], 'bottom': []}
        
        for slide in prs.slides:
            for shape in slide.shapes:
                # Position analysis
                left_inches = shape.left.inches if shape.left else 0
                top_inches = shape.top.inches if shape.top else 0
                width_inches = shape.width.inches if shape.width else 0
                height_inches = shape.height.inches if shape.height else 0
                
                positions.append((left_inches, top_inches))
                sizes.append((width_inches, height_inches))
                
                margins['left'].append(left_inches)
                margins['top'].append(top_inches)
        
        # Calculate layout patterns
        if positions:
            avg_margins = {
                'left': np.mean(margins['left']) if margins['left'] else 0,
                'top': np.mean(margins['top']) if margins['top'] else 0
            }
            
            # Find common positioning patterns using clustering
            common_positions = self._find_common_positions(positions)
            common_sizes = self._find_common_sizes(sizes)
        else:
            avg_margins = {'left': 0, 'top': 0}
            common_positions = []
            common_sizes = []
        
        return {
            'average_margins': avg_margins,
            'common_positions': common_positions,
            'common_sizes': common_sizes,
            'total_shapes': len(positions)
        }
    
    def _find_common_positions(self, positions: List[Tuple[float, float]]) -> List[Tuple[float, float]]:
        """Use clustering to find common positioning patterns"""
        if len(positions) < 2:
            return positions
        
        try:
            # Normalize positions for clustering
            positions_array = np.array(positions)
            scaler = StandardScaler()
            positions_scaled = scaler.fit_transform(positions_array)
            
            # Use k-means to find 3-5 common positions
            n_clusters = min(5, len(positions))
            kmeans = KMeans(n_clusters=n_clusters, random_state=42, n_init=10)
            clusters = kmeans.fit(positions_scaled)
            
            # Transform centroids back to original space
            centroids_original = scaler.inverse_transform(clusters.cluster_centers_)
            return [tuple(centroid) for centroid in centroids_original]
        except Exception:
            # Fallback: return most common positions manually
            position_counter = Counter(positions)
            return [pos for pos, count in position_counter.most_common(5)]
    
    def _find_common_sizes(self, sizes: List[Tuple[float, float]]) -> List[Tuple[float, float]]:
        """Find common size patterns"""
        if not sizes:
            return []
        
        size_counter = Counter(sizes)
        return [size for size, count in size_counter.most_common(5)]
    
    def _analyze_text_hierarchy(self, prs: Presentation) -> Dict[str, Any]:
        """Analyze text hierarchy patterns (title, subtitle, body, etc.)"""
        hierarchy_patterns = {
            'title': {'fonts': defaultdict(int), 'sizes': [], 'positions': []},
            'subtitle': {'fonts': defaultdict(int), 'sizes': [], 'positions': []},
            'body': {'fonts': defaultdict(int), 'sizes': [], 'positions': []},
            'bullet': {'fonts': defaultdict(int), 'sizes': [], 'positions': []}
        }
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # Classify text based on position and size
                    text_type = self._classify_text_type(shape)
                    
                    if text_type in hierarchy_patterns:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name:
                                    hierarchy_patterns[text_type]['fonts'][run.font.name] += 1
                                if run.font.size:
                                    hierarchy_patterns[text_type]['sizes'].append(run.font.size.pt)
                        
                        hierarchy_patterns[text_type]['positions'].append(
                            (shape.left.inches if shape.left else 0, 
                             shape.top.inches if shape.top else 0)
                        )
        
        # Process patterns to find most common styles for each hierarchy level
        processed_hierarchy = {}
        for text_type, patterns in hierarchy_patterns.items():
            if patterns['fonts']:
                processed_hierarchy[text_type] = {
                    'primary_font': max(patterns['fonts'].items(), key=lambda x: x[1])[0],
                    'common_sizes': Counter(patterns['sizes']).most_common(3),
                    'typical_positions': patterns['positions'][:5]  # Top 5 positions
                }
        
        return processed_hierarchy
    
    def _classify_text_type(self, shape: BaseShape) -> str:
        """Classify text based on position, size, and content"""
        if not hasattr(shape, 'text_frame') or not shape.text_frame:
            return 'unknown'
        
        # Simple heuristic classification
        top_inches = shape.top.inches if shape.top else 0
        text_length = len(shape.text_frame.text)
        
        # Title: usually at top and shorter
        if top_inches < 2 and text_length < 100:
            return 'title'
        # Subtitle: near top but longer than title
        elif top_inches < 3 and text_length < 200:
            return 'subtitle'
        # Body: longer text, middle area
        elif text_length > 50:
            return 'body'
        # Default to body
        else:
            return 'body'
    
    def _analyze_theme(self, prs: Presentation) -> Dict[str, Any]:
        """Analyze theme colors and overall design theme"""
        # This is a simplified theme analysis
        # In a full implementation, we would parse the theme XML
        return {
            'slide_width': prs.slide_width.inches,
            'slide_height': prs.slide_height.inches,
            'aspect_ratio': round(prs.slide_width.inches / prs.slide_height.inches, 2),
            'total_slides': len(prs.slides)
        }
    
    def _analyze_shapes(self, prs: Presentation) -> Dict[str, Any]:
        """Analyze shape usage patterns"""
        shape_types = defaultdict(int)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                shape_type = str(type(shape).__name__)
                shape_types[shape_type] += 1
        
        return {
            'shape_distribution': dict(shape_types),
            'total_shapes': sum(shape_types.values())
        }
    
    def _calculate_consistency_score(self, analysis_results: Dict[str, Any]) -> float:
        """Calculate a consistency score (0-1) based on style analysis"""
        score_factors = []
        
        # Font consistency (fewer fonts = higher score)
        font_count = len(analysis_results['fonts']['font_usage'])
        font_score = max(0, 1 - (font_count - 1) * 0.1)  # Penalty for each additional font
        score_factors.append(font_score)
        
        # Color consistency (reasonable palette size)
        color_count = analysis_results['colors']['total_unique_colors']
        color_score = max(0, 1 - (color_count - 5) * 0.05)  # Penalty after 5 colors
        score_factors.append(color_score)
        
        # Size consistency (fewer size variations = higher score)
        size_count = len(analysis_results['fonts']['size_patterns'])
        size_score = max(0, 1 - (size_count - 3) * 0.1)  # Penalty after 3 sizes
        score_factors.append(size_score)
        
        return round(np.mean(score_factors), 2)
    
    def create_style_profile(self, analysis_results: Dict[str, Any], profile_name: str = None) -> str:
        """Create a reusable style profile from analysis results"""
        if not profile_name:
            profile_name = f"profile_{len(self.style_profiles)}"
        
        try:
            # Extract font hierarchy
            font_data = analysis_results['fonts']
            hierarchy_data = analysis_results.get('text_hierarchy', {})
            
            # Create text hierarchy profile
            text_hierarchy = self._create_text_hierarchy_profile(font_data, hierarchy_data)
            
            # Create color palette
            color_palette = self._create_color_palette(analysis_results['colors'])
            
            # Create layout patterns
            layout_patterns = self._create_layout_patterns(analysis_results['layouts'])
            
            # Create complete style profile
            style_profile = StyleProfile(
                name=profile_name,
                description=f"Style profile extracted from {analysis_results['file_path']}",
                font_hierarchy=text_hierarchy,
                color_palette=color_palette,
                layout_patterns=layout_patterns,
                slide_layouts={},  # Will be populated in advanced implementation
                theme_colors={},   # Will be populated in advanced implementation
                created_from=analysis_results['file_path'],
                confidence_score=analysis_results['consistency_score']
            )
            
            self.style_profiles[profile_name] = style_profile
            logger.info(f"Created style profile: {profile_name}")
            return profile_name
            
        except Exception as e:
            logger.error(f"Error creating style profile: {e}")
            raise
    
    def _create_text_hierarchy_profile(self, font_data: Dict, hierarchy_data: Dict) -> TextHierarchy:
        """Create text hierarchy from analyzed data"""
        primary_font = font_data.get('primary_font', 'Calibri')
        common_sizes = font_data.get('common_sizes', [(18, 1)])
        
        # Default font profile
        default_size = common_sizes[0][0] if common_sizes else 18
        
        title_style = FontProfile(
            family=primary_font,
            size_pt=hierarchy_data.get('title', {}).get('common_sizes', [(24, 1)])[0][0] if hierarchy_data.get('title', {}).get('common_sizes') else default_size + 6,
            bold=True,
            italic=False,
            underline=False,
            color_rgb=None
        )
        
        body_style = FontProfile(
            family=primary_font,
            size_pt=default_size,
            bold=False,
            italic=False,
            underline=False,
            color_rgb=None
        )
        
        return TextHierarchy(
            title_style=title_style,
            subtitle_style=None,  # Will be enhanced in advanced implementation
            body_style=body_style,
            bullet_style=None,
            caption_style=None
        )
    
    def _create_color_palette(self, color_data: Dict) -> List[ColorProfile]:
        """Create color palette from analyzed data"""
        primary_colors = color_data.get('primary_palette', [])
        color_palette = []
        
        for color_hex, frequency in primary_colors[:10]:  # Top 10 colors
            # Convert hex to RGB
            hex_clean = color_hex.lstrip('#')
            rgb = tuple(int(hex_clean[i:i+2], 16) for i in (0, 2, 4))
            
            # Determine context (simplified)
            contexts = color_data.get('color_contexts', {}).get(color_hex, ['unknown'])
            primary_context = max(set(contexts), key=contexts.count)
            
            color_profile = ColorProfile(
                rgb=rgb,
                hex_code=color_hex,
                usage_context=primary_context,
                frequency=frequency
            )
            color_palette.append(color_profile)
        
        return color_palette
    
    def _create_layout_patterns(self, layout_data: Dict) -> LayoutProfile:
        """Create layout patterns from analyzed data"""
        margins = layout_data.get('average_margins', {'left': 1, 'top': 1})
        common_sizes = layout_data.get('common_sizes', [])
        
        return LayoutProfile(
            average_left_margin=margins['left'],
            average_top_margin=margins['top'],
            common_widths=[size[0] for size in common_sizes],
            common_heights=[size[1] for size in common_sizes],
            grid_alignment={},  # Will be enhanced
            shape_distribution={}  # Will be enhanced
        )
    
    def save_style_profile(self, profile_name: str, file_path: str) -> bool:
        """Save a style profile to JSON file"""
        if profile_name not in self.style_profiles:
            raise ValueError(f"Style profile '{profile_name}' not found")
        
        try:
            profile = self.style_profiles[profile_name]
            profile_dict = asdict(profile)
            
            with open(file_path, 'w') as f:
                json.dump(profile_dict, f, indent=2, default=str)
            
            logger.info(f"Saved style profile '{profile_name}' to {file_path}")
            return True
        except Exception as e:
            logger.error(f"Error saving style profile: {e}")
            raise
    
    def load_style_profile(self, file_path: str) -> str:
        """Load a style profile from JSON file"""
        try:
            with open(file_path, 'r') as f:
                profile_dict = json.load(f)
            
            # Reconstruct the StyleProfile object
            profile_name = profile_dict['name']
            
            # This would need proper deserialization in full implementation
            # For now, store as dict
            self.style_profiles[profile_name] = profile_dict
            
            logger.info(f"Loaded style profile '{profile_name}' from {file_path}")
            return profile_name
        except Exception as e:
            logger.error(f"Error loading style profile: {e}")
            raise
    
    def get_style_profile(self, profile_name: str) -> Optional[StyleProfile]:
        """Get a style profile by name"""
        return self.style_profiles.get(profile_name)
    
    def list_style_profiles(self) -> List[str]:
        """List all available style profile names"""
        return list(self.style_profiles.keys())

# Example usage and testing
if __name__ == "__main__":
    analyzer = StyleAnalyzer()
    
    # This would be used to analyze an existing presentation
    # analysis = analyzer.analyze_presentation_style("example.pptx")
    # profile_name = analyzer.create_style_profile(analysis, "corporate_template")
    # analyzer.save_style_profile(profile_name, "corporate_style.json")
    
    print("Style Analysis Engine initialized successfully") 
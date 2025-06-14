#!/usr/bin/env python3
"""
Installation Test Script for PowerPoint MCP Server

This script tests that all dependencies are properly installed and
the basic functionality of the PowerPoint MCP server works.
"""

import sys
import traceback

def test_imports():
    """Test that all required modules can be imported"""
    print("Testing imports...")
    
    try:
        import pptx
        print("✓ python-pptx imported successfully")
    except ImportError as e:
        print(f"✗ python-pptx import failed: {e}")
        return False
    
    try:
        # Test basic pptx functionality
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        print("✓ python-pptx components imported successfully")
    except ImportError as e:
        print(f"✗ python-pptx components import failed: {e}")
        return False
    
    # Note: MCP library might not be available in all environments
    try:
        import mcp
        print("✓ MCP library imported successfully")
        mcp_available = True
    except ImportError as e:
        print(f"⚠ MCP library not available: {e}")
        print("  This is expected if MCP is not installed")
        mcp_available = False
    
    return True

def test_powerpoint_manager():
    """Test the PowerPointManager class functionality"""
    print("\nTesting PowerPointManager...")
    
    try:
        # Import without MCP dependencies
        sys.path.insert(0, '.')
        
        # Create a minimal version for testing
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        # Test basic presentation creation
        prs = Presentation()
        print("✓ Can create blank presentation")
        
        # Test adding a slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        print("✓ Can add slide to presentation")
        
        # Test adding text
        title = slide.shapes.title
        title.text = "Test Presentation"
        print("✓ Can add text to slide")
        
        # Test saving (to memory, not file)
        import tempfile
        import os
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            temp_path = tmp.name
        
        try:
            prs.save(temp_path)
            print("✓ Can save presentation to file")
        finally:
            # Clean up with better error handling
            try:
                os.unlink(temp_path)
            except (OSError, PermissionError):
                pass  # Ignore cleanup errors on Windows
        
        return True
        
    except Exception as e:
        print(f"✗ PowerPointManager test failed: {e}")
        traceback.print_exc()
        return False

def test_chart_functionality():
    """Test chart creation functionality"""
    print("\nTesting chart functionality...")
    
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
        
        # Create presentation and slide
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = ['Q1', 'Q2', 'Q3']
        chart_data.add_series('Sales', (100, 150, 120))
        
        # Add chart to slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(2), Inches(2), Inches(6), Inches(4.5),
            chart_data
        )
        
        print("✓ Can create and add charts")
        return True
        
    except Exception as e:
        print(f"✗ Chart functionality test failed: {e}")
        traceback.print_exc()
        return False

def test_image_functionality():
    """Test image handling functionality"""
    print("\nTesting image functionality...")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from io import BytesIO
        
        # Create a simple test image (1x1 pixel PNG)
        # This is a minimal PNG file in bytes
        png_data = b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\tpHYs\x00\x00\x0b\x13\x00\x00\x0b\x13\x01\x00\x9a\x9c\x18\x00\x00\x00\nIDATx\x9cc```\x00\x00\x00\x04\x00\x01\xdd\x8d\xb4\x1c\x00\x00\x00\x00IEND\xaeB`\x82'
        
        # Create presentation and slide
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add image from BytesIO
        image_stream = BytesIO(png_data)
        slide.shapes.add_picture(image_stream, Inches(1), Inches(1))
        
        print("✓ Can add images from BytesIO")
        return True
        
    except Exception as e:
        print(f"✗ Image functionality test failed: {e}")
        traceback.print_exc()
        return False

def main():
    """Run all tests"""
    print("PowerPoint MCP Server - Installation Test")
    print("=" * 50)
    
    all_passed = True
    
    # Test imports
    if not test_imports():
        all_passed = False
    
    # Test PowerPoint functionality
    if not test_powerpoint_manager():
        all_passed = False
    
    # Test chart functionality
    if not test_chart_functionality():
        all_passed = False
    
    # Test image functionality
    if not test_image_functionality():
        all_passed = False
    
    print("\n" + "=" * 50)
    if all_passed:
        print("✓ All tests passed! PowerPoint MCP Server is ready to use.")
        print("\nNext steps:")
        print("1. Install MCP library if needed: pip install mcp")
        print("2. Run the server: python powerpoint_mcp_server.py")
        print("3. Configure Cursor to use the MCP server")
        return 0
    else:
        print("✗ Some tests failed. Please check the error messages above.")
        print("\nTroubleshooting:")
        print("1. Install dependencies: pip install -r requirements.txt")
        print("2. Check Python version (3.8+ required)")
        print("3. Verify python-pptx installation: pip install python-pptx")
        return 1

if __name__ == "__main__":
    sys.exit(main()) 
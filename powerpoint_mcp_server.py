#!/usr/bin/env python3
"""
PowerPoint MCP Server

A Model Context Protocol server for PowerPoint presentation manipulation using python-pptx.
Provides tools for creating, editing, and managing PowerPoint presentations programmatically.

Based on research of existing MCP servers and python-pptx best practices.
Optimized for use with Cursor IDE and AI-assisted development workflows.
"""

import asyncio
import json
import logging
import os
import sys
import tempfile
import platform
import time
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
from urllib.request import urlopen
from urllib.parse import urlparse

# Configure logging early
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("powerpoint-mcp-server")

# Add Pydantic for input validation
try:
    from pydantic import BaseModel, field_validator, ValidationError
    PYDANTIC_AVAILABLE = True
except ImportError:
    PYDANTIC_AVAILABLE = False
    logger.warning("Pydantic not available - input validation will be disabled")

# Windows-specific COM imports for screenshot functionality
if platform.system() == "Windows":
    try:
        import win32com.client
        import pythoncom
        WIN32_COM_AVAILABLE = True
    except ImportError:
        WIN32_COM_AVAILABLE = False
        logger.warning("win32com not available - screenshot functionality will be disabled")
else:
    WIN32_COM_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError as e:
    print(f"python-pptx library not found: {e}")
    print("Please install with: pip install python-pptx")
    sys.exit(1)

try:
    from mcp.server import Server
    from mcp.server.models import InitializationOptions
    from mcp.server.stdio import stdio_server
    from mcp.types import (
        Resource,
        Tool,
        TextContent,
        ImageContent,
        EmbeddedResource,
        LoggingLevel
    )
except ImportError as e:
    print(f"MCP library not found: {e}")
    print("Please install with: pip install mcp")
    sys.exit(1)

# Import style analysis capabilities
try:
    from style_analysis import StyleAnalyzer
    STYLE_ANALYSIS_AVAILABLE = True
except ImportError as e:
    logger.warning(f"Style analysis not available: {e}")
    StyleAnalyzer = None
    STYLE_ANALYSIS_AVAILABLE = False

# =============================================================================
# LEAN PROMPT ARCHITECTURE
# =============================================================================

# Core system prompt - kept under 400 tokens for efficiency
CORE_SYSTEM_PROMPT = """You are a PowerPoint expert focused on professional, accessible presentations. 
Key principles: clean layouts, consistent styling, clear hierarchy, readable fonts. 
Always prioritize visual appeal and professional standards."""

# Focused tool guidance - concise, specific tips only
FOCUSED_TOOL_GUIDANCE = {
    "create_presentation": "Use blank slides by default. Apply consistent themes early.",
    "add_text_box": "Title text: 28-36pt, body: 18-24pt. Max 6 bullets per slide.",
    "add_image": "Maintain aspect ratio. Use placeholders for missing images.",
    "add_chart": "Choose chart type that tells the data story. Keep labels readable.",
    "create_color_palette": "Use high contrast. Consider accessibility and brand alignment.",
    "apply_typography_style": "Maintain hierarchy: title > heading > body > caption.",
    "create_master_slide_theme": "Design for consistency and reusability across slides.",
    "create_template": "Use clear placeholders. Test with varied data sets.",
    "bulk_generate_presentations": "Maintain quality standards while optimizing efficiency.",
    "save_presentation": "Verify file path accessibility before saving.",
    "load_presentation": "Check file format compatibility and integrity.",
    "add_slide": "Use layout index 6 (blank) for maximum flexibility."
}

# Success criteria generator
def get_success_criteria(tool_name: str, **kwargs) -> str:
    """Generate specific success criteria based on tool and parameters"""
    criteria = {
        "create_presentation": "Presentation object created and assigned unique ID.",
        "add_text_box": f"Text box added to slide with font size {kwargs.get('font_size', 18)}pt.",
        "add_image": f"Image positioned at ({kwargs.get('left', 1)}, {kwargs.get('top', 1)}) inches.",
        "add_chart": f"{kwargs.get('chart_type', 'chart').title()} chart with {len(kwargs.get('categories', []))} categories created.",
        "save_presentation": f"Presentation saved to {kwargs.get('file_path', 'specified path')}.",
        "load_presentation": f"Presentation loaded from {kwargs.get('file_path', 'file')} successfully.",
        "add_slide": f"New slide added using layout {kwargs.get('layout_index', 6)}."
    }
    return criteria.get(tool_name, f"{tool_name} operation completed successfully.")

# Generate focused prompt for specific tool operation
def get_focused_prompt(tool_name: str, **kwargs) -> str:
    """Generate lean, focused prompt for specific tool operation"""
    base = CORE_SYSTEM_PROMPT
    tool_specific = FOCUSED_TOOL_GUIDANCE.get(tool_name, "")
    success_criteria = get_success_criteria(tool_name, **kwargs)
    
    return f"{base}\n\n{tool_specific}\n\nSuccess criteria: {success_criteria}"

# =============================================================================
# ENHANCED SUCCESS MESSAGE FORMATTING
# =============================================================================

def format_success_message(tool_name: str, **kwargs) -> str:
    """Generate specific, actionable success messages instead of generic ones"""
    
    if tool_name == "add_text_box":
        slide_idx = kwargs.get('slide_index', 0)
        font_size = kwargs.get('font_size', 18)
        text_preview = kwargs.get('text', '')[:50] + ('...' if len(kwargs.get('text', '')) > 50 else '')
        bold_italic = []
        if kwargs.get('bold'): bold_italic.append('bold')
        if kwargs.get('italic'): bold_italic.append('italic')
        styling = f" ({', '.join(bold_italic)})" if bold_italic else ""
        
        return f"✅ Added text box to slide {slide_idx + 1}: \"{text_preview}\" ({font_size}pt{styling})"
    
    elif tool_name == "add_image":
        slide_idx = kwargs.get('slide_index', 0)
        image_source = kwargs.get('image_source', '')
        image_name = os.path.basename(image_source) if image_source else 'image'
        width = kwargs.get('width')
        height = kwargs.get('height')
        dimensions = f" ({width}×{height})" if width and height else ""
        
        return f"✅ Added image to slide {slide_idx + 1}: {image_name}{dimensions}"
    
    elif tool_name == "add_chart":
        slide_idx = kwargs.get('slide_index', 0)
        chart_type = kwargs.get('chart_type', 'chart')
        categories = kwargs.get('categories', [])
        series_data = kwargs.get('series_data', {})
        series_names = list(series_data.keys()) if series_data else []
        
        return f"✅ Added {chart_type} chart to slide {slide_idx + 1}: {len(categories)} categories, {len(series_names)} series ({', '.join(series_names[:3])}{'...' if len(series_names) > 3 else ''})"
    
    elif tool_name == "save_presentation":
        file_path = kwargs.get('file_path', '')
        file_name = os.path.basename(file_path) if file_path else 'presentation'
        file_size = ""
        try:
            if file_path and os.path.exists(file_path):
                size_bytes = os.path.getsize(file_path)
                if size_bytes > 1024*1024:
                    file_size = f" ({size_bytes/(1024*1024):.1f}MB)"
                else:
                    file_size = f" ({size_bytes/1024:.0f}KB)"
        except:
            pass
        
        return f"✅ Saved presentation: {file_name}{file_size} → Ready for use!"
    
    elif tool_name == "create_presentation":
        prs_id = kwargs.get('presentation_id', 'new')
        template_info = f" from template" if kwargs.get('template_path') else ""
        
        return f"✅ Created presentation {prs_id}{template_info} → Ready to add slides!"
    
    elif tool_name == "add_slide":
        slide_count = kwargs.get('slide_count', 'unknown')
        layout_type = kwargs.get('layout_index', 6)
        layout_name = {6: 'blank', 0: 'title', 1: 'title+content'}.get(layout_type, f'layout_{layout_type}')
        
        return f"✅ Added slide {slide_count} ({layout_name} layout) → Ready for content!"
    
    # Default fallback
    return f"✅ {tool_name.replace('_', ' ').title()} completed successfully"

def format_chart_success(chart_type: str, slide_index: int, categories: list, series_names: list) -> str:
    """Format detailed chart creation success message"""
    category_summary = f"{len(categories)} periods" if len(categories) <= 10 else f"{len(categories)} periods (large dataset)"
    series_summary = f"{len(series_names)} series"
    
    if len(series_names) <= 3:
        series_detail = f"({', '.join(series_names)})"
    else:
        series_detail = f"({', '.join(series_names[:2])}, +{len(series_names)-2} more)"
    
    return f"✅ {chart_type.title()} chart added to slide {slide_index + 1}: {category_summary}, {series_summary} {series_detail}"

# =============================================================================
# RESULT QUALITY KNOBS - USER PREFERENCE SYSTEM
# =============================================================================

if PYDANTIC_AVAILABLE:
    class ResultPreferences(BaseModel):
        verbosity: str = "normal"  # "brief", "normal", "verbose"
        tone: str = "professional"  # "executive", "technical", "friendly", "professional"
        include_preview: bool = True
        include_tips: bool = False
        
        @field_validator('verbosity')
        @classmethod
        def validate_verbosity(cls, v):
            valid_levels = ["brief", "normal", "verbose"]
            if v not in valid_levels:
                raise ValueError(f"Verbosity must be one of: {', '.join(valid_levels)}")
            return v
            
        @field_validator('tone')
        @classmethod
        def validate_tone(cls, v):
            valid_tones = ["executive", "technical", "friendly", "professional"]
            if v not in valid_tones:
                raise ValueError(f"Tone must be one of: {', '.join(valid_tones)}")
            return v

def format_response_with_preferences(tool_name: str, base_message: str, preferences: Optional[Dict] = None, **kwargs) -> str:
    """Format response according to user preferences"""
    if not preferences:
        return base_message
    
    try:
        if PYDANTIC_AVAILABLE:
            prefs = ResultPreferences(**preferences)
        else:
            # Fallback without validation
            prefs = type('Prefs', (), {
                'verbosity': preferences.get('verbosity', 'normal'),
                'tone': preferences.get('tone', 'professional'),
                'include_preview': preferences.get('include_preview', True),
                'include_tips': preferences.get('include_tips', False)
            })()
    except:
        return base_message  # Return base message if preferences are invalid
    
    # Adjust verbosity
    if prefs.verbosity == "brief":
        # Strip details, keep only core message
        if ":" in base_message:
            brief_msg = base_message.split(":")[0] + " ✓"
        else:
            brief_msg = base_message.split("→")[0].strip() if "→" in base_message else base_message
        response = brief_msg
    elif prefs.verbosity == "verbose":
        # Add extra details and context
        response = base_message
        if tool_name == "add_text_box" and kwargs.get('text'):
            response += f"\n💡 Tip: You can adjust positioning with left/top parameters"
        elif tool_name == "add_chart" and kwargs.get('categories'):
            response += f"\n💡 Tip: Try different chart types (column, bar, line, pie) for better data visualization"
        elif tool_name == "save_presentation":
            response += f"\n💡 Next: Open the file in PowerPoint to present or share"
    else:
        response = base_message
    
    # Adjust tone
    if prefs.tone == "executive":
        response = response.replace("✅", "✓").replace("💡 Tip:", "Note:")
        response = response.replace("Ready for", "Available for")
    elif prefs.tone == "friendly":
        response = response.replace("✅", "🎉").replace("added", "successfully added")
        if "💡" not in response and prefs.include_tips:
            response += " 🎊"
    elif prefs.tone == "technical":
        # Add technical details where available
        if tool_name == "add_chart" and kwargs.get('series_data'):
            data_points = sum(len(series) for series in kwargs['series_data'].values())
            response += f" [{data_points} data points total]"
    
    return response

# =============================================================================
# OPERATION LOGGING & INSTRUMENTATION
# =============================================================================

class OperationLogger:
    """Track operation metrics and usage patterns for continuous improvement"""
    
    def __init__(self, log_file: str = "mcp_operations.jsonl"):
        self.log_file = log_file
        self.session_start = datetime.utcnow()
        
    def log_operation(self, tool: str, success: bool, latency_ms: int, 
                     args_summary: str = "", error: str = None, **metadata):
        """Log an operation with timing and success metrics"""
        log_entry = {
            "timestamp": datetime.utcnow().isoformat(),
            "session_start": self.session_start.isoformat(),
            "tool": tool,
            "success": success,
            "latency_ms": latency_ms,
            "args_summary": args_summary,
            "error": error,
            **metadata
        }
        
        try:
            with open(self.log_file, "a", encoding='utf-8') as f:
                f.write(json.dumps(log_entry) + "\n")
        except Exception as e:
            logger.warning(f"Failed to write operation log: {e}")
    
    def log_validation_error(self, tool: str, error: str, args_summary: str = ""):
        """Log validation errors to track common input issues"""
        self.log_operation(
            tool=tool,
            success=False,
            latency_ms=0,
            args_summary=args_summary,
            error=f"Validation: {error}",
            error_type="validation"
        )
    
    def get_operation_summary(self, hours: int = 24) -> Dict[str, Any]:
        """Get operation summary for the last N hours"""
        try:
            cutoff_time = datetime.utcnow().timestamp() - (hours * 3600)
            
            operations = []
            if os.path.exists(self.log_file):
                with open(self.log_file, "r", encoding='utf-8') as f:
                    for line in f:
                        try:
                            entry = json.loads(line.strip())
                            entry_time = datetime.fromisoformat(entry['timestamp']).timestamp()
                            if entry_time >= cutoff_time:
                                operations.append(entry)
                        except:
                            continue
            
            if not operations:
                return {"message": "No operations in the specified time period"}
            
            # Calculate metrics
            total_ops = len(operations)
            successful_ops = sum(1 for op in operations if op['success'])
            success_rate = (successful_ops / total_ops) * 100 if total_ops > 0 else 0
            
            # Tool usage
            tool_usage = {}
            for op in operations:
                tool = op['tool']
                tool_usage[tool] = tool_usage.get(tool, 0) + 1
            
            # Error analysis
            errors = [op for op in operations if not op['success']]
            error_types = {}
            for error in errors:
                error_type = error.get('error_type', 'unknown')
                error_types[error_type] = error_types.get(error_type, 0) + 1
            
            # Latency stats
            latencies = [op['latency_ms'] for op in operations if op['latency_ms'] > 0]
            avg_latency = sum(latencies) / len(latencies) if latencies else 0
            
            return {
                "period_hours": hours,
                "total_operations": total_ops,
                "success_rate": round(success_rate, 1),
                "average_latency_ms": round(avg_latency, 1),
                "tool_usage": sorted(tool_usage.items(), key=lambda x: x[1], reverse=True),
                "error_types": error_types,
                "recent_errors": [{"tool": e['tool'], "error": e['error']} for e in errors[-5:]]
            }
        except Exception as e:
            return {"error": f"Failed to generate summary: {e}"}

# Global operation logger instance
operation_logger = OperationLogger()

def log_tool_execution(func):
    """Decorator to automatically log tool execution metrics"""
    async def wrapper(tool_name: str, *args, **kwargs):
        start_time = time.time()
        args_summary = f"args_count={len(kwargs)}"
        
        try:
            result = await func(tool_name, *args, **kwargs)
            latency_ms = int((time.time() - start_time) * 1000)
            
            # Create args summary
            key_args = []
            if 'slide_index' in kwargs:
                key_args.append(f"slide={kwargs['slide_index']}")
            if 'chart_type' in kwargs:
                key_args.append(f"type={kwargs['chart_type']}")
            if 'text' in kwargs:
                text_preview = kwargs['text'][:20] + "..." if len(kwargs['text']) > 20 else kwargs['text']
                key_args.append(f"text=\"{text_preview}\"")
            
            args_summary = ", ".join(key_args) if key_args else args_summary
            
            operation_logger.log_operation(
                tool=tool_name,
                success=True,
                latency_ms=latency_ms,
                args_summary=args_summary
            )
            
            return result
        except Exception as e:
            latency_ms = int((time.time() - start_time) * 1000)
            operation_logger.log_operation(
                tool=tool_name,
                success=False,
                latency_ms=latency_ms,
                args_summary=args_summary,
                error=str(e),
                error_type="execution"
            )
            raise
    
    return wrapper

# =============================================================================
# INPUT VALIDATION MODELS
# =============================================================================

if PYDANTIC_AVAILABLE:
    class CreatePresentationRequest(BaseModel):
        template_path: Optional[str] = None
        
        @field_validator('template_path')
        @classmethod
        def validate_template_path(cls, v):
            if v and not os.path.exists(v):
                raise ValueError(f"Template file not found: {v}")
            return v

    class AddTextBoxRequest(BaseModel):
        presentation_id: str
        slide_index: int
        text: str
        left: float = 1
        top: float = 1
        width: float = 8
        height: float = 1
        font_size: int = 18
        bold: bool = False
        italic: bool = False
        
        @field_validator('slide_index')
        @classmethod
        def validate_slide_index(cls, v):
            if v < 0:
                raise ValueError("Slide index cannot be negative")
            return v
            
        @field_validator('font_size')
        @classmethod
        def validate_font_size(cls, v):
            if v < 8 or v > 72:
                raise ValueError("Font size must be between 8 and 72 points")
            return v
            
        @field_validator('text')
        @classmethod
        def validate_text(cls, v):
            if not v.strip():
                raise ValueError("Text cannot be empty")
            return v

    class AddImageRequest(BaseModel):
        presentation_id: str
        slide_index: int
        image_source: str
        left: float = 1
        top: float = 1
        width: Optional[float] = None
        height: Optional[float] = None
        
        @field_validator('slide_index')
        @classmethod
        def validate_slide_index(cls, v):
            if v < 0:
                raise ValueError("Slide index cannot be negative")
            return v
            
        @field_validator('image_source')
        @classmethod
        def validate_image_source(cls, v):
            # Check if it's a URL or local file
            if v.startswith(('http://', 'https://')):
                return v  # URL validation could be added here
            elif os.path.exists(v):
                # Check if it's a valid image file
                valid_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg']
                if not any(v.lower().endswith(ext) for ext in valid_extensions):
                    raise ValueError(f"Unsupported image format. Supported: {', '.join(valid_extensions)}")
                return v
            else:
                raise ValueError(f"Image file not found: {v}")

    class AddChartRequest(BaseModel):
        presentation_id: str
        slide_index: int
        chart_type: str
        categories: List[str]
        series_data: Dict[str, List[float]]
        left: float = 2
        top: float = 2
        width: float = 6
        height: float = 4.5
        
        @field_validator('slide_index')
        @classmethod
        def validate_slide_index(cls, v):
            if v < 0:
                raise ValueError("Slide index cannot be negative")
            return v
            
        @field_validator('chart_type')
        @classmethod
        def validate_chart_type(cls, v):
            valid_types = ["column", "bar", "line", "pie"]
            if v not in valid_types:
                raise ValueError(f"Chart type must be one of: {', '.join(valid_types)}")
            return v
            
        @field_validator('categories')
        @classmethod
        def validate_categories(cls, v):
            if not v:
                raise ValueError("Categories cannot be empty")
            if len(v) > 20:
                raise ValueError("Maximum 20 categories allowed for readability")
            return v
            
        @field_validator('series_data')
        @classmethod
        def validate_series_data(cls, v):
            if not v:
                raise ValueError("Series data cannot be empty")
            if len(v) > 10:
                raise ValueError("Maximum 10 data series allowed for readability")
            
            # Check all series have same length
            lengths = [len(data) for data in v.values()]
            if len(set(lengths)) > 1:
                raise ValueError("All data series must have the same length")
            return v

    class SavePresentationRequest(BaseModel):
        presentation_id: str
        file_path: str
        
        @field_validator('file_path')
        @classmethod
        def validate_file_path(cls, v):
            # Check if path ends with .pptx
            if not v.lower().endswith('.pptx'):
                raise ValueError("File path must end with .pptx extension")
            
            # Note: Directory existence will be checked after path resolution in save_presentation
            # This allows for simple filenames to be resolved to workspace directory
            return v

    def validate_request(tool_name: str, arguments: Dict[str, Any]) -> Dict[str, Any]:
        """Validate request arguments using Pydantic models"""
        if not PYDANTIC_AVAILABLE:
            return arguments  # Skip validation if Pydantic not available
            
        validation_models = {
            "create_presentation": CreatePresentationRequest,
            "add_text_box": AddTextBoxRequest,
            "add_image": AddImageRequest,
            "add_chart": AddChartRequest,
            "save_presentation": SavePresentationRequest
        }
        
        model_class = validation_models.get(tool_name)
        if model_class:
            try:
                validated = model_class(**arguments)
                return validated.model_dump()
            except ValidationError as e:
                error_messages = []
                for error in e.errors():
                    field = " -> ".join(str(x) for x in error["loc"])
                    message = error["msg"]
                    error_messages.append(f"{field}: {message}")
                raise ValueError(f"Invalid input: {'; '.join(error_messages)}")
        
        return arguments  # No validation model found, return as-is

else:
    def validate_request(tool_name: str, arguments: Dict[str, Any]) -> Dict[str, Any]:
        """Fallback validation when Pydantic is not available"""
        return arguments

# Logger already configured above

# Server instance
server = Server("powerpoint-mcp-server")

class PowerPointManager:
    """Manages PowerPoint presentation operations"""
    
    def __init__(self):
        self.presentations: Dict[str, Presentation] = {}
        self.temp_files: List[str] = []
        # Initialize style analyzer if available
        self.style_analyzer = StyleAnalyzer() if STYLE_ANALYSIS_AVAILABLE else None
        
        # Phase 1: Professional Formatting & Layout
        self.layout_grids: Dict[str, Dict] = {}  # Store grid configurations per presentation
        self.color_palettes: Dict[str, Dict] = {}  # Store color palettes per presentation
        self.typography_profiles: Dict[str, Dict] = {}  # Store typography settings per presentation
        
        # Phase 2: Content Automation & Templates
        self.templates: Dict[str, Dict] = {}  # Store template definitions
        self.template_data_sources: Dict[str, Dict] = {}  # Store data source configurations
        self.generated_presentations: Dict[str, List[str]] = {}  # Track bulk generated presentations
        
        # Predefined professional color palettes
        self.predefined_palettes = {
            "corporate_blue": {
                "primary": RGBColor(0, 82, 147),
                "secondary": RGBColor(66, 139, 202),
                "accent": RGBColor(245, 245, 245),
                "text_dark": RGBColor(51, 51, 51),
                "text_light": RGBColor(255, 255, 255)
            },
            "modern_green": {
                "primary": RGBColor(46, 125, 50),
                "secondary": RGBColor(102, 187, 106),
                "accent": RGBColor(232, 245, 233),
                "text_dark": RGBColor(33, 33, 33),
                "text_light": RGBColor(255, 255, 255)
            },
            "professional_gray": {
                "primary": RGBColor(96, 125, 139),
                "secondary": RGBColor(144, 164, 174),
                "accent": RGBColor(236, 239, 241),
                "text_dark": RGBColor(55, 71, 79),
                "text_light": RGBColor(255, 255, 255)
            }
        }
        
        # Professional shape library (using valid MSO_SHAPE enum values)
        self.shape_library = {
            "arrows": [
                MSO_SHAPE.BLOCK_ARC,
                MSO_SHAPE.LEFT_ARROW,
                MSO_SHAPE.RIGHT_ARROW,
                MSO_SHAPE.UP_ARROW,
                MSO_SHAPE.DOWN_ARROW
            ],
            "callouts": [
                MSO_SHAPE.RECTANGULAR_CALLOUT,
                MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
                MSO_SHAPE.OVAL_CALLOUT,
                MSO_SHAPE.CLOUD_CALLOUT
            ],
            "geometric": [
                MSO_SHAPE.RECTANGLE,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                MSO_SHAPE.OVAL,
                MSO_SHAPE.RIGHT_TRIANGLE,  # Using RIGHT_TRIANGLE instead of TRIANGLE
                MSO_SHAPE.DIAMOND,
                MSO_SHAPE.HEXAGON,
                MSO_SHAPE.OCTAGON
            ]
        }
    
    def get_focused_prompt_for_operation(self, tool_name: str, **kwargs) -> str:
        """Get lean, focused prompt for specific tool operation"""
        return get_focused_prompt(tool_name, **kwargs)
    
    def log_operation_guidance(self, tool_name: str, **kwargs):
        """Log guidance for the current operation"""
        guidance = self.get_focused_prompt_for_operation(tool_name, **kwargs)
        logger.info(f"Operation guidance for {tool_name}: {guidance.strip()}")
    
    def create_presentation(self, template_path: Optional[str] = None) -> str:
        """Create a new presentation, optionally from a template"""
        try:
            # Log guidance for presentation creation
            self.log_operation_guidance("create_presentation", template_path=template_path)
            
            if template_path and os.path.exists(template_path):
                prs = Presentation(template_path)
                logger.info(f"Created presentation from template: {template_path}")
            else:
                prs = Presentation()
                logger.info("Created new blank presentation")
            
            # Generate unique ID for this presentation
            prs_id = f"prs_{len(self.presentations)}"
            self.presentations[prs_id] = prs
            return prs_id
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            raise
    
    def load_presentation(self, file_path: str) -> str:
        """Load an existing presentation"""
        try:
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"Presentation file not found: {file_path}")
            
            prs = Presentation(file_path)
            prs_id = f"prs_{len(self.presentations)}"
            self.presentations[prs_id] = prs
            logger.info(f"Loaded presentation: {file_path}")
            return prs_id
        except Exception as e:
            logger.error(f"Error loading presentation: {e}")
            raise
    
    def save_presentation(self, prs_id: str, file_path: str) -> bool:
        """Save a presentation to file"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            # Handle relative paths and simple filenames by making them relative to workspace
            resolved_path = self._resolve_file_path(file_path)
            
            # Ensure the target directory exists
            target_dir = os.path.dirname(resolved_path)
            if target_dir and not os.path.exists(target_dir):
                os.makedirs(target_dir, exist_ok=True)
                logger.info(f"Created directory: {target_dir}")
            
            self.presentations[prs_id].save(resolved_path)
            logger.info(f"Saved presentation {prs_id} to {resolved_path}")
            return True
        except Exception as e:
            logger.error(f"Error saving presentation: {e}")
            raise

    async def save_presentation_async(self, prs_id: str, file_path: str) -> bool:
        """Save a presentation to file asynchronously to prevent blocking"""
        return await asyncio.to_thread(self.save_presentation, prs_id, file_path)
    
    def add_slide(self, prs_id: str, layout_index: int = 6) -> int:
        """Add a new slide to the presentation"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            prs = self.presentations[prs_id]
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            slide_index = len(prs.slides) - 1
            
            logger.info(f"Added slide {slide_index} to presentation {prs_id}")
            return slide_index
        except Exception as e:
            logger.error(f"Error adding slide: {e}")
            raise
    
    def add_text_box(self, prs_id: str, slide_index: int, text: str, 
                     left: float = 1, top: float = 1, width: float = 8, height: float = 1,
                     font_size: int = 18, bold: bool = False, italic: bool = False) -> bool:
        """Add a text box to a slide"""
        try:
            # Log guidance for text formatting
            self.log_operation_guidance("add_text_box", font_size=font_size, text_length=len(text))
            
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            prs = self.presentations[prs_id]
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide index {slide_index} out of range")
            
            slide = prs.slides[slide_index]
            txt_box = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            txt_frame = txt_box.text_frame
            txt_frame.text = text
            
            # Apply formatting
            for paragraph in txt_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = bold
                    run.font.italic = italic
            
            logger.info(f"Added text box to slide {slide_index} in presentation {prs_id}")
            
            # Apply post-processing to fix common issues
            self.post_process_slide(prs_id, slide_index)
            
            return True
        except Exception as e:
            logger.error(f"Error adding text box: {e}")
            raise
    
    def add_image(self, prs_id: str, slide_index: int, image_source: str,
                  left: float = 1, top: float = 1, width: Optional[float] = None, height: Optional[float] = None) -> bool:
        """Add an image to a slide from file path or URL"""
        try:
            # Log guidance for image addition
            self.log_operation_guidance("add_image", is_url=image_source.startswith(('http://', 'https://')))
            
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            prs = self.presentations[prs_id]
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide index {slide_index} out of range")
            
            slide = prs.slides[slide_index]
            
            # Handle URL or file path
            if image_source.startswith(('http://', 'https://')):
                # Download image from URL
                image_data = BytesIO(urlopen(image_source).read())
                if width and height:
                    slide.shapes.add_picture(image_data, Inches(left), Inches(top), Inches(width), Inches(height))
                else:
                    slide.shapes.add_picture(image_data, Inches(left), Inches(top))
            else:
                # Load from file path
                if not os.path.exists(image_source):
                    raise FileNotFoundError(f"Image file not found: {image_source}")
                
                if width and height:
                    slide.shapes.add_picture(image_source, Inches(left), Inches(top), Inches(width), Inches(height))
                else:
                    slide.shapes.add_picture(image_source, Inches(left), Inches(top))
            
            logger.info(f"Added image to slide {slide_index} in presentation {prs_id}")
            return True
        except Exception as e:
            logger.error(f"Error adding image: {e}")
            raise
    
    def add_chart(self, prs_id: str, slide_index: int, chart_type: str, 
                  categories: List[str], series_data: Dict[str, List[float]],
                  left: float = 2, top: float = 2, width: float = 6, height: float = 4.5) -> bool:
        """Add a chart to a slide"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            prs = self.presentations[prs_id]
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide index {slide_index} out of range")
            
            slide = prs.slides[slide_index]
            
            # Create chart data
            chart_data = CategoryChartData()
            chart_data.categories = categories
            
            for series_name, data in series_data.items():
                chart_data.add_series(series_name, data)
            
            # Map chart type string to enum
            chart_type_map = {
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE
            }
            
            chart_type_enum = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            slide.shapes.add_chart(
                chart_type_enum,
                Inches(left), Inches(top), Inches(width), Inches(height),
                chart_data
            )
            
            logger.info(f"Added {chart_type} chart to slide {slide_index} in presentation {prs_id}")
            return True
        except Exception as e:
            logger.error(f"Error adding chart: {e}")
            raise
    
    def extract_text(self, prs_id: str) -> List[Dict[str, Any]]:
        """Extract all text content from a presentation"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            prs = self.presentations[prs_id]
            extracted_content = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_content = {
                    'slide_index': slide_idx,
                    'text_content': []
                }
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = ''.join(run.text for run in paragraph.runs)
                            if text.strip():
                                slide_content['text_content'].append(text.strip())
                
                extracted_content.append(slide_content)
            
            logger.info(f"Extracted text from presentation {prs_id}")
            return extracted_content
        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            raise
    
    def get_presentation_info(self, prs_id: str) -> Dict[str, Any]:
        """Get information about a presentation"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            prs = self.presentations[prs_id]
            
            info = {
                'presentation_id': prs_id,
                'slide_count': len(prs.slides),
                'slide_layouts_count': len(prs.slide_layouts),
                'slide_master_count': len(prs.slide_masters)
            }
            
            return info
        except Exception as e:
            logger.error(f"Error getting presentation info: {e}")
            raise
    
    def screenshot_slides(self, file_path: str, output_dir: Optional[str] = None, 
                         image_format: str = "PNG", width: int = 1920, height: int = 1080) -> List[str]:
        """Screenshot each slide of a PowerPoint presentation (Windows only)
        
        Args:
            file_path: Path to the PowerPoint file
            output_dir: Directory to save screenshots (defaults to temp directory)
            image_format: Image format (PNG, JPG, etc.)
            width: Screenshot width in pixels
            height: Screenshot height in pixels
            
        Returns:
            List of paths to the generated screenshot files
        """
        if not WIN32_COM_AVAILABLE:
            raise RuntimeError("Screenshot feature is only available on Windows with win32com installed")
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            # Initialize COM
            pythoncom.CoInitialize()
            
            # Create PowerPoint application instance
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = True  # Make visible for screenshot
            
            # Open the presentation
            presentation = ppt_app.Presentations.Open(os.path.abspath(file_path))
            
            # Set up output directory
            if output_dir is None:
                output_dir = tempfile.mkdtemp(prefix="ppt_screenshots_")
            else:
                os.makedirs(output_dir, exist_ok=True)
            
            screenshot_paths = []
            
            # Export each slide as image
            for i, slide in enumerate(presentation.Slides):
                slide_num = i + 1
                output_file = os.path.join(output_dir, f"slide_{slide_num:03d}.{image_format.lower()}")
                
                # Export slide as image
                slide.Export(output_file, image_format, width, height)
                screenshot_paths.append(output_file)
                
                logger.info(f"Exported slide {slide_num} to {output_file}")
            
            # Close presentation and quit PowerPoint
            presentation.Close()
            ppt_app.Quit()
            
            # Add to temp files for cleanup if using temp directory
            if output_dir.startswith(tempfile.gettempdir()):
                self.temp_files.extend(screenshot_paths)
                self.temp_files.append(output_dir)
            
            logger.info(f"Successfully created {len(screenshot_paths)} slide screenshots")
            return screenshot_paths
            
        except Exception as e:
            logger.error(f"Error creating slide screenshots: {e}")
            # Try to cleanup COM objects
            try:
                if 'presentation' in locals():
                    presentation.Close()
                if 'ppt_app' in locals():
                    ppt_app.Quit()
            except:
                pass
            finally:
                pythoncom.CoUninitialize()
            raise
        finally:
            # Ensure COM is uninitialized
            try:
                pythoncom.CoUninitialize()
            except:
                pass

    async def screenshot_slides_async(self, file_path: str, output_dir: Optional[str] = None, 
                                    image_format: str = "PNG", width: int = 1920, height: int = 1080) -> List[str]:
        """Take screenshots asynchronously to prevent blocking the event loop"""
        return await asyncio.to_thread(self.screenshot_slides, file_path, output_dir, image_format, width, height)
    
    def analyze_presentation_style(self, file_path: str) -> Dict[str, Any]:
        """Analyze presentation style patterns for learning and reuse"""
        if not STYLE_ANALYSIS_AVAILABLE:
            raise RuntimeError("Style analysis not available. Please install required dependencies.")
        
        try:
            return self.style_analyzer.analyze_presentation_style(file_path)
        except Exception as e:
            logger.error(f"Error analyzing presentation style: {e}")
            raise
    
    def create_style_profile(self, analysis_results: Dict[str, Any], profile_name: str = None) -> str:
        """Create a reusable style profile from analysis results"""
        if not STYLE_ANALYSIS_AVAILABLE:
            raise RuntimeError("Style analysis not available. Please install required dependencies.")
        
        try:
            return self.style_analyzer.create_style_profile(analysis_results, profile_name)
        except Exception as e:
            logger.error(f"Error creating style profile: {e}")
            raise
    
    def apply_style_profile(self, prs_id: str, profile_name: str) -> bool:
        """Apply a style profile to an existing presentation"""
        if not STYLE_ANALYSIS_AVAILABLE:
            raise RuntimeError("Style analysis not available. Please install required dependencies.")
        
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            profile = self.style_analyzer.get_style_profile(profile_name)
            if not profile:
                raise ValueError(f"Style profile '{profile_name}' not found")
            
            # This would contain the full implementation of applying styles
            # For now, this is a placeholder that shows the intended functionality
            logger.info(f"Applying style profile '{profile_name}' to presentation {prs_id}")
            
            # TODO: Implement style application logic
            # - Apply font hierarchy to text elements
            # - Apply color palette to shapes and text
            # - Adjust layout patterns and positioning
            # - Ensure consistency across slides
            
            return True
        except Exception as e:
            logger.error(f"Error applying style profile: {e}")
            raise
    
    def save_style_profile(self, profile_name: str, file_path: str) -> bool:
        """Save a style profile to JSON file"""
        if not STYLE_ANALYSIS_AVAILABLE:
            raise RuntimeError("Style analysis not available. Please install required dependencies.")
        
        try:
            return self.style_analyzer.save_style_profile(profile_name, file_path)
        except Exception as e:
            logger.error(f"Error saving style profile: {e}")
            raise
    
    def load_style_profile(self, file_path: str) -> str:
        """Load a style profile from JSON file"""
        if not STYLE_ANALYSIS_AVAILABLE:
            raise RuntimeError("Style analysis not available. Please install required dependencies.")
        
        try:
            return self.style_analyzer.load_style_profile(file_path)
        except Exception as e:
            logger.error(f"Error loading style profile: {e}")
            raise
    
    def list_style_profiles(self) -> List[str]:
        """List all available style profiles"""
        if not STYLE_ANALYSIS_AVAILABLE:
            return []
        
        return self.style_analyzer.list_style_profiles()

    def post_process_slide(self, prs_id: str, slide_index: int):
        """Apply quality checks and fixes after slide operations"""
        if prs_id not in self.presentations:
            return
            
        try:
            prs = self.presentations[prs_id]
            if slide_index >= len(prs.slides):
                return
                
            slide = prs.slides[slide_index]
            
            # Fix 1: Clean up placeholder fills that cause green rectangles
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        if hasattr(shape, 'fill'):
                            shape.fill.background()
                except Exception:
                    pass  # Skip if shape doesn't support fill operations
            
            # Fix 2: Ensure proper bullet formatting in text boxes
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip().startswith('•'):
                                # Ensure paragraph level is set for bullet points
                                if not hasattr(paragraph, 'level') or paragraph.level is None:
                                    paragraph.level = 0
                except Exception:
                    pass  # Skip if paragraph doesn't support level operations
                    
            logger.info(f"Applied post-processing to slide {slide_index} in presentation {prs_id}")
            
        except Exception as e:
            logger.warning(f"Post-processing failed for slide {slide_index}: {e}")

    def cleanup(self):
        """Clean up temporary files"""
        for temp_file in self.temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as e:
                logger.warning(f"Could not remove temp file {temp_file}: {e}")
        self.temp_files.clear()

    # Phase 1: Grid-Based Positioning
    def create_layout_grid(self, prs_id: str, columns: int, rows: int, 
                          margins: Dict[str, float] = None) -> bool:
        """Create a layout grid for professional alignment and spacing"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            # Default margins in inches
            if margins is None:
                margins = {"left": 0.5, "right": 0.5, "top": 0.5, "bottom": 0.5}
            
            # Get slide dimensions (standard is 10x7.5 inches)
            slide_width = 10.0
            slide_height = 7.5
            
            # Calculate grid cell dimensions
            grid_width = slide_width - margins["left"] - margins["right"]
            grid_height = slide_height - margins["top"] - margins["bottom"]
            
            cell_width = grid_width / columns
            cell_height = grid_height / rows
            
            grid_config = {
                "columns": columns,
                "rows": rows,
                "margins": margins,
                "cell_width": cell_width,
                "cell_height": cell_height,
                "grid_width": grid_width,
                "grid_height": grid_height
            }
            
            self.layout_grids[prs_id] = grid_config
            logger.info(f"Created {columns}x{rows} grid for presentation {prs_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error creating layout grid: {e}")
            raise
    
    def snap_to_grid(self, prs_id: str, slide_index: int, shape_id: str, 
                     grid_position: tuple, alignment: str = "top-left") -> bool:
        """Snap a shape to grid position"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            if prs_id not in self.layout_grids:
                raise ValueError(f"No grid defined for presentation {prs_id}")
            
            grid = self.layout_grids[prs_id]
            col, row = grid_position
            
            if col >= grid["columns"] or row >= grid["rows"]:
                raise ValueError(f"Grid position ({col}, {row}) out of bounds")
            
            # Calculate position based on grid
            left = grid["margins"]["left"] + (col * grid["cell_width"])
            top = grid["margins"]["top"] + (row * grid["cell_height"])
            
            # Get the shape and update its position
            prs = self.presentations[prs_id]
            slide = prs.slides[slide_index]
            
            # Find shape by ID (assuming shape_id is the index for now)
            shape_index = int(shape_id)
            if shape_index >= len(slide.shapes):
                raise ValueError(f"Shape {shape_id} not found")
            
            shape = slide.shapes[shape_index]
            shape.left = Inches(left)
            shape.top = Inches(top)
            
            logger.info(f"Snapped shape {shape_id} to grid position ({col}, {row})")
            return True
            
        except Exception as e:
            logger.error(f"Error snapping to grid: {e}")
            raise
    
    def distribute_shapes(self, prs_id: str, slide_index: int, shape_ids: List[str], 
                         distribution_type: str) -> bool:
        """Distribute shapes evenly"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            prs = self.presentations[prs_id]
            slide = prs.slides[slide_index]
            
            shapes = []
            for shape_id in shape_ids:
                shape_index = int(shape_id)
                if shape_index < len(slide.shapes):
                    shapes.append(slide.shapes[shape_index])
            
            if len(shapes) < 2:
                raise ValueError("Need at least 2 shapes to distribute")
            
            if distribution_type == "horizontal":
                shapes.sort(key=lambda s: s.left)
                total_width = shapes[-1].left + shapes[-1].width - shapes[0].left
                space_between = (total_width - sum(s.width for s in shapes)) / (len(shapes) - 1)
                
                current_left = shapes[0].left
                for i, shape in enumerate(shapes[1:], 1):
                    current_left += shapes[i-1].width + space_between
                    shape.left = int(current_left)
                    
            elif distribution_type == "vertical":
                shapes.sort(key=lambda s: s.top)
                total_height = shapes[-1].top + shapes[-1].height - shapes[0].top
                space_between = (total_height - sum(s.height for s in shapes)) / (len(shapes) - 1)
                
                current_top = shapes[0].top
                for i, shape in enumerate(shapes[1:], 1):
                    current_top += shapes[i-1].height + space_between
                    shape.top = int(current_top)
            
            logger.info(f"Distributed {len(shapes)} shapes {distribution_type}ly")
            return True
            
        except Exception as e:
            logger.error(f"Error distributing shapes: {e}")
            raise

    # Phase 1: Color Palette Management
    def create_color_palette(self, prs_id: str, palette_name: str, 
                           colors: Dict[str, str] = None) -> bool:
        """Create a custom color palette for brand consistency"""
        try:
            # Log guidance for color palette creation
            self.log_operation_guidance("create_color_palette", palette_name=palette_name, has_custom_colors=colors is not None)
            
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            if colors is None and palette_name in self.predefined_palettes:
                # Use predefined palette
                palette = self.predefined_palettes[palette_name]
            elif colors:
                # Create custom palette from hex colors
                palette = {}
                for role, hex_color in colors.items():
                    hex_color = hex_color.lstrip('#')
                    rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                    palette[role] = RGBColor(*rgb)
            else:
                raise ValueError("Either use predefined palette or provide custom colors")
            
            if prs_id not in self.color_palettes:
                self.color_palettes[prs_id] = {}
            
            self.color_palettes[prs_id][palette_name] = palette
            logger.info(f"Created color palette '{palette_name}' for presentation {prs_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error creating color palette: {e}")
            raise
    
    def apply_color_palette(self, prs_id: str, palette_name: str) -> bool:
        """Apply color palette to presentation elements"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            if prs_id not in self.color_palettes or palette_name not in self.color_palettes[prs_id]:
                raise ValueError(f"Color palette '{palette_name}' not found")
            
            palette = self.color_palettes[prs_id][palette_name]
            prs = self.presentations[prs_id]
            
            # Apply colors to all slides
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        # Apply text colors
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if 'text_dark' in palette:
                                    run.font.color.rgb = palette['text_dark']
                    
                    if hasattr(shape, 'fill'):
                        # Apply shape fills
                        if 'primary' in palette:
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = palette['primary']
            
            logger.info(f"Applied color palette '{palette_name}' to presentation {prs_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error applying color palette: {e}")
            raise

    # Phase 1: Typography System
    def create_typography_profile(self, prs_id: str, profile_name: str, 
                                typography_config: Dict) -> bool:
        """Create a typography profile with style hierarchies"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            # Default typography hierarchy
            default_config = {
                "title": {"font_name": "Calibri", "font_size": 44, "bold": True, "color": "text_dark"},
                "subtitle": {"font_name": "Calibri", "font_size": 32, "bold": False, "color": "text_dark"},
                "heading": {"font_name": "Calibri", "font_size": 24, "bold": True, "color": "primary"},
                "body": {"font_name": "Calibri", "font_size": 18, "bold": False, "color": "text_dark"},
                "caption": {"font_name": "Calibri", "font_size": 14, "bold": False, "color": "secondary"}
            }
            
            # Merge with provided config
            config = {**default_config, **typography_config}
            
            if prs_id not in self.typography_profiles:
                self.typography_profiles[prs_id] = {}
            
            self.typography_profiles[prs_id][profile_name] = config
            logger.info(f"Created typography profile '{profile_name}' for presentation {prs_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error creating typography profile: {e}")
            raise
    
    def apply_typography_style(self, prs_id: str, slide_index: int, shape_id: str, 
                             style_type: str, profile_name: str = "default") -> bool:
        """Apply typography style to a text shape"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            if prs_id not in self.typography_profiles or profile_name not in self.typography_profiles[prs_id]:
                raise ValueError(f"Typography profile '{profile_name}' not found")
            
            profile = self.typography_profiles[prs_id][profile_name]
            if style_type not in profile:
                raise ValueError(f"Style type '{style_type}' not found in profile")
            
            style = profile[style_type]
            
            prs = self.presentations[prs_id]
            slide = prs.slides[slide_index]
            shape_index = int(shape_id)
            
            if shape_index >= len(slide.shapes):
                raise ValueError(f"Shape {shape_id} not found")
            
            shape = slide.shapes[shape_index]
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                raise ValueError("Shape does not contain text")
            
            # Apply typography styles
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = style["font_name"]
                    run.font.size = Pt(style["font_size"])
                    run.font.bold = style["bold"]
                    
                    # Apply color if palette is available
                    if prs_id in self.color_palettes:
                        for palette_name, palette in self.color_palettes[prs_id].items():
                            if style["color"] in palette:
                                run.font.color.rgb = palette[style["color"]]
                                break
            
            logger.info(f"Applied typography style '{style_type}' to shape {shape_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error applying typography style: {e}")
            raise

    # Phase 1: Shape Libraries
    def add_professional_shape(self, prs_id: str, slide_index: int, shape_category: str, 
                             shape_name: str, left: float = 1, top: float = 1, 
                             width: float = 2, height: float = 2) -> bool:
        """Add a professional shape from the shape library"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            if shape_category not in self.shape_library:
                raise ValueError(f"Shape category '{shape_category}' not found")
            
            prs = self.presentations[prs_id]
            slide = prs.slides[slide_index]
            
            # Find shape by name or use first shape in category
            shapes_in_category = self.shape_library[shape_category]
            if shape_name.isdigit():
                shape_index = int(shape_name)
                if shape_index >= len(shapes_in_category):
                    raise ValueError(f"Shape index {shape_index} out of range")
                shape_type = shapes_in_category[shape_index]
            else:
                # Try to find by name (simplified)
                shape_type = shapes_in_category[0]  # Default to first shape
            
            # Add the shape
            shape = slide.shapes.add_shape(
                shape_type,
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )
            
            logger.info(f"Added professional shape from category '{shape_category}' to slide {slide_index}")
            return True
            
        except Exception as e:
            logger.error(f"Error adding professional shape: {e}")
            raise
    
    def list_shape_library(self) -> Dict[str, List[str]]:
        """List available professional shapes"""
        try:
            # Convert MSO_SHAPE enum values to readable names
            readable_library = {}
            for category, shapes in self.shape_library.items():
                readable_library[category] = [str(shape).split('.')[-1] for shape in shapes]
            
            return readable_library
            
        except Exception as e:
            logger.error(f"Error listing shape library: {e}")
            raise

    # Phase 1: Master Slide Management
    def create_master_slide_theme(self, prs_id: str, theme_name: str, 
                                 theme_config: Dict[str, Any]) -> bool:
        """Create a master slide theme with consistent formatting"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            prs = self.presentations[prs_id]
            
            # Store theme configuration for later application
            if not hasattr(self, 'master_themes'):
                self.master_themes = {}
            if prs_id not in self.master_themes:
                self.master_themes[prs_id] = {}
            
            # Default theme configuration
            default_theme = {
                "background_color": "#ffffff",
                "title_font": {
                    "name": "Calibri",
                    "size": 44,
                    "color": "#333333",
                    "bold": True
                },
                "content_font": {
                    "name": "Calibri", 
                    "size": 18,
                    "color": "#333333",
                    "bold": False
                },
                "accent_color": "#0052a3",
                "layout_margins": {
                    "left": 0.75,
                    "right": 0.75,
                    "top": 0.75,
                    "bottom": 0.75
                }
            }
            
            # Merge with provided config
            merged_config = {**default_theme, **theme_config}
            self.master_themes[prs_id][theme_name] = merged_config
            
            logger.info(f"Created master slide theme '{theme_name}' for presentation {prs_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error creating master slide theme: {e}")
            raise
    
    def apply_master_theme(self, prs_id: str, theme_name: str) -> bool:
        """Apply master theme to all slides in presentation"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            if not hasattr(self, 'master_themes') or prs_id not in self.master_themes or theme_name not in self.master_themes[prs_id]:
                raise ValueError(f"Master theme '{theme_name}' not found")
            
            theme = self.master_themes[prs_id][theme_name]
            prs = self.presentations[prs_id]
            
            # Apply theme to all slides
            for slide in prs.slides:
                # Apply background color if specified
                if "background_color" in theme:
                    hex_color = theme["background_color"].lstrip('#')
                    rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(*rgb)
                
                # Apply typography to text shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        # Determine if this is likely a title or content based on position
                        is_title = shape.top < Inches(2)  # Rough heuristic
                        
                        font_config = theme["title_font"] if is_title else theme["content_font"]
                        
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = font_config["name"]
                                run.font.size = Pt(font_config["size"])
                                run.font.bold = font_config["bold"]
                                
                                # Apply color
                                hex_color = font_config["color"].lstrip('#')
                                rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                                run.font.color.rgb = RGBColor(*rgb)
            
            logger.info(f"Applied master theme '{theme_name}' to presentation {prs_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error applying master theme: {e}")
            raise
    
    def list_master_themes(self, prs_id: str) -> List[str]:
        """List available master themes for a presentation"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            if not hasattr(self, 'master_themes') or prs_id not in self.master_themes:
                return []
            
            return list(self.master_themes[prs_id].keys())
            
        except Exception as e:
            logger.error(f"Error listing master themes: {e}")
            raise

    def set_slide_layout_template(self, prs_id: str, slide_index: int, 
                                 template_config: Dict[str, Any]) -> bool:
        """Apply a layout template to a specific slide"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            prs = self.presentations[prs_id]
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide index {slide_index} out of range")
            
            slide = prs.slides[slide_index]
            
            # Clear existing content if specified
            if template_config.get("clear_existing", False):
                # Remove all shapes except placeholders
                shapes_to_remove = []
                for shape in slide.shapes:
                    if not hasattr(shape, 'is_placeholder') or not shape.is_placeholder:
                        shapes_to_remove.append(shape)
                
                for shape in shapes_to_remove:
                    slide.shapes._element.remove(shape._element)
            
            # Apply template layout
            layout_type = template_config.get("layout_type", "title_content")
            
            if layout_type == "title_content":
                # Add title
                if "title" in template_config:
                    title_box = slide.shapes.add_textbox(
                        Inches(0.75), Inches(0.75), Inches(8.5), Inches(1.5)
                    )
                    title_frame = title_box.text_frame
                    title_frame.text = template_config["title"]
                    
                    # Apply title formatting
                    for paragraph in title_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(32)
                            run.font.bold = True
                
                # Add content
                if "content" in template_config:
                    content_box = slide.shapes.add_textbox(
                        Inches(0.75), Inches(2.5), Inches(8.5), Inches(4.5)
                    )
                    content_frame = content_box.text_frame
                    content_frame.text = template_config["content"]
            
            elif layout_type == "two_column":
                # Left column
                if "left_content" in template_config:
                    left_box = slide.shapes.add_textbox(
                        Inches(0.75), Inches(1.5), Inches(4), Inches(5.5)
                    )
                    left_frame = left_box.text_frame
                    left_frame.text = template_config["left_content"]
                
                # Right column
                if "right_content" in template_config:
                    right_box = slide.shapes.add_textbox(
                        Inches(5.25), Inches(1.5), Inches(4), Inches(5.5)
                    )
                    right_frame = right_box.text_frame
                    right_frame.text = template_config["right_content"]
            
            logger.info(f"Applied layout template to slide {slide_index}")
            return True
            
        except Exception as e:
            logger.error(f"Error setting slide layout template: {e}")
            raise
    
    # Phase 2: Content Automation & Templates Methods
    
    def create_template(self, template_config: Dict[str, Any]) -> str:
        """Create a reusable template with placeholders and rules"""
        try:
            template_id = f"template_{len(self.templates)}"
            
            # Template configuration structure:
            # {
            #   "name": "Monthly Report Template",
            #   "description": "Template for monthly business reports",
            #   "slides": [
            #     {
            #       "layout_type": "title_slide",
            #       "elements": [
            #         {
            #           "type": "text",
            #           "content": "{{title}}",
            #           "position": {"left": 1, "top": 2},
            #           "formatting": {"font_size": 32, "bold": True}
            #         }
            #       ]
            #     }
            #   ],
            #   "conditional_logic": {...},
            #   "data_mapping": {...}
            # }
            
            self.templates[template_id] = {
                "config": template_config,
                "created_at": str(Path().cwd()),
                "usage_count": 0
            }
            
            logger.info(f"Created template {template_id}: {template_config.get('name', 'Unnamed')}")
            return template_id
            
        except Exception as e:
            logger.error(f"Error creating template: {e}")
            raise
    
    def apply_template(self, template_id: str, data: Dict[str, Any]) -> str:
        """Apply a template with data substitution to create a new presentation"""
        try:
            if template_id not in self.templates:
                raise ValueError(f"Template {template_id} not found")
            
            template = self.templates[template_id]
            template_config = template["config"]
            
            # Create new presentation
            prs_id = self.create_presentation()
            
            # Process each slide in the template
            for slide_config in template_config.get("slides", []):
                slide_index = self.add_slide(prs_id, layout_index=6)
                
                # Apply conditional logic
                if not self._should_include_slide(slide_config, data):
                    continue
                
                # Process elements on the slide
                for element in slide_config.get("elements", []):
                    self._process_template_element(prs_id, slide_index, element, data)
            
            # Update template usage count
            self.templates[template_id]["usage_count"] += 1
            
            logger.info(f"Applied template {template_id} to create presentation {prs_id}")
            return prs_id
            
        except Exception as e:
            logger.error(f"Error applying template: {e}")
            raise
    
    def update_template_content(self, prs_id: str, updates: Dict[str, Any]) -> bool:
        """Update content in an existing presentation using template data"""
        try:
            if prs_id not in self.presentations:
                raise ValueError(f"Presentation {prs_id} not found")
            
            prs = self.presentations[prs_id]
            
            # Apply updates to slides
            for slide_index, slide_updates in updates.items():
                if isinstance(slide_index, str) and slide_index.isdigit():
                    slide_index = int(slide_index)
                
                if slide_index >= len(prs.slides):
                    continue
                
                slide = prs.slides[slide_index]
                
                # Update text content
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        # Replace placeholders in existing text
                        original_text = shape.text_frame.text
                        updated_text = self._substitute_variables(original_text, slide_updates)
                        if updated_text != original_text:
                            shape.text_frame.text = updated_text
            
            logger.info(f"Updated content in presentation {prs_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error updating template content: {e}")
            raise
    
    def bulk_generate_presentations(self, template_id: str, data_sets: List[Dict[str, Any]], 
                                   output_config: Dict[str, Any] = None) -> List[str]:
        """Generate multiple presentations from a template and data sets"""
        try:
            if template_id not in self.templates:
                raise ValueError(f"Template {template_id} not found")
            
            generated_presentations = []
            
            for i, data_set in enumerate(data_sets):
                # Create presentation from template
                prs_id = self.apply_template(template_id, data_set)
                generated_presentations.append(prs_id)
                
                # Save if output configuration provided
                if output_config and output_config.get("auto_save", False):
                    output_path = output_config.get("output_path", "")
                    filename = f"{output_path}/presentation_{i+1}.pptx"
                    self.save_presentation(prs_id, filename)
            
            # Track bulk generation
            bulk_id = f"bulk_{len(self.generated_presentations)}"
            self.generated_presentations[bulk_id] = generated_presentations
            
            logger.info(f"Bulk generated {len(generated_presentations)} presentations")
            return generated_presentations
            
        except Exception as e:
            logger.error(f"Error in bulk generation: {e}")
            raise
    
    def map_data_source(self, source_config: Dict[str, Any]) -> str:
        """Configure a data source for template mapping"""
        try:
            source_id = f"source_{len(self.template_data_sources)}"
            
            # Data source configuration:
            # {
            #   "type": "json|csv|excel|api",
            #   "source": "file_path_or_url",
            #   "mapping": {
            #     "title": "data.report_title",
            #     "date": "data.report_date",
            #     "metrics": "data.performance_metrics"
            #   },
            #   "refresh_interval": 3600  # seconds
            # }
            
            self.template_data_sources[source_id] = source_config
            
            logger.info(f"Mapped data source {source_id}: {source_config.get('type', 'unknown')}")
            return source_id
            
        except Exception as e:
            logger.error(f"Error mapping data source: {e}")
            raise
    
    def list_templates(self) -> List[Dict[str, Any]]:
        """List all available templates"""
        try:
            template_list = []
            for template_id, template in self.templates.items():
                template_info = {
                    "id": template_id,
                    "name": template["config"].get("name", "Unnamed Template"),
                    "description": template["config"].get("description", ""),
                    "slides_count": len(template["config"].get("slides", [])),
                    "usage_count": template.get("usage_count", 0),
                    "created_at": template.get("created_at", "")
                }
                template_list.append(template_info)
            
            return template_list
            
        except Exception as e:
            logger.error(f"Error listing templates: {e}")
            raise
    
    def delete_template(self, template_id: str) -> bool:
        """Delete a template"""
        try:
            if template_id not in self.templates:
                raise ValueError(f"Template {template_id} not found")
            
            del self.templates[template_id]
            logger.info(f"Deleted template {template_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error deleting template: {e}")
            raise
    
    def _resolve_file_path(self, file_path: str) -> str:
        """Resolve file path to be relative to workspace directory instead of current working directory"""
        # If it's already an absolute path, use it as-is
        if os.path.isabs(file_path):
            return file_path
        
        # Get the directory name from the path
        directory = os.path.dirname(file_path)
        filename = os.path.basename(file_path)
        
        # If no directory specified (just a filename), use workspace directory
        if not directory:
            # Try to detect workspace directory from common indicators
            workspace_dir = self._detect_workspace_directory()
            resolved_path = os.path.join(workspace_dir, filename)
            logger.info(f"Resolved simple filename '{filename}' to workspace path: {resolved_path}")
            return resolved_path
        
        # For relative paths with directories, make them relative to workspace
        workspace_dir = self._detect_workspace_directory()
        resolved_path = os.path.join(workspace_dir, file_path)
        logger.info(f"Resolved relative path '{file_path}' to workspace path: {resolved_path}")
        return resolved_path
    
    def _detect_workspace_directory(self) -> str:
        """Detect the workspace directory by looking for common project indicators"""
        current_dir = os.getcwd()
        
        # Common project files that indicate a workspace root
        workspace_indicators = [
            'package.json', 'requirements.txt', '.git', 'README.md', 
            'pyproject.toml', 'setup.py', 'pom.xml', 'Cargo.toml'
        ]
        
        # Start from current directory and walk up the tree
        check_dir = current_dir
        while check_dir != os.path.dirname(check_dir):  # Stop at root
            for indicator in workspace_indicators:
                if os.path.exists(os.path.join(check_dir, indicator)):
                    logger.info(f"Detected workspace directory: {check_dir}")
                    return check_dir
            check_dir = os.path.dirname(check_dir)
        
        # Fallback to current directory if no workspace indicators found
        logger.info(f"Using current directory as workspace: {current_dir}")
        return current_dir

    def _substitute_variables(self, text: str, data: Dict[str, Any]) -> str:
        """Replace placeholders in text with actual data"""
        try:
            import re
            
            # Find all {{variable}} patterns
            pattern = r'\{\{([^}]+)\}\}'
            matches = re.findall(pattern, text)
            
            result_text = text
            
            for match in matches:
                variable_name = match.strip()
                
                # Support nested data access (e.g., data.user.name)
                value = self._get_nested_value(data, variable_name)
                
                if value is not None:
                    # Replace placeholder with actual value
                    placeholder = f"{{{{{variable_name}}}}}"
                    result_text = result_text.replace(placeholder, str(value))
            
            return result_text
            
        except Exception as e:
            logger.error(f"Error substituting variables: {e}")
            return text
    
    def _get_nested_value(self, data: Dict[str, Any], key_path: str) -> Any:
        """Get value from nested dictionary using dot notation"""
        try:
            keys = key_path.split('.')
            value = data
            
            for key in keys:
                if isinstance(value, dict) and key in value:
                    value = value[key]
                else:
                    return None
            
            return value
            
        except Exception:
            return None
    
    def _should_include_slide(self, slide_config: Dict[str, Any], data: Dict[str, Any]) -> bool:
        """Evaluate conditional logic to determine if slide should be included"""
        try:
            conditions = slide_config.get("conditional_logic", {})
            
            if not conditions:
                return True  # No conditions, include slide
            
            # Support various condition types
            if "if" in conditions:
                condition = conditions["if"]
                return self._evaluate_condition(condition, data)
            
            return True
            
        except Exception as e:
            logger.error(f"Error evaluating slide conditions: {e}")
            return True
    
    def _evaluate_condition(self, condition: Dict[str, Any], data: Dict[str, Any]) -> bool:
        """Evaluate a single condition"""
        try:
            if "field" in condition and "operator" in condition and "value" in condition:
                field_value = self._get_nested_value(data, condition["field"])
                operator = condition["operator"]
                expected_value = condition["value"]
                
                if operator == "equals":
                    return field_value == expected_value
                elif operator == "not_equals":
                    return field_value != expected_value
                elif operator == "greater_than":
                    return float(field_value) > float(expected_value)
                elif operator == "less_than":
                    return float(field_value) < float(expected_value)
                elif operator == "contains":
                    return expected_value in str(field_value)
                elif operator == "exists":
                    return field_value is not None
            
            return True
            
        except Exception as e:
            logger.error(f"Error evaluating condition: {e}")
            return True
    
    def _process_template_element(self, prs_id: str, slide_index: int, element: Dict[str, Any], data: Dict[str, Any]):
        """Process a single template element (text, image, chart, etc.)"""
        try:
            element_type = element.get("type", "text")
            
            if element_type == "text":
                # Process text element
                content = self._substitute_variables(element.get("content", ""), data)
                position = element.get("position", {})
                formatting = element.get("formatting", {})
                
                self.add_text_box(
                    prs_id, slide_index, content,
                    left=position.get("left", 1),
                    top=position.get("top", 1),
                    width=position.get("width", 8),
                    height=position.get("height", 1),
                    font_size=formatting.get("font_size", 18),
                    bold=formatting.get("bold", False),
                    italic=formatting.get("italic", False)
                )
            
            elif element_type == "image":
                # Process image element
                image_source = self._substitute_variables(element.get("source", ""), data)
                position = element.get("position", {})
                
                if image_source and image_source != element.get("source", ""):
                    self.add_image(
                        prs_id, slide_index, image_source,
                        left=position.get("left", 1),
                        top=position.get("top", 1),
                        width=position.get("width"),
                        height=position.get("height")
                    )
            
            elif element_type == "chart":
                # Process chart element
                chart_data = element.get("data", {})
                categories = self._get_nested_value(data, chart_data.get("categories", ""))
                series_data = self._get_nested_value(data, chart_data.get("series", ""))
                
                if categories and series_data:
                    position = element.get("position", {})
                    self.add_chart(
                        prs_id, slide_index, 
                        chart_type=element.get("chart_type", "column"),
                        categories=categories,
                        series_data=series_data,
                        left=position.get("left", 2),
                        top=position.get("top", 2),
                        width=position.get("width", 6),
                        height=position.get("height", 4.5)
                    )
            
        except Exception as e:
            logger.error(f"Error processing template element: {e}")

# Global PowerPoint manager instance
ppt_manager = PowerPointManager()

@server.list_resources()
async def handle_list_resources() -> List[Resource]:
    """List available PowerPoint presentations as resources"""
    resources = []
    
    for prs_id, prs in ppt_manager.presentations.items():
        resources.append(
            Resource(
                uri=f"powerpoint://{prs_id}",
                name=f"PowerPoint Presentation {prs_id}",
                description=f"PowerPoint presentation with {len(prs.slides)} slides",
                mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        )
    
    return resources

@server.read_resource()
async def handle_read_resource(uri: str) -> str:
    """Read a PowerPoint presentation resource"""
    if not uri.startswith("powerpoint://"):
        raise ValueError(f"Unsupported URI scheme: {uri}")
    
    prs_id = uri.replace("powerpoint://", "")
    
    if prs_id not in ppt_manager.presentations:
        raise ValueError(f"Presentation {prs_id} not found")
    
    # Extract and return presentation information
    info = ppt_manager.get_presentation_info(prs_id)
    text_content = ppt_manager.extract_text(prs_id)
    
    result = {
        'presentation_info': info,
        'content': text_content
    }
    
    return json.dumps(result, indent=2)

@server.list_tools()
async def handle_list_tools() -> List[Tool]:
    """List available PowerPoint manipulation tools"""
    return [
        Tool(
            name="create_presentation",
            description="Create a new PowerPoint presentation, optionally from a template",
            inputSchema={
                "type": "object",
                "properties": {
                    "template_path": {
                        "type": "string",
                        "description": "Optional path to template file"
                    }
                },
                "examples": [
                    {},
                    {
                        "template_path": "./templates/corporate_template.pptx"
                    }
                ]
            }
        ),
        Tool(
            name="load_presentation",
            description="Load an existing PowerPoint presentation from file",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the PowerPoint file to load"
                    }
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="save_presentation",
            description="Save a PowerPoint presentation to file",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {
                        "type": "string",
                        "description": "ID of the presentation to save"
                    },
                    "file_path": {
                        "type": "string",
                        "description": "Path where to save the presentation"
                    }
                },
                "required": ["presentation_id", "file_path"],
                "examples": [
                    {
                        "presentation_id": "prs_0",
                        "file_path": "./output/quarterly_report.pptx"
                    }
                ]
            }
        ),
        Tool(
            name="add_slide",
            description="Add a new slide to a presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {
                        "type": "string",
                        "description": "ID of the presentation"
                    },
                    "layout_index": {
                        "type": "integer",
                        "description": "Slide layout index (0=title, 1=title+content, 6=blank, etc.)",
                        "default": 6
                    }
                },
                "required": ["presentation_id"]
            }
        ),
        Tool(
            name="add_text_box",
            description="Add a text box to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"},
                    "slide_index": {"type": "integer"},
                    "text": {"type": "string"},
                    "left": {"type": "number", "default": 1},
                    "top": {"type": "number", "default": 1},
                    "width": {"type": "number", "default": 8},
                    "height": {"type": "number", "default": 1},
                    "font_size": {"type": "integer", "default": 18},
                    "bold": {"type": "boolean", "default": False},
                    "italic": {"type": "boolean", "default": False}
                },
                "required": ["presentation_id", "slide_index", "text"],
                "examples": [
                    {
                        "presentation_id": "prs_0",
                        "slide_index": 0,
                        "text": "Quarterly Results Overview",
                        "font_size": 32,
                        "bold": true
                    },
                    {
                        "presentation_id": "prs_0",
                        "slide_index": 1,
                        "text": "• Revenue increased 25% year-over-year\n• New customer acquisition up 40%\n• Profit margins improved to 18%",
                        "top": 2,
                        "font_size": 20
                    }
                ]
            }
        ),
        Tool(
            name="add_image",
            description="Add an image to a slide from file path or URL",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"},
                    "slide_index": {"type": "integer"},
                    "image_source": {"type": "string", "description": "File path or URL to image"},
                    "left": {"type": "number", "default": 1},
                    "top": {"type": "number", "default": 1},
                    "width": {"type": "number"},
                    "height": {"type": "number"}
                },
                "required": ["presentation_id", "slide_index", "image_source"],
                "examples": [
                    {
                        "presentation_id": "prs_0",
                        "slide_index": 2,
                        "image_source": "./assets/company_logo.png",
                        "left": 8,
                        "top": 0.5,
                        "width": 2,
                        "height": 1
                    },
                    {
                        "presentation_id": "prs_0",
                        "slide_index": 3,
                        "image_source": "https://example.com/graph.png",
                        "left": 2,
                        "top": 2.5,
                        "width": 6,
                        "height": 4
                    }
                ]
            }
        ),
        Tool(
            name="add_chart",
            description="Add a chart to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"},
                    "slide_index": {"type": "integer"},
                    "chart_type": {"type": "string", "enum": ["column", "bar", "line", "pie"]},
                    "categories": {"type": "array", "items": {"type": "string"}},
                    "series_data": {
                        "type": "object",
                        "description": "Dictionary with series names as keys and data arrays as values"
                    },
                    "left": {"type": "number", "default": 2},
                    "top": {"type": "number", "default": 2},
                    "width": {"type": "number", "default": 6},
                    "height": {"type": "number", "default": 4.5}
                },
                "required": ["presentation_id", "slide_index", "chart_type", "categories", "series_data"],
                "examples": [
                    {
                        "presentation_id": "prs_0",
                        "slide_index": 1,
                        "chart_type": "column",
                        "categories": ["Q1", "Q2", "Q3", "Q4"],
                        "series_data": {
                            "Revenue": [100, 150, 120, 180],
                            "Profit": [20, 30, 25, 40]
                        }
                    }
                ]
            }
        ),
        Tool(
            name="extract_text",
            description="Extract all text content from a presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"}
                },
                "required": ["presentation_id"]
            }
        ),
        Tool(
            name="get_presentation_info",
            description="Get information about a presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string"}
                },
                "required": ["presentation_id"]
            }
        ),
        Tool(
            name="screenshot_slides",
            description="Screenshot each slide of a PowerPoint presentation for vision review (Windows only)",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the PowerPoint file"
                    },
                    "output_dir": {
                        "type": "string",
                        "description": "Directory to save screenshots (optional, defaults to temp directory)"
                    },
                    "image_format": {
                        "type": "string",
                        "description": "Image format (PNG, JPG, etc.)",
                        "default": "PNG"
                    },
                    "width": {
                        "type": "integer",
                        "description": "Screenshot width in pixels",
                        "default": 1920
                    },
                    "height": {
                        "type": "integer",
                        "description": "Screenshot height in pixels",  
                        "default": 1080
                    }
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="analyze_presentation_style",
            description="Analyze a PowerPoint presentation to extract style patterns and create learning data",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the PowerPoint file to analyze"
                    }
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="create_style_profile",
            description="Create a reusable style profile from analysis results",
            inputSchema={
                "type": "object",
                "properties": {
                    "analysis_results": {
                        "type": "object",
                        "description": "Style analysis results from analyze_presentation_style"
                    },
                    "profile_name": {
                        "type": "string",
                        "description": "Optional name for the style profile"
                    }
                },
                "required": ["analysis_results"]
            }
        ),
        Tool(
            name="apply_style_profile",
            description="Apply a learned style profile to an existing presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {
                        "type": "string",
                        "description": "ID of the presentation to style"
                    },
                    "profile_name": {
                        "type": "string",
                        "description": "Name of the style profile to apply"
                    }
                },
                "required": ["presentation_id", "profile_name"]
            }
        ),
        Tool(
            name="save_style_profile",
            description="Save a style profile to JSON file for reuse",
            inputSchema={
                "type": "object",
                "properties": {
                    "profile_name": {
                        "type": "string",
                        "description": "Name of the style profile to save"
                    },
                    "file_path": {
                        "type": "string",
                        "description": "Path where to save the style profile JSON"
                    }
                },
                "required": ["profile_name", "file_path"]
            }
        ),
        Tool(
            name="load_style_profile",
            description="Load a style profile from JSON file",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the style profile JSON file"
                    }
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="list_style_profiles",
            description="List all available style profiles",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        ),
        Tool(
            name="create_from_json",
            description="Create a presentation from structured JSON data (schema-driven approach)",
            inputSchema={
                "type": "object",
                "properties": {
                    "json_data": {
                        "type": "object",
                        "description": "Structured presentation data"
                    },
                    "template_path": {
                        "type": "string",
                        "description": "Optional template file path"
                    }
                },
                "required": ["json_data"]
            }
        ),
        # Phase 1: Professional Formatting & Layout Tools
        Tool(
            name="create_layout_grid",
            description="Create a layout grid for professional alignment and spacing",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "columns": {"type": "integer", "description": "Number of grid columns"},
                    "rows": {"type": "integer", "description": "Number of grid rows"},
                    "margins": {
                        "type": "object",
                        "description": "Grid margins in inches",
                        "properties": {
                            "left": {"type": "number", "default": 0.5},
                            "right": {"type": "number", "default": 0.5},
                            "top": {"type": "number", "default": 0.5},
                            "bottom": {"type": "number", "default": 0.5}
                        }
                    }
                },
                "required": ["presentation_id", "columns", "rows"]
            }
        ),
        Tool(
            name="snap_to_grid",
            description="Snap a shape to grid position for professional alignment",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "slide_index": {"type": "integer", "description": "Index of the slide"},
                    "shape_id": {"type": "string", "description": "ID of the shape to position"},
                    "grid_position": {
                        "type": "array",
                        "description": "Grid position as [column, row]",
                        "items": {"type": "integer"},
                        "minItems": 2,
                        "maxItems": 2
                    },
                    "alignment": {"type": "string", "default": "top-left", "description": "Shape alignment within grid cell"}
                },
                "required": ["presentation_id", "slide_index", "shape_id", "grid_position"]
            }
        ),
        Tool(
            name="distribute_shapes",
            description="Distribute shapes evenly for professional spacing",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "slide_index": {"type": "integer", "description": "Index of the slide"},
                    "shape_ids": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of shape IDs to distribute"
                    },
                    "distribution_type": {
                        "type": "string",
                        "enum": ["horizontal", "vertical"],
                        "description": "Direction of distribution"
                    }
                },
                "required": ["presentation_id", "slide_index", "shape_ids", "distribution_type"]
            }
        ),
        Tool(
            name="create_color_palette",
            description="Create a color palette for brand consistency",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "palette_name": {"type": "string", "description": "Name of the color palette"},
                    "colors": {
                        "type": "object",
                        "description": "Custom colors as role: hex_color pairs",
                        "additionalProperties": {"type": "string"}
                    }
                },
                "required": ["presentation_id", "palette_name"]
            }
        ),
        Tool(
            name="apply_color_palette",
            description="Apply color palette to presentation elements",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "palette_name": {"type": "string", "description": "Name of the color palette to apply"}
                },
                "required": ["presentation_id", "palette_name"]
            }
        ),
        Tool(
            name="create_typography_profile",
            description="Create a typography profile with style hierarchies",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "profile_name": {"type": "string", "description": "Name of the typography profile"},
                    "typography_config": {
                        "type": "object",
                        "description": "Typography configuration for different text levels",
                        "additionalProperties": {
                            "type": "object",
                            "properties": {
                                "font_name": {"type": "string"},
                                "font_size": {"type": "integer"},
                                "bold": {"type": "boolean"},
                                "color": {"type": "string"}
                            }
                        }
                    }
                },
                "required": ["presentation_id", "profile_name", "typography_config"]
            }
        ),
        Tool(
            name="apply_typography_style",
            description="Apply typography style to a text shape",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "slide_index": {"type": "integer", "description": "Index of the slide"},
                    "shape_id": {"type": "string", "description": "ID of the text shape"},
                    "style_type": {
                        "type": "string",
                        "enum": ["title", "subtitle", "heading", "body", "caption"],
                        "description": "Type of typography style to apply"
                    },
                    "profile_name": {
                        "type": "string",
                        "default": "default",
                        "description": "Name of the typography profile to use"
                    }
                },
                "required": ["presentation_id", "slide_index", "shape_id", "style_type"]
            }
        ),
        Tool(
            name="add_professional_shape",
            description="Add a professional shape from the shape library",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "slide_index": {"type": "integer", "description": "Index of the slide"},
                    "shape_category": {
                        "type": "string",
                        "enum": ["arrows", "callouts", "geometric"],
                        "description": "Category of professional shape"
                    },
                    "shape_name": {"type": "string", "description": "Name or index of the shape within category"},
                    "left": {"type": "number", "default": 1, "description": "Left position in inches"},
                    "top": {"type": "number", "default": 1, "description": "Top position in inches"},
                    "width": {"type": "number", "default": 2, "description": "Width in inches"},
                    "height": {"type": "number", "default": 2, "description": "Height in inches"}
                },
                "required": ["presentation_id", "slide_index", "shape_category", "shape_name"]
            }
        ),
        Tool(
            name="list_shape_library",
            description="List available professional shapes in the shape library",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        ),
        Tool(
            name="create_master_slide_theme",
            description="Create a master slide theme with consistent formatting",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "theme_name": {"type": "string", "description": "Name of the master slide theme"},
                    "theme_config": {
                        "type": "object",
                        "description": "Theme configuration for the master slide",
                        "properties": {
                            "background_color": {"type": "string", "description": "Hex color code for the background"},
                            "title_font": {
                                "type": "object",
                                "description": "Font configuration for title text",
                                "properties": {
                                    "name": {"type": "string", "description": "Font name"},
                                    "size": {"type": "integer", "description": "Font size"},
                                    "color": {"type": "string", "description": "Font color"},
                                    "bold": {"type": "boolean", "description": "Font bold flag"}
                                }
                            },
                            "content_font": {
                                "type": "object",
                                "description": "Font configuration for content text",
                                "properties": {
                                    "name": {"type": "string", "description": "Font name"},
                                    "size": {"type": "integer", "description": "Font size"},
                                    "color": {"type": "string", "description": "Font color"},
                                    "bold": {"type": "boolean", "description": "Font bold flag"}
                                }
                            },
                            "accent_color": {"type": "string", "description": "Hex color code for accent color"},
                            "layout_margins": {
                                "type": "object",
                                "description": "Layout margins for the master slide",
                                "properties": {
                                    "left": {"type": "number", "description": "Left margin"},
                                    "right": {"type": "number", "description": "Right margin"},
                                    "top": {"type": "number", "description": "Top margin"},
                                    "bottom": {"type": "number", "description": "Bottom margin"}
                                }
                            }
                        }
                    }
                },
                "required": ["presentation_id", "theme_name", "theme_config"]
            }
        ),
        Tool(
            name="apply_master_theme",
            description="Apply master theme to all slides in presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "theme_name": {"type": "string", "description": "Name of the master theme to apply"}
                },
                "required": ["presentation_id", "theme_name"]
            }
        ),
        Tool(
            name="list_master_themes",
            description="List available master themes for a presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"}
                },
                "required": ["presentation_id"]
            }
        ),
        Tool(
            name="set_slide_layout_template",
            description="Apply a layout template to a specific slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation"},
                    "slide_index": {"type": "integer", "description": "Index of the slide"},
                    "template_config": {
                        "type": "object",
                        "description": "Layout template configuration",
                        "properties": {
                            "clear_existing": {"type": "boolean", "description": "Flag to clear existing content"},
                            "layout_type": {"type": "string", "enum": ["title_content", "two_column"], "description": "Type of layout template"},
                            "title": {"type": "string", "description": "Title for title_content layout"},
                            "content": {"type": "string", "description": "Content for title_content layout"},
                            "left_content": {"type": "string", "description": "Left content for two_column layout"},
                            "right_content": {"type": "string", "description": "Right content for two_column layout"}
                        }
                    }
                },
                "required": ["presentation_id", "slide_index", "template_config"]
            }
        ),
        # Phase 2: Content Automation & Templates Tools
        Tool(
            name="create_template",
            description="Create a reusable template with placeholders and rules",
            inputSchema={
                "type": "object",
                "properties": {
                    "template_config": {
                        "type": "object",
                        "description": "Template configuration with slides, elements, and logic",
                        "properties": {
                            "name": {"type": "string", "description": "Template name"},
                            "description": {"type": "string", "description": "Template description"},
                            "slides": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "layout_type": {"type": "string", "description": "Slide layout type"},
                                        "elements": {
                                            "type": "array",
                                            "items": {
                                                "type": "object",
                                                "properties": {
                                                    "type": {"type": "string", "enum": ["text", "image", "chart"]},
                                                    "content": {"type": "string", "description": "Content with {{placeholders}}"},
                                                    "position": {"type": "object", "description": "Element position"},
                                                    "formatting": {"type": "object", "description": "Element formatting"}
                                                }
                                            }
                                        },
                                        "conditional_logic": {"type": "object", "description": "Conditions for slide inclusion"}
                                    }
                                }
                            }
                        }
                    }
                },
                "required": ["template_config"]
            }
        ),
        Tool(
            name="apply_template",
            description="Apply a template with data substitution to create a new presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "template_id": {"type": "string", "description": "ID of the template to apply"},
                    "data": {
                        "type": "object", 
                        "description": "Data to substitute into template placeholders",
                        "additionalProperties": True
                    }
                },
                "required": ["template_id", "data"]
            }
        ),
        Tool(
            name="update_template_content",
            description="Update content in an existing presentation using template data",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {"type": "string", "description": "ID of the presentation to update"},
                    "updates": {
                        "type": "object",
                        "description": "Updates per slide index",
                        "additionalProperties": {
                            "type": "object",
                            "description": "Data to update on the slide"
                        }
                    }
                },
                "required": ["presentation_id", "updates"]
            }
        ),
        Tool(
            name="bulk_generate_presentations",
            description="Generate multiple presentations from a template and data sets",
            inputSchema={
                "type": "object",
                "properties": {
                    "template_id": {"type": "string", "description": "ID of the template to use"},
                    "data_sets": {
                        "type": "array",
                        "items": {"type": "object"},
                        "description": "Array of data objects for each presentation"
                    },
                    "output_config": {
                        "type": "object",
                        "properties": {
                            "auto_save": {"type": "boolean", "description": "Automatically save generated presentations"},
                            "output_path": {"type": "string", "description": "Directory to save presentations"}
                        }
                    }
                },
                "required": ["template_id", "data_sets"]
            }
        ),
        Tool(
            name="map_data_source",
            description="Configure a data source for template mapping",
            inputSchema={
                "type": "object",
                "properties": {
                    "source_config": {
                        "type": "object",
                        "properties": {
                            "type": {"type": "string", "enum": ["json", "csv", "excel", "api"], "description": "Data source type"},
                            "source": {"type": "string", "description": "File path or URL to data source"},
                            "mapping": {
                                "type": "object",
                                "description": "Field mappings from data to template placeholders",
                                "additionalProperties": {"type": "string"}
                            },
                            "refresh_interval": {"type": "integer", "description": "Refresh interval in seconds"}
                        },
                        "required": ["type", "source"]
                    }
                },
                "required": ["source_config"]
            }
        ),
        Tool(
            name="list_templates", 
            description="List all available templates",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        ),
        Tool(
            name="delete_template",
            description="Delete a template",
            inputSchema={
                "type": "object",
                "properties": {
                    "template_id": {"type": "string", "description": "ID of the template to delete"}
                },
                "required": ["template_id"]
            }
        ),
        
        # Analytics and Monitoring Tools
        Tool(
            name="get_operation_analytics",
            description="Get server operation analytics and usage patterns",
            inputSchema={
                "type": "object",
                "properties": {
                    "hours": {
                        "type": "integer",
                        "default": 24,
                        "description": "Number of hours to analyze (default: 24)"
                    },
                    "format": {
                        "type": "string",
                        "enum": ["summary", "detailed"],
                        "default": "summary",
                        "description": "Analysis format level"
                    }
                }
            }
        )
    ]

@server.call_tool()
async def handle_call_tool(name: str, arguments: Dict[str, Any]) -> List[TextContent]:
    """Handle tool calls for PowerPoint operations"""
    start_time = time.time()
    
    try:
        # Extract user preferences if provided
        user_preferences = arguments.pop('_preferences', None)
        
        # Validate input arguments before processing
        try:
            validated_args = validate_request(name, arguments)
        except ValueError as e:
            # Log validation error
            operation_logger.log_validation_error(name, str(e))
            
            return [TextContent(
                type="text",
                text=f"❌ Invalid input: {str(e)}"
            )]
        
        if name == "create_presentation":
            template_path = validated_args.get("template_path")
            prs_id = ppt_manager.create_presentation(template_path)
            
            success_msg = format_success_message("create_presentation", 
                                               presentation_id=prs_id, 
                                               template_path=template_path)
            
            # Apply user preferences if provided
            if user_preferences:
                success_msg = format_response_with_preferences("create_presentation", success_msg, user_preferences, **validated_args)
            
            # Log successful operation
            latency_ms = int((time.time() - start_time) * 1000)
            operation_logger.log_operation("create_presentation", True, latency_ms, f"template={bool(template_path)}")
            
            return [TextContent(
                type="text",
                text=success_msg
            )]
        
        elif name == "load_presentation":
            file_path = arguments["file_path"]
            prs_id = ppt_manager.load_presentation(file_path)
            return [TextContent(
                type="text",
                text=f"Loaded presentation with ID: {prs_id}"
            )]
        
        elif name == "save_presentation":
            prs_id = validated_args["presentation_id"]
            file_path = validated_args["file_path"]
            success = await ppt_manager.save_presentation_async(prs_id, file_path)
            
            # Return rich content with file access
            return [
                TextContent(
                    type="text",
                    text=f"✅ Presentation {prs_id} saved successfully to {file_path}"
                ),
                EmbeddedResource(
                    uri=f"file://{os.path.abspath(file_path)}",
                    mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            ]
        
        elif name == "add_slide":
            prs_id = arguments["presentation_id"]
            layout_index = arguments.get("layout_index", 6)
            slide_index = ppt_manager.add_slide(prs_id, layout_index)
            return [TextContent(
                type="text",
                text=f"Added slide {slide_index} to presentation {prs_id}"
            )]
        
        elif name == "add_text_box":
            result = ppt_manager.add_text_box(
                validated_args["presentation_id"],
                validated_args["slide_index"],
                validated_args["text"],
                validated_args.get("left", 1),
                validated_args.get("top", 1),
                validated_args.get("width", 8),
                validated_args.get("height", 1),
                validated_args.get("font_size", 18),
                validated_args.get("bold", False),
                validated_args.get("italic", False)
            )
            
            success_msg = format_success_message("add_text_box", **validated_args)
            
            return [TextContent(
                type="text",
                text=success_msg
            )]
        
        elif name == "add_image":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            image_source = validated_args["image_source"]
            left = validated_args.get("left", 1)
            top = validated_args.get("top", 1)
            width = validated_args.get("width")
            height = validated_args.get("height")
            
            result = ppt_manager.add_image(prs_id, slide_index, image_source, left, top, width, height)
            
            # Create rich response with image preview for local files
            response = [
                TextContent(
                    type="text",
                    text=f"✅ Image added to slide {slide_index} at position ({left}, {top}) inches"
                )
            ]
            
            # Add embedded resource if it's a local file
            if not image_source.startswith(('http://', 'https://')) and os.path.exists(image_source):
                # Determine mime type based on file extension
                ext = os.path.splitext(image_source)[1].lower()
                mime_types = {
                    '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
                    '.png': 'image/png', '.gif': 'image/gif',
                    '.bmp': 'image/bmp', '.svg': 'image/svg+xml'
                }
                mime_type = mime_types.get(ext, 'image/png')
                
                response.append(EmbeddedResource(
                    uri=f"file://{os.path.abspath(image_source)}",
                    mimeType=mime_type
                ))
            
            return response
        
        elif name == "add_chart":
            prs_id = validated_args["presentation_id"]
            slide_index = validated_args["slide_index"]
            chart_type = validated_args["chart_type"]
            categories = validated_args["categories"]
            series_data = validated_args["series_data"]
            left = validated_args.get("left", 2)
            top = validated_args.get("top", 2)
            width = validated_args.get("width", 6)
            height = validated_args.get("height", 4.5)
            
            result = ppt_manager.add_chart(prs_id, slide_index, chart_type, categories, series_data, left, top, width, height)
            
            # Create detailed success message
            success_msg = format_success_message("add_chart", **validated_args)
            
            return [TextContent(
                type="text",
                text=success_msg
            )]
        
        elif name == "extract_text":
            prs_id = arguments["presentation_id"]
            content = ppt_manager.extract_text(prs_id)
            return [TextContent(
                type="text",
                text=json.dumps(content, indent=2)
            )]
        
        elif name == "get_presentation_info":
            prs_id = arguments["presentation_id"]
            info = ppt_manager.get_presentation_info(prs_id)
            return [TextContent(
                type="text",
                text=json.dumps(info, indent=2)
            )]
        
        elif name == "screenshot_slides":
            file_path = arguments["file_path"]
            output_dir = arguments.get("output_dir")
            image_format = arguments.get("image_format", "PNG")
            width = arguments.get("width", 1920)
            height = arguments.get("height", 1080)
            
            screenshot_paths = await ppt_manager.screenshot_slides_async(
                file_path, output_dir, image_format, width, height
            )
            
            result_info = {
                "total_slides": len(screenshot_paths),
                "screenshot_paths": screenshot_paths,
                "image_format": image_format,
                "dimensions": f"{width}x{height}",
                "output_directory": os.path.dirname(screenshot_paths[0]) if screenshot_paths else None
            }
            
            return [TextContent(
                type="text",
                text=f"Successfully created {len(screenshot_paths)} slide screenshots.\n" +
                     json.dumps(result_info, indent=2)
            )]
        
        elif name == "analyze_presentation_style":
            file_path = arguments["file_path"]
            analysis = ppt_manager.analyze_presentation_style(file_path)
            return [TextContent(
                type="text",
                text=json.dumps(analysis, indent=2, default=str)
            )]
        
        elif name == "create_style_profile":
            analysis_results = arguments["analysis_results"]
            profile_name = arguments.get("profile_name")
            profile_id = ppt_manager.create_style_profile(analysis_results, profile_name)
            return [TextContent(
                type="text",
                text=f"Created style profile: {profile_id}"
            )]
        
        elif name == "apply_style_profile":
            prs_id = arguments["presentation_id"]
            profile_name = arguments["profile_name"]
            success = ppt_manager.apply_style_profile(prs_id, profile_name)
            return [TextContent(
                type="text",
                text=f"Applied style profile '{profile_name}' to presentation {prs_id}"
            )]
        
        elif name == "save_style_profile":
            profile_name = arguments["profile_name"]
            file_path = arguments["file_path"]
            success = ppt_manager.save_style_profile(profile_name, file_path)
            return [TextContent(
                type="text",
                text=f"Saved style profile '{profile_name}' to {file_path}"
            )]
        
        elif name == "load_style_profile":
            file_path = arguments["file_path"]
            profile_name = ppt_manager.load_style_profile(file_path)
            return [TextContent(
                type="text",
                text=f"Loaded style profile: {profile_name}"
            )]
        
        elif name == "list_style_profiles":
            profiles = ppt_manager.list_style_profiles()
            return [TextContent(
                type="text",
                text=json.dumps(profiles, indent=2)
            )]

        elif name == "create_from_json":
            json_data = arguments["json_data"]
            template_path = arguments.get("template_path")
            
            # Create presentation from JSON schema (simplified implementation)
            prs_id = ppt_manager.create_presentation(template_path)
            
            # Process JSON data to create slides
            # This is a simplified version - could be expanded based on specific schema
            if isinstance(json_data, dict):
                for slide_name, slide_data in json_data.items():
                    slide_index = ppt_manager.add_slide(prs_id)
                    
                    if "title" in slide_data:
                        ppt_manager.add_text_box(prs_id, slide_index, slide_data["title"], 
                                               top=1, font_size=24, bold=True)
                    
                    if "content" in slide_data:
                        ppt_manager.add_text_box(prs_id, slide_index, slide_data["content"], 
                                               top=2.5, height=4)
            
            return [TextContent(
                type="text",
                text=f"Created presentation from JSON data with ID: {prs_id}"
            )]
        
        # Phase 1: Professional Formatting & Layout Tool Handlers
        elif name == "create_layout_grid":
            prs_id = arguments["presentation_id"]
            columns = arguments["columns"]
            rows = arguments["rows"]
            margins = arguments.get("margins")
            success = ppt_manager.create_layout_grid(prs_id, columns, rows, margins)
            return [TextContent(
                type="text",
                text=f"Created {columns}x{rows} layout grid for presentation {prs_id}"
            )]
        
        elif name == "snap_to_grid":
            prs_id = arguments["presentation_id"]
            slide_index = arguments["slide_index"]
            shape_id = arguments["shape_id"]
            grid_position = tuple(arguments["grid_position"])
            alignment = arguments.get("alignment", "top-left")
            success = ppt_manager.snap_to_grid(prs_id, slide_index, shape_id, grid_position, alignment)
            return [TextContent(
                type="text",
                text=f"Snapped shape {shape_id} to grid position {grid_position}"
            )]
        
        elif name == "distribute_shapes":
            prs_id = arguments["presentation_id"]
            slide_index = arguments["slide_index"]
            shape_ids = arguments["shape_ids"]
            distribution_type = arguments["distribution_type"]
            success = ppt_manager.distribute_shapes(prs_id, slide_index, shape_ids, distribution_type)
            return [TextContent(
                type="text",
                text=f"Distributed {len(shape_ids)} shapes {distribution_type}ly"
            )]
        
        elif name == "create_color_palette":
            prs_id = arguments["presentation_id"]
            palette_name = arguments["palette_name"]
            colors = arguments.get("colors")
            success = ppt_manager.create_color_palette(prs_id, palette_name, colors)
            return [TextContent(
                type="text",
                text=f"Created color palette '{palette_name}' for presentation {prs_id}"
            )]
        
        elif name == "apply_color_palette":
            prs_id = arguments["presentation_id"]
            palette_name = arguments["palette_name"]
            success = ppt_manager.apply_color_palette(prs_id, palette_name)
            return [TextContent(
                type="text",
                text=f"Applied color palette '{palette_name}' to presentation {prs_id}"
            )]
        
        elif name == "create_typography_profile":
            prs_id = arguments["presentation_id"]
            profile_name = arguments["profile_name"]
            typography_config = arguments["typography_config"]
            success = ppt_manager.create_typography_profile(prs_id, profile_name, typography_config)
            return [TextContent(
                type="text",
                text=f"Created typography profile '{profile_name}' for presentation {prs_id}"
            )]
        
        elif name == "apply_typography_style":
            prs_id = arguments["presentation_id"]
            slide_index = arguments["slide_index"]
            shape_id = arguments["shape_id"]
            style_type = arguments["style_type"]
            profile_name = arguments.get("profile_name", "default")
            success = ppt_manager.apply_typography_style(prs_id, slide_index, shape_id, style_type, profile_name)
            return [TextContent(
                type="text",
                text=f"Applied typography style '{style_type}' to shape {shape_id}"
            )]
        
        elif name == "add_professional_shape":
            prs_id = arguments["presentation_id"]
            slide_index = arguments["slide_index"]
            shape_category = arguments["shape_category"]
            shape_name = arguments["shape_name"]
            left = arguments.get("left", 1)
            top = arguments.get("top", 1)
            width = arguments.get("width", 2)
            height = arguments.get("height", 2)
            success = ppt_manager.add_professional_shape(prs_id, slide_index, shape_category, shape_name, left, top, width, height)
            return [TextContent(
                type="text",
                text=f"Added professional {shape_category} shape '{shape_name}' to slide {slide_index}"
            )]
        
        elif name == "list_shape_library":
            library = ppt_manager.list_shape_library()
            return [TextContent(
                type="text",
                text=json.dumps(library, indent=2)
            )]
        
        elif name == "create_master_slide_theme":
            prs_id = arguments["presentation_id"]
            theme_name = arguments["theme_name"]
            theme_config = arguments["theme_config"]
            success = ppt_manager.create_master_slide_theme(prs_id, theme_name, theme_config)
            return [TextContent(
                type="text",
                text=f"Created master slide theme '{theme_name}' for presentation {prs_id}"
            )]
        
        elif name == "apply_master_theme":
            prs_id = arguments["presentation_id"]
            theme_name = arguments["theme_name"]
            success = ppt_manager.apply_master_theme(prs_id, theme_name)
            return [TextContent(
                type="text",
                text=f"Applied master theme '{theme_name}' to presentation {prs_id}"
            )]
        
        elif name == "list_master_themes":
            prs_id = arguments["presentation_id"]
            themes = ppt_manager.list_master_themes(prs_id)
            return [TextContent(
                type="text",
                text=json.dumps(themes, indent=2)
            )]
        
        elif name == "set_slide_layout_template":
            prs_id = arguments["presentation_id"]
            slide_index = arguments["slide_index"]
            template_config = arguments["template_config"]
            success = ppt_manager.set_slide_layout_template(prs_id, slide_index, template_config)
            return [TextContent(
                type="text",
                text=f"Applied layout template to slide {slide_index}"
            )]
        
        # Phase 2: Content Automation & Templates Tool Handlers
        elif name == "create_template":
            template_config = arguments["template_config"]
            template_id = ppt_manager.create_template(template_config)
            return [TextContent(
                type="text",
                text=f"Created template with ID: {template_id}"
            )]
        
        elif name == "apply_template":
            template_id = arguments["template_id"]
            data = arguments["data"]
            prs_id = ppt_manager.apply_template(template_id, data)
            return [TextContent(
                type="text",
                text=f"Applied template {template_id} to create presentation {prs_id}"
            )]
        
        elif name == "update_template_content":
            prs_id = arguments["presentation_id"]
            updates = arguments["updates"]
            success = ppt_manager.update_template_content(prs_id, updates)
            return [TextContent(
                type="text",
                text=f"Updated content in presentation {prs_id}"
            )]
        
        elif name == "bulk_generate_presentations":
            template_id = arguments["template_id"]
            data_sets = arguments["data_sets"]
            output_config = arguments.get("output_config")
            presentation_ids = ppt_manager.bulk_generate_presentations(template_id, data_sets, output_config)
            return [TextContent(
                type="text",
                text=f"Bulk generated {len(presentation_ids)} presentations from template {template_id}: {presentation_ids}"
            )]
        
        elif name == "map_data_source":
            source_config = arguments["source_config"]
            source_id = ppt_manager.map_data_source(source_config)
            return [TextContent(
                type="text",
                text=f"Mapped data source with ID: {source_id}"
            )]
        
        elif name == "list_templates":
            templates = ppt_manager.list_templates()
            return [TextContent(
                type="text",
                text=json.dumps(templates, indent=2)
            )]
        
        elif name == "delete_template":
            template_id = arguments["template_id"]
            success = ppt_manager.delete_template(template_id)
            return [TextContent(
                type="text",
                text=f"Deleted template {template_id}"
            )]
            
        elif name == "get_operation_analytics":
            hours = arguments.get("hours", 24)
            format_type = arguments.get("format", "summary")
            
            analytics = operation_logger.get_operation_summary(hours)
            
            if format_type == "summary":
                # Create a concise summary
                if "error" in analytics:
                    summary_text = f"Analytics Error: {analytics['error']}"
                elif "message" in analytics:
                    summary_text = analytics["message"]
                else:
                    summary_text = f"""📊 Server Analytics (Last {hours}h)
✅ Operations: {analytics['total_operations']} ({analytics['success_rate']}% success)
⚡ Avg Response: {analytics['average_latency_ms']}ms
🔥 Top Tools: {', '.join([f"{tool}({count})" for tool, count in analytics['tool_usage'][:3]])}
❌ Error Types: {', '.join([f"{err_type}({count})" for err_type, count in analytics['error_types'].items()]) if analytics['error_types'] else 'None'}"""
            else:
                # Return detailed JSON
                summary_text = json.dumps(analytics, indent=2)
            
            return [TextContent(
                type="text",
                text=summary_text
            )]
        
        else:
            raise ValueError(f"Unknown tool: {name}")
    
    except Exception as e:
        # Log failed operation
        latency_ms = int((time.time() - start_time) * 1000)
        operation_logger.log_operation(name, False, latency_ms, error=str(e), error_type="execution")
        
        logger.error(f"Error in tool {name}: {e}")
        return [TextContent(
            type="text",
            text=f"Error: {str(e)}"
        )]

async def main():
    """Main entry point for the PowerPoint MCP server"""
    # Initialize server options
    options = InitializationOptions(
        server_name="powerpoint-mcp-server",
        server_version="1.0.0",
        capabilities={
            "resources": {},
            "tools": {},
            "logging": {}
        }
    )
    
    try:
        async with stdio_server() as (read_stream, write_stream):
            await server.run(
                read_stream,
                write_stream,
                options
            )
    except KeyboardInterrupt:
        logger.info("Server interrupted by user")
    except Exception as e:
        logger.error(f"Server error: {e}")
    finally:
        # Cleanup
        ppt_manager.cleanup()
        logger.info("PowerPoint MCP server shutdown complete")

if __name__ == "__main__":
    asyncio.run(main()) 